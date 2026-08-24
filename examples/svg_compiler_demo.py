"""SVG Compiler 全能力演示

生成一个 8 页 PPTX，每页展示不同的 SVG 编译能力：
1. 基础形状 + 纯色填充
2. 线性渐变 + 径向渐变
3. 路径命令 (M/L/C/S/Q/A/Z)
4. 文本 + tspan + 中文
5. Transform (translate/scale/rotate/matrix)
6. ClipPath + evenodd
7. 组合图表 (漏斗 + 数据)
8. 赛博朋克风格仪表盘
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pptx_designer.tools.svg import svg_chart


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title(slide, title_text):
    """添加页标题"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)


def demo_basic_shapes(slide):
    """演示 1: 基础形状"""
    add_title(slide, "1. Basic Shapes — 基础形状")
    svg = '''<svg viewBox="0 0 800 400" xmlns="http://www.w3.org/2000/svg">
        <!-- 矩形 -->
        <rect x="30" y="30" width="200" height="150" rx="10" fill="#4472C4"/>
        <rect x="30" y="220" width="200" height="150" ry="15" fill="#2E75B6"/>

        <!-- 圆形 + 椭圆 -->
        <circle cx="420" cy="105" r="80" fill="#E74C3C"/>
        <ellipse cx="420" cy="300" rx="100" ry="60" fill="#2ECC71"/>

        <!-- 多边形 -->
        <polygon points="700,30 780,180 620,180" fill="#F39C12"/>
        <polygon points="700,220 760,320 640,320 640,220" fill="#9B59B6"/>

        <!-- 线条 -->
        <line x1="280" y1="200" x2="380" y2="200" stroke="#333" stroke-width="3"/>
        <line x1="280" y1="220" x2="380" y2="220" stroke="#E74C3C" stroke-width="2" stroke-dasharray="8,4"/>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.2, w=11.5, h=5.8)


def demo_gradients(slide):
    """演示 2: 渐变"""
    add_title(slide, "2. Gradients — 线性渐变 + 径向渐变")
    svg = '''<svg viewBox="0 0 800 400" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="g1" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0%" stop-color="#667eea"/>
                <stop offset="100%" stop-color="#764ba2"/>
            </linearGradient>
            <linearGradient id="g2" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="#f093fb"/>
                <stop offset="100%" stop-color="#f5576c"/>
            </linearGradient>
            <linearGradient id="g3" x1="0" y1="0" x2="1" y2="1">
                <stop offset="0%" stop-color="#4facfe"/>
                <stop offset="100%" stop-color="#00f2fe"/>
            </linearGradient>
            <radialGradient id="r1" cx="50%" cy="50%" r="50%">
                <stop offset="0%" stop-color="#FFF"/>
                <stop offset="100%" stop-color="#4472C4"/>
            </radialGradient>
            <radialGradient id="r2" cx="30%" cy="30%" r="60%">
                <stop offset="0%" stop-color="#FFD700"/>
                <stop offset="50%" stop-color="#FF6B35"/>
                <stop offset="100%" stop-color="#C62828"/>
            </radialGradient>
        </defs>

        <!-- 线性渐变矩形 -->
        <rect x="30" y="30" width="240" height="160" rx="8" fill="url(#g1)"/>
        <rect x="30" y="210" width="240" height="160" rx="8" fill="url(#g2)"/>
        <rect x="300" y="30" width="240" height="340" rx="8" fill="url(#g3)"/>

        <!-- 径向渐变圆 -->
        <circle cx="650" cy="120" r="90" fill="url(#r1)"/>
        <circle cx="650" cy="300" r="90" fill="url(#r2)"/>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.2, w=11.5, h=5.8)


def demo_paths(slide):
    """演示 3: 路径命令"""
    add_title(slide, "3. Path Commands — 贝塞尔曲线 + 弧线")
    svg = '''<svg viewBox="0 0 800 400" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="wave" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0%" stop-color="#00b4db"/>
                <stop offset="100%" stop-color="#0083b0"/>
            </linearGradient>
        </defs>

        <!-- 贝塞尔曲线: S 命令 (平滑三次) -->
        <path d="M30,200 C130,50 270,350 400,200 S670,50 770,200"
              fill="none" stroke="url(#wave)" stroke-width="4"/>

        <!-- Q 命令 (二次贝塞尔) -->
        <path d="M30,350 Q200,250 400,350 T770,350"
              fill="none" stroke="#E74C3C" stroke-width="3"/>

        <!-- A 命令 (弧线) -->
        <path d="M100,50 A60,60 0 1,1 220,50" fill="#2ECC71" fill-opacity="0.7"/>
        <path d="M300,50 A80,40 0 0,1 460,50" fill="#F39C12" fill-opacity="0.7"/>

        <!-- 复杂路径: 心形 -->
        <path d="M400,380 C400,380 250,280 250,200 C250,140 300,100 350,100
                 C380,100 400,120 400,150 C400,120 420,100 450,100
                 C500,100 550,140 550,200 C550,280 400,380 400,380 Z"
              fill="#E74C3C" fill-opacity="0.8"/>

        <!-- Z 闭合路径 -->
        <path d="M650,100 L750,250 L550,250 Z" fill="#9B59B6" fill-opacity="0.6"/>
        <path d="M650,280 L720,400 L580,400 Z" fill="#3498DB" fill-opacity="0.6"/>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.2, w=11.5, h=5.8)


def demo_text(slide):
    """演示 4: 文本 + tspan + 中文"""
    add_title(slide, "4. Text & Tspan — 多语言文本排版")
    svg = '''<svg viewBox="0 0 800 400" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="txtGrad" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0%" stop-color="#667eea"/>
                <stop offset="100%" stop-color="#764ba2"/>
            </linearGradient>
        </defs>

        <!-- 大标题 -->
        <text x="400" y="60" text-anchor="middle" font-size="36" font-weight="bold"
              fill="url(#txtGrad)" font-family="Arial">SVG Compiler Demo</text>

        <!-- 中文文本 -->
        <text x="400" y="110" text-anchor="middle" font-size="24" fill="#333"
              font-family="Microsoft YaHei">SVG 编译器能力演示</text>

        <!-- 多行 tspan -->
        <text x="50" y="170" font-size="18" fill="#555">
            <tspan font-weight="bold" fill="#E74C3C">Features:</tspan>
            <tspan dx="10" fill="#333">Shapes, Paths, Gradients, Transforms</tspan>
        </text>
        <text x="50" y="200" font-size="18" fill="#555">
            <tspan font-weight="bold" fill="#2ECC71">优势:</tspan>
            <tspan dx="10" fill="#333">完全可编辑, 原生 PPTX 形状, 无图片嵌入</tspan>
        </text>
        <text x="50" y="230" font-size="18" fill="#555">
            <tspan font-weight="bold" fill="#3498DB">Use Cases:</tspan>
            <tspan dx="10" fill="#333">架构图, 流程图, 数据可视化, 仪表盘</tspan>
        </text>

        <!-- 混合语言段落 -->
        <text x="400" y="300" text-anchor="middle" font-size="20" fill="#666">
            <tspan>PowerPoint native shapes</tspan>
            <tspan dx="5" fill="#E74C3C">+</tspan>
            <tspan dx="5" fill="#666">SVG flexibility</tspan>
            <tspan dx="5" fill="#E74C3C">=</tspan>
            <tspan dx="5" fill="#2ECC71" font-weight="bold">Beautiful diagrams</tspan>
        </text>

        <!-- 底部注释 -->
        <text x="400" y="370" text-anchor="middle" font-size="14" fill="#999">
            Generated by pptx-designer • MIT License • github.com/sunchaokun/pptx-designer
        </text>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.2, w=11.5, h=5.8)


def demo_transforms(slide):
    """演示 5: Transform"""
    add_title(slide, "5. Transforms — 平移/缩放/旋转/矩阵")
    svg = '''<svg viewBox="0 0 800 400" xmlns="http://www.w3.org/2000/svg">
        <!-- Translate -->
        <g transform="translate(80, 80)">
            <rect x="0" y="0" width="120" height="80" fill="#4472C4" rx="4"/>
            <text x="60" y="45" text-anchor="middle" fill="#FFF" font-size="14">translate</text>
        </g>

        <!-- Scale -->
        <g transform="translate(300, 60) scale(0.8)">
            <rect x="0" y="0" width="150" height="100" fill="#E74C3C" rx="4"/>
            <text x="75" y="55" text-anchor="middle" fill="#FFF" font-size="14">scale(0.8)</text>
        </g>

        <!-- Rotate -->
        <g transform="translate(550, 100) rotate(15)">
            <rect x="-60" y="-30" width="120" height="60" fill="#2ECC71" rx="4"/>
            <text x="0" y="5" text-anchor="middle" fill="#FFF" font-size="14">rotate(15°)</text>
        </g>

        <!-- 组合变换 -->
        <g transform="translate(200, 250) rotate(10) scale(1.2)">
            <rect x="0" y="0" width="140" height="90" fill="#9B59B6" rx="4"/>
            <text x="70" y="50" text-anchor="middle" fill="#FFF" font-size="13">rotate+scale</text>
        </g>

        <!-- Matrix -->
        <g transform="translate(500, 230) matrix(1, 0.2, -0.2, 1, 0, 0)">
            <rect x="0" y="0" width="130" height="80" fill="#F39C12" rx="4"/>
            <text x="65" y="45" text-anchor="middle" fill="#FFF" font-size="14">matrix skew</text>
        </g>

        <!-- 嵌套变换 -->
        <g transform="translate(100, 300)">
            <rect x="0" y="0" width="80" height="60" fill="#1ABC9C" rx="4"/>
            <g transform="translate(90, 10) rotate(20)">
                <rect x="0" y="0" width="80" height="60" fill="#E67E22" rx="4"/>
                <g transform="translate(90, 10) rotate(20)">
                    <rect x="0" y="0" width="80" height="60" fill="#E74C3C" rx="4"/>
                </g>
            </g>
        </g>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.2, w=11.5, h=5.8)


def demo_clippath(slide):
    """演示 6: ClipPath + evenodd"""
    add_title(slide, "6. ClipPath & EvenOdd — 裁剪与孔洞")
    svg = '''<svg viewBox="0 0 800 400" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="ocean" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="#667eea"/>
                <stop offset="100%" stop-color="#764ba2"/>
            </linearGradient>
            <linearGradient id="sunset" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="#f093fb"/>
                <stop offset="100%" stop-color="#f5576c"/>
            </linearGradient>
            <clipPath id="circle-clip">
                <circle cx="200" cy="200" r="140"/>
            </clipPath>
            <clipPath id="star-clip">
                <polygon points="100,10 127,90 210,90 142,140 167,220 100,170 33,220 58,140 -10,90 73,90"/>
            </clipPath>
        </defs>

        <!-- 圆形裁剪 -->
        <g clip-path="url(#circle-clip)">
            <rect x="0" y="0" width="800" height="400" fill="url(#ocean)"/>
            <circle cx="150" cy="150" r="40" fill="#FFF" fill-opacity="0.3"/>
            <circle cx="250" cy="250" r="30" fill="#FFF" fill-opacity="0.2"/>
            <circle cx="180" cy="280" r="25" fill="#FFF" fill-opacity="0.25"/>
        </g>
        <circle cx="200" cy="200" r="140" fill="none" stroke="#333" stroke-width="2"/>

        <!-- EvenOdd 孔洞效果 -->
        <path d="M450,50 L750,50 L750,350 L450,350 Z
                 M520,120 A40,40 0 1,0 600,120 A40,40 0 1,0 520,120 Z
                 M520,250 A30,30 0 1,0 580,250 A30,30 0 1,0 520,250 Z"
              fill="url(#sunset)" fill-rule="evenodd"/>

        <!-- 裁剪后的复杂图形 -->
        <g clip-path="url(#star-clip)" transform="translate(0, 20)">
            <rect x="0" y="0" width="800" height="400" fill="#2ECC71" fill-opacity="0.6"/>
            <circle cx="100" cy="100" r="80" fill="#F39C12" fill-opacity="0.5"/>
        </g>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.2, w=11.5, h=5.8)


def demo_funnel(slide):
    """演示 7: 漏斗图表"""
    add_title(slide, "7. Funnel Chart — 漏斗转化图")
    svg = '''<svg viewBox="0 0 800 500" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="f1" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0%" stop-color="#667eea"/><stop offset="100%" stop-color="#764ba2"/>
            </linearGradient>
            <linearGradient id="f2" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0%" stop-color="#4facfe"/><stop offset="100%" stop-color="#00f2fe"/>
            </linearGradient>
            <linearGradient id="f3" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0%" stop-color="#43e97b"/><stop offset="100%" stop-color="#38f9d7"/>
            </linearGradient>
            <linearGradient id="f4" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0%" stop-color="#fa709a"/><stop offset="100%" stop-color="#fee140"/>
            </linearGradient>
        </defs>

        <!-- 漏斗层级 -->
        <polygon points="60,30 740,30 680,100 120,100" fill="url(#f1)"/>
        <polygon points="120,115 680,115 620,185 180,185" fill="url(#f2)"/>
        <polygon points="180,200 620,200 560,270 240,270" fill="url(#f3)"/>
        <polygon points="240,285 560,285 500,355 300,355" fill="url(#f4)"/>

        <!-- 数据标签 -->
        <text x="400" y="75" text-anchor="middle" fill="#FFF" font-size="18" font-weight="bold">Impressions  10,000</text>
        <text x="400" y="160" text-anchor="middle" fill="#FFF" font-size="18" font-weight="bold">Clicks  4,500</text>
        <text x="400" y="245" text-anchor="middle" fill="#FFF" font-size="18" font-weight="bold">Signups  1,800</text>
        <text x="400" y="330" text-anchor="middle" fill="#FFF" font-size="18" font-weight="bold">Conversions  600</text>

        <!-- 转化率箭头 -->
        <text x="770" y="75" fill="#667eea" font-size="14">100%</text>
        <text x="770" y="160" fill="#4facfe" font-size="14">45%</text>
        <text x="770" y="245" fill="#43e97b" font-size="14">18%</text>
        <text x="770" y="330" fill="#fa709a" font-size="14">6%</text>

        <!-- 底部说明 -->
        <text x="400" y="420" text-anchor="middle" fill="#666" font-size="16">
            Conversion Funnel • Jan 2026
        </text>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.0, w=11.5, h=6.0)


def demo_dashboard(slide):
    """演示 8: 赛博朋克仪表盘"""
    add_title(slide, "8. Cyberpunk Dashboard — 暗色主题仪表盘")
    svg = '''<svg viewBox="0 0 800 450" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="neon-blue" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="#00f2fe"/><stop offset="100%" stop-color="#4facfe"/>
            </linearGradient>
            <linearGradient id="neon-pink" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="#f093fb"/><stop offset="100%" stop-color="#f5576c"/>
            </linearGradient>
            <linearGradient id="neon-green" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="#43e97b"/><stop offset="100%" stop-color="#38f9d7"/>
            </linearGradient>
            <radialGradient id="glow" cx="50%" cy="50%" r="50%">
                <stop offset="0%" stop-color="#4facfe" stop-opacity="0.3"/>
                <stop offset="100%" stop-color="#4facfe" stop-opacity="0"/>
            </radialGradient>
        </defs>

        <!-- 背景 -->
        <rect x="0" y="0" width="800" height="450" fill="#0D1117" rx="8"/>

        <!-- 网格线 -->
        <line x1="0" y1="100" x2="800" y2="100" stroke="#21262D" stroke-width="1"/>
        <line x1="0" y1="200" x2="800" y2="200" stroke="#21262D" stroke-width="1"/>
        <line x1="0" y1="300" x2="800" y2="300" stroke="#21262D" stroke-width="1"/>
        <line x1="0" y1="400" x2="800" y2="400" stroke="#21262D" stroke-width="1"/>

        <!-- KPI 卡片 1 -->
        <rect x="20" y="20" width="170" height="100" rx="8" fill="#161B22" stroke="#30363D"/>
        <text x="30" y="50" fill="#8B949E" font-size="12">REVENUE</text>
        <text x="30" y="80" fill="#00f2fe" font-size="28" font-weight="bold">$2.4M</text>
        <text x="30" y="105" fill="#3FB950" font-size="11">+12.5% ▲</text>

        <!-- KPI 卡片 2 -->
        <rect x="210" y="20" width="170" height="100" rx="8" fill="#161B22" stroke="#30363D"/>
        <text x="220" y="50" fill="#8B949E" font-size="12">USERS</text>
        <text x="220" y="80" fill="#f093fb" font-size="28" font-weight="bold">18,420</text>
        <text x="220" y="105" fill="#3FB950" font-size="11">+8.2% ▲</text>

        <!-- KPI 卡片 3 -->
        <rect x="400" y="20" width="170" height="100" rx="8" fill="#161B22" stroke="#30363D"/>
        <text x="410" y="50" fill="#8B949E" font-size="12">CONVERSION</text>
        <text x="410" y="80" fill="#43e97b" font-size="28" font-weight="bold">3.2%</text>
        <text x="410" y="105" fill="#F85149" font-size="11">-0.3% ▼</text>

        <!-- KPI 卡片 4 -->
        <rect x="590" y="20" width="190" height="100" rx="8" fill="#161B22" stroke="#30363D"/>
        <text x="600" y="50" fill="#8B949E" font-size="12">RESPONSE TIME</text>
        <text x="600" y="80" fill="#FFA657" font-size="28" font-weight="bold">42ms</text>
        <text x="600" y="105" fill="#3FB950" font-size="11">-5ms ▼</text>

        <!-- 面积图 -->
        <path d="M20,400 L20,250 L100,220 L200,280 L300,180 L400,200 L500,140 L600,160 L700,100 L780,120 L780,400 Z"
              fill="url(#neon-blue)" fill-opacity="0.2"/>
        <path d="M20,250 L100,220 L200,280 L300,180 L400,200 L500,140 L600,160 L700,100 L780,120"
              fill="none" stroke="#4facfe" stroke-width="2"/>

        <!-- 柱状图 -->
        <rect x="30" y="320" width="30" height="70" fill="#f093fb" fill-opacity="0.6" rx="2"/>
        <rect x="80" y="290" width="30" height="100" fill="#f093fb" fill-opacity="0.7" rx="2"/>
        <rect x="130" y="250" width="30" height="140" fill="#f093fb" fill-opacity="0.8" rx="2"/>
        <rect x="180" y="270" width="30" height="120" fill="#f093fb" fill-opacity="0.7" rx="2"/>
        <rect x="230" y="220" width="30" height="170" fill="#f093fb" fill-opacity="0.9" rx="2"/>
        <rect x="280" y="200" width="30" height="190" fill="#f093fb" rx="2"/>

        <!-- 环形进度 -->
        <circle cx="700" cy="350" r="60" fill="none" stroke="#21262D" stroke-width="12"/>
        <circle cx="700" cy="350" r="60" fill="none" stroke="#43e97b" stroke-width="12"
                stroke-dasharray="280 380" stroke-linecap="round"
                transform="rotate(-90, 700, 350)"/>
        <text x="700" y="355" text-anchor="middle" fill="#43e97b" font-size="20" font-weight="bold">74%</text>
        <text x="700" y="375" text-anchor="middle" fill="#8B949E" font-size="10">QUOTA</text>

        <!-- 底部状态栏 -->
        <rect x="0" y="420" width="800" height="30" fill="#161B22"/>
        <circle cx="20" cy="435" r="4" fill="#3FB950"/>
        <text x="30" y="440" fill="#8B949E" font-size="10">SYSTEM STATUS: OPERATIONAL</text>
        <text x="780" y="440" text-anchor="end" fill="#8B949E" font-size="10">Last updated: 2026-08-24 14:32 UTC</text>
    </svg>'''
    svg_chart(slide, svg, x=0.8, y=1.0, w=11.5, h=6.2)


def main():
    prs = new_prs()

    # Slide 1: Basic Shapes
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_basic_shapes(slide1)

    # Slide 2: Gradients
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_gradients(slide2)

    # Slide 3: Paths
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_paths(slide3)

    # Slide 4: Text
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_text(slide4)

    # Slide 5: Transforms
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_transforms(slide5)

    # Slide 6: ClipPath
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_clippath(slide6)

    # Slide 7: Funnel
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_funnel(slide7)

    # Slide 8: Dashboard
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    demo_dashboard(slide8)

    output = "E:/pptx-designer/output/svg_compiler_demo.pptx"
    prs.save(output)
    print(f"Generated: {output}")
    print(f"Total slides: {len(prs.slides)} (each demonstrates different SVG capabilities)")


if __name__ == "__main__":
    main()
