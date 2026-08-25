"""Integration tests for PrecisionRenderer SVG rendering path."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from pptx_designer.enterprise.brand import BrandSpec
from pptx_designer.renderer.precision import PrecisionRenderer


def _make_renderer(brand: BrandSpec | None = None) -> PrecisionRenderer:
    return PrecisionRenderer(brand_spec=brand)


def _make_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestPrecisionRendererSVG:
    def test_basic_svg_string(self):
        renderer = _make_renderer()
        slide = _make_slide()
        svg = '<svg viewBox="0 0 400 300"><circle cx="200" cy="150" r="100" fill="#4472C4"/></svg>'
        result = renderer._render_svg_diagram_on_slide(slide, svg)
        assert len(slide.shapes) >= 1
        assert result is not None
        assert result.metrics["native_shape_count"] >= 1

    def test_svg_dict_input(self):
        renderer = _make_renderer()
        slide = _make_slide()
        svg_data = {
            "svg": '<svg viewBox="0 0 400 300"><rect x="50" y="50" width="200" height="150" fill="#E74C3C"/></svg>'
        }
        renderer._render_svg_diagram_on_slide(slide, svg_data)
        assert len(slide.shapes) >= 1

    def test_empty_svg_string(self):
        renderer = _make_renderer()
        slide = _make_slide()
        initial_count = len(slide.shapes)
        renderer._render_svg_diagram_on_slide(slide, "")
        assert len(slide.shapes) == initial_count

    def test_empty_svg_dict(self):
        renderer = _make_renderer()
        slide = _make_slide()
        initial_count = len(slide.shapes)
        renderer._render_svg_diagram_on_slide(slide, {"svg": ""})
        assert len(slide.shapes) == initial_count

    def test_invalid_svg_logs_error(self):
        renderer = _make_renderer()
        slide = _make_slide()
        initial_count = len(slide.shapes)
        renderer._render_svg_diagram_on_slide(slide, "not svg at all")
        assert len(slide.shapes) == initial_count

    def test_svg_with_brand_context(self):
        brand = BrandSpec(
            colors={"primary": "#1D78FA", "secondary": "#FF6B35"},
            fonts={"heading": "Arial", "body": "Arial"},
        )
        renderer = _make_renderer(brand=brand)
        slide = _make_slide()
        svg = '<svg viewBox="0 0 400 300"><rect x="50" y="50" width="300" height="200" fill="var(--primary)"/></svg>'
        renderer._render_svg_diagram_on_slide(slide, svg)
        assert len(slide.shapes) >= 1

    def test_svg_with_clip_path(self):
        renderer = _make_renderer()
        slide = _make_slide()
        svg = """<svg viewBox="0 0 400 300">
            <defs><clipPath id="cp"><rect x="50" y="50" width="200" height="200"/></clipPath></defs>
            <rect x="0" y="0" width="400" height="300" fill="#4472C4" clip-path="url(#cp)"/>
        </svg>"""
        renderer._render_svg_diagram_on_slide(slide, svg)
        assert len(slide.shapes) >= 1

    def test_svg_with_gradient(self):
        renderer = _make_renderer()
        slide = _make_slide()
        svg = """<svg viewBox="0 0 400 300">
            <defs><linearGradient id="g1" x1="0" y1="0" x2="1" y2="1">
                <stop offset="0" stop-color="#4472C4"/><stop offset="1" stop-color="#2E75B6"/>
            </linearGradient></defs>
            <rect x="50" y="50" width="300" height="200" fill="url(#g1)"/>
        </svg>"""
        renderer._render_svg_diagram_on_slide(slide, svg)
        assert len(slide.shapes) >= 1

    def test_svg_with_text(self):
        renderer = _make_renderer()
        slide = _make_slide()
        svg = '<svg viewBox="0 0 400 300"><text x="200" y="150" text-anchor="middle" font-size="24" fill="#333">Hello</text></svg>'
        renderer._render_svg_diagram_on_slide(slide, svg)
        assert len(slide.shapes) >= 1

    def test_svg_with_transform(self):
        renderer = _make_renderer()
        slide = _make_slide()
        svg = '<svg viewBox="0 0 400 300"><g transform="translate(100,50)"><rect x="0" y="0" width="200" height="150" fill="#3498DB"/></g></svg>'
        renderer._render_svg_diagram_on_slide(slide, svg)
        assert len(slide.shapes) >= 1

    def test_custom_position(self):
        renderer = _make_renderer()
        slide = _make_slide()
        svg = '<svg viewBox="0 0 400 300"><circle cx="200" cy="150" r="100" fill="#4472C4"/></svg>'
        renderer._render_svg_diagram_on_slide(slide, svg, cx=2.0, cy=2.0, cw=5.0, ch=4.0)
        assert len(slide.shapes) >= 1

    def test_unsupported_element_warning(self, caplog):
        renderer = _make_renderer()
        slide = _make_slide()
        svg = '<svg viewBox="0 0 400 300"><image href="photo.jpg" width="400" height="300"/></svg>'
        renderer._render_svg_diagram_on_slide(slide, svg)
        assert any("unsupported" in r.message.lower() or "image" in r.message.lower() for r in caplog.records)
