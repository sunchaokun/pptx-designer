"""Unit tests for the SVG sanitizer module."""

from __future__ import annotations

import pytest
from lxml import etree

from pptx_designer.compiler._sanitizer import (
    SVG_NS,
    _expand_style,
    _fix_self_closing_lxml,
    _infer_viewbox,
    _strip_unwanted,
    _walk_expand,
    sanitize,
)
from pptx_designer.compiler._errors import SVGCompileError


SVG = f"{{{SVG_NS}}}"

_SVG_OPEN = '<svg xmlns="http://www.w3.org/2000/svg"'


def _parse_el(tag: str, attrs: str = "", style: str = "") -> etree._Element:
    """Parse a single SVG element for isolated testing."""
    s_attr = f' style="{style}"' if style else ""
    return etree.fromstring(f'<{tag} xmlns="http://www.w3.org/2000/svg"{attrs}{s_attr}/>'.encode())


def _find(el: etree._Element, tag: str) -> etree._Element | None:
    """Find first descendant by local name."""
    for desc in el.iter():
        local = desc.tag.split("}")[-1] if "}" in desc.tag else desc.tag
        if local == tag:
            return desc
    return None


def _findall(el: etree._Element, tag: str) -> list[etree._Element]:
    """Find all descendants by local name."""
    results = []
    for desc in el.iter():
        local = desc.tag.split("}")[-1] if "}" in desc.tag else desc.tag
        if local == tag:
            results.append(desc)
    return results


class TestExpandStyle:
    def test_extracts_fill(self):
        el = _parse_el("rect", style="fill: #FF0000")
        _expand_style(el)
        assert el.get("fill") == "#FF0000"
        assert el.get("style") is None

    def test_extracts_multiple_properties(self):
        el = _parse_el("rect", style="fill: red; stroke: blue; stroke-width: 2")
        _expand_style(el)
        assert el.get("fill") == "red"
        assert el.get("stroke") == "blue"
        assert el.get("stroke-width") == "2"

    def test_does_not_override_existing_attr(self):
        el = etree.fromstring(
            b'<rect xmlns="http://www.w3.org/2000/svg" fill="green" style="fill: red"/>'
        )
        _expand_style(el)
        assert el.get("fill") == "green"

    def test_no_style_attr(self):
        el = etree.fromstring(b'<rect xmlns="http://www.w3.org/2000/svg" fill="blue"/>')
        _expand_style(el)
        assert el.get("fill") == "blue"

    def test_empty_style_preserved(self):
        el = etree.fromstring(b'<rect xmlns="http://www.w3.org/2000/svg" style=""/>')
        _expand_style(el)
        assert el.get("style") == ""

    def test_expands_font_properties(self):
        el = _parse_el("text", style="font-size: 14px; font-family: Arial; font-weight: bold")
        _expand_style(el)
        assert el.get("font-size") == "14px"
        assert el.get("font-family") == "Arial"
        assert el.get("font-weight") == "bold"

    def test_expands_opacity(self):
        el = _parse_el("rect", style="opacity: 0.5")
        _expand_style(el)
        assert el.get("opacity") == "0.5"

    def test_expands_transform(self):
        el = _parse_el("g", style="transform: translate(10px, 20px)")
        _expand_style(el)
        assert el.get("transform") == "translate(10px, 20px)"


class TestFixSelfClosing:
    def test_fixes_self_closing_rect(self):
        root = etree.fromstring(b'<svg xmlns="http://www.w3.org/2000/svg"><rect></rect></svg>')
        _fix_self_closing_lxml(root)
        rect = _find(root, "rect")
        assert rect is not None
        assert rect.text is None

    def test_preserves_child_elements(self):
        root = etree.fromstring(
            b'<svg xmlns="http://www.w3.org/2000/svg"><rect><title>box</title></rect></svg>'
        )
        _fix_self_closing_lxml(root)
        rect = _find(root, "rect")
        assert rect is not None
        title = _find(rect, "title")
        assert title is not None

    def test_skips_non_self_closing(self):
        root = etree.fromstring(b'<svg xmlns="http://www.w3.org/2000/svg"><g></g></svg>')
        _fix_self_closing_lxml(root)
        g = _find(root, "g")
        assert g is not None


class TestInferViewbox:
    def test_already_has_viewbox(self):
        root = etree.fromstring(b'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"></svg>')
        _infer_viewbox(root)
        assert root.get("viewBox") == "0 0 100 100"

    def test_infers_from_width_height(self):
        root = etree.fromstring(b'<svg xmlns="http://www.w3.org/2000/svg" width="400" height="300"></svg>')
        _infer_viewbox(root)
        assert root.get("viewBox") == "0 0 400 300"

    def test_strips_units(self):
        root = etree.fromstring(b'<svg xmlns="http://www.w3.org/2000/svg" width="400px" height="300px"></svg>')
        _infer_viewbox(root)
        assert root.get("viewBox") == "0 0 400 300"

    def test_no_width_height(self):
        root = etree.fromstring(b'<svg xmlns="http://www.w3.org/2000/svg"></svg>')
        _infer_viewbox(root)
        assert root.get("viewBox") is None


class TestStripUnwanted:
    def test_removes_style_element(self):
        svg = f'{_SVG_OPEN}><style>.cls{{fill:red}}</style><rect/></svg>'
        root = etree.fromstring(svg.encode())
        _strip_unwanted(root)
        assert len(_findall(root, "style")) == 0
        assert _find(root, "rect") is not None

    def test_removes_script_element(self):
        svg = f"{_SVG_OPEN}><script>alert(1)</script><rect/></svg>"
        root = etree.fromstring(svg.encode())
        _strip_unwanted(root)
        assert len(_findall(root, "script")) == 0

    def test_removes_nested_style(self):
        svg = f'{_SVG_OPEN}><g><style>.x{{fill:blue}}</style></g></svg>'
        root = etree.fromstring(svg.encode())
        _strip_unwanted(root)
        assert len(_findall(root, "style")) == 0

    def test_preserves_other_elements(self):
        svg = f"{_SVG_OPEN}><defs></defs><g></g><rect/></svg>"
        root = etree.fromstring(svg.encode())
        _strip_unwanted(root)
        assert _find(root, "defs") is not None
        assert _find(root, "g") is not None
        assert _find(root, "rect") is not None


class TestWalkExpand:
    def test_expands_nested_styles(self):
        svg = f'{_SVG_OPEN}><g style="opacity: 0.5"><rect style="fill: red"/></g></svg>'
        root = etree.fromstring(svg.encode())
        _walk_expand(root)
        g = _find(root, "g")
        rect = _find(root, "rect")
        assert g is not None
        assert g.get("opacity") == "0.5"
        assert rect is not None
        assert rect.get("fill") == "red"


class TestSanitize:
    def test_basic_svg(self):
        svg = f'{_SVG_OPEN} viewBox="0 0 100 100"><rect x="0" y="0" width="50" height="50" fill="red"/></svg>'
        root = sanitize(svg)
        assert root is not None
        rect = _find(root, "rect")
        assert rect is not None

    def test_empty_input(self):
        with pytest.raises(SVGCompileError):
            sanitize("")

    def test_whitespace_only(self):
        with pytest.raises(SVGCompileError):
            sanitize("   ")

    def test_invalid_xml_recovers(self):
        svg = f"{_SVG_OPEN}><rect></svg>"
        root = sanitize(svg)
        assert root is not None

    def test_strips_style_element(self):
        svg = f'{_SVG_OPEN}><style>.cls{{fill:red}}</style><rect class="cls"/></svg>'
        root = sanitize(svg)
        assert len(_findall(root, "style")) == 0

    def test_expands_inline_style(self):
        svg = f'{_SVG_OPEN}><rect style="fill: blue"/></svg>'
        root = sanitize(svg)
        rect = _find(root, "rect")
        assert rect is not None
        assert rect.get("fill") == "blue"

    def test_infers_viewbox(self):
        svg = f'{_SVG_OPEN} width="400" height="300"><rect/></svg>'
        root = sanitize(svg)
        assert root.get("viewBox") == "0 0 400 300"

    def test_strips_script(self):
        svg = f"{_SVG_OPEN}><script>alert(1)</script><rect/></svg>"
        root = sanitize(svg)
        assert len(_findall(root, "script")) == 0

    def test_full_compilation(self):
        """End-to-end: sanitized SVG compiles successfully."""
        from pptx_designer.compiler import SVGCompiler
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        svg = f'{_SVG_OPEN} viewBox="0 0 400 300"><circle cx="200" cy="150" r="100" fill="#4472C4"/></svg>'
        result = SVGCompiler().compile(svg, slide, (1, 1, 8, 6))
        assert result.shape_count >= 1
