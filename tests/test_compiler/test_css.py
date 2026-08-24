"""Unit tests for the CSS style block parser."""

from __future__ import annotations

from lxml import etree

from pptx_designer.compiler._css import (
    _match_selector,
    _parse_css_block,
    apply_css_blocks,
)
from pptx_designer.compiler._sanitizer import sanitize


SVG_NS = "http://www.w3.org/2000/svg"
SVG = f"{{{SVG_NS}}}"


def _find(el: etree._Element, tag: str) -> etree._Element | None:
    for desc in el.iter():
        local = desc.tag.split("}")[-1] if "}" in desc.tag else desc.tag
        if local == tag:
            return desc
    return None


def _findall(el: etree._Element, tag: str) -> list[etree._Element]:
    results = []
    for desc in el.iter():
        local = desc.tag.split("}")[-1] if "}" in desc.tag else desc.tag
        if local == tag:
            results.append(desc)
    return results


class TestParseCssBlock:
    def test_single_selector_single_prop(self):
        rules = _parse_css_block(".box { fill: red; }")
        assert len(rules) == 1
        assert rules[0] == (".box", {"fill": "red"})

    def test_single_selector_multiple_props(self):
        rules = _parse_css_block(".box { fill: red; stroke: blue; stroke-width: 2; }")
        assert len(rules) == 1
        assert rules[0] == (".box", {"fill": "red", "stroke": "blue", "stroke-width": "2"})

    def test_multiple_selectors(self):
        rules = _parse_css_block(".a { fill: red; } .b { fill: blue; }")
        assert len(rules) == 2
        assert rules[0][0] == ".a"
        assert rules[1][0] == ".b"

    def test_comma_separated_selectors(self):
        rules = _parse_css_block(".a, .b { fill: red; }")
        assert len(rules) == 2
        selectors = {r[0] for r in rules}
        assert selectors == {".a", ".b"}

    def test_id_selector(self):
        rules = _parse_css_block("#icon { opacity: 0.5; }")
        assert rules[0] == ("#icon", {"opacity": "0.5"})

    def test_tag_selector(self):
        rules = _parse_css_block("rect { stroke-width: 2; }")
        assert rules[0] == ("rect", {"stroke-width": "2"})

    def test_empty_block(self):
        rules = _parse_css_block(".box { }")
        assert rules == []

    def test_comments_stripped(self):
        rules = _parse_css_block("/* comment */ .x { fill: red; } /* another */")
        assert len(rules) == 1
        assert rules[0][1] == {"fill": "red"}

    def test_multiline(self):
        css = """
        .cls-1 {
            fill: #333;
            stroke: #000;
        }
        """
        rules = _parse_css_block(css)
        assert len(rules) == 1
        assert rules[0][1]["fill"] == "#333"
        assert rules[0][1]["stroke"] == "#000"


class TestMatchSelector:
    def _make_tree(self) -> etree._Element:
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<rect id="bg" class="box primary" x="0" y="0" width="400" height="300"/>'
            '<circle id="dot" class="accent" cx="200" cy="150" r="50"/>'
            '<g class="group"><path id="arrow" d="M0,0 L10,10"/></g>'
            '<text class="label">Hello</text>'
            "</svg>"
        )
        return etree.fromstring(svg.encode())

    def test_tag_selector(self):
        root = self._make_tree()
        matched = _match_selector(root, "rect")
        assert len(matched) == 1
        assert matched.pop().get("id") == "bg"

    def test_class_selector(self):
        root = self._make_tree()
        matched = _match_selector(root, ".box")
        assert len(matched) == 1
        assert matched.pop().get("id") == "bg"

    def test_class_selector_multiple_matches(self):
        root = self._make_tree()
        matched = _match_selector(root, ".primary")
        assert len(matched) == 1

    def test_id_selector(self):
        root = self._make_tree()
        matched = _match_selector(root, "#dot")
        assert len(matched) == 1
        assert matched.pop().tag.endswith("circle")

    def test_class_on_nested_element(self):
        root = self._make_tree()
        matched = _match_selector(root, ".accent")
        assert len(matched) == 1
        assert matched.pop().get("id") == "dot"

    def test_no_match(self):
        root = self._make_tree()
        matched = _match_selector(root, ".nonexistent")
        assert len(matched) == 0


class TestApplyCssBlocks:
    def test_class_selector_applies_fill(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>.box { fill: #3366FF; }</style>'
            '<rect class="box" x="10" y="10" width="100" height="50"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        assert rect is not None
        assert rect.get("fill") == "#3366FF"
        # <style> should be removed
        assert _find(root, "style") is None

    def test_id_selector(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>#icon { stroke: #000; stroke-width: 2; }</style>'
            '<circle id="icon" cx="200" cy="150" r="50"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        circle = _find(root, "circle")
        assert circle.get("stroke") == "#000"
        assert circle.get("stroke-width") == "2"

    def test_tag_selector(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>rect { stroke-width: 3; }</style>'
            '<rect x="10" y="10" width="100" height="50" fill="red"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        assert rect.get("stroke-width") == "3"
        assert rect.get("fill") == "red"  # not overridden

    def test_stylesheet_overrides_presentation_attribute(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>.box { fill: red; }</style>'
            '<rect class="box" fill="green" x="10" y="10" width="100" height="50"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        # SVG presentation attributes have lower cascade precedence than a
        # stylesheet rule. Actual inline ``style=`` is tested separately.
        assert rect.get("fill") == "red"

    def test_css_var_definition(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>:root { --primary: #1D78FA; }</style>'
            '<rect fill="var(--primary)" x="10" y="10" width="100" height="50"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        assert rect.get("fill") == "#1D78FA"

    def test_later_rule_wins_with_same_specificity(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<style>.a { fill: red; } .b { fill: blue; }</style>'
            '<rect class="a b" width="10" height="10"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        assert _find(root, "rect").get("fill") == "blue"

    def test_class_overrides_tag_and_presentation_attribute(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<style>rect { fill: red; } .accent { fill: blue; }</style>'
            '<rect class="accent" fill="green" width="10" height="10"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        assert _find(root, "rect").get("fill") == "blue"

    def test_inline_style_and_important_precedence(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<style>.accent { fill: blue !important; }</style>'
            '<rect class="accent" style="fill: green" width="10" height="10"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        assert _find(root, "rect").get("fill") == "blue"

    def test_inline_important_overrides_stylesheet_important(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<style>.accent { fill: blue !important; }</style>'
            '<rect class="accent" style="fill: green !important" width="10" height="10"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        assert _find(root, "rect").get("fill") == "green"

    def test_opacity_is_not_inherited(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<g style="opacity: 0.3"><rect width="10" height="10"/></g>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        assert _find(root, "rect").get("opacity") is None

    def test_inheritance_fill(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>.parent { fill: blue; }</style>'
            '<g class="parent">'
            '<rect x="10" y="10" width="100" height="50"/>'
            "</g>"
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        assert rect is not None
        assert rect.get("fill") == "blue"

    def test_group_opacity_is_not_inherited(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>.faded { opacity: 0.3; }</style>'
            '<g class="faded">'
            '<rect x="10" y="10" width="100" height="50"/>'
            '<circle cx="200" cy="150" r="50"/>'
            "</g>"
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        group = _find(root, "g")
        rect = _find(root, "rect")
        circle = _find(root, "circle")
        assert group.get("opacity") == "0.3"
        assert rect.get("opacity") is None
        assert circle.get("opacity") is None

    def test_child_overrides_inherited(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>.parent { fill: blue; } .child { fill: red; }</style>'
            '<g class="parent">'
            '<rect class="child" x="10" y="10" width="100" height="50"/>'
            "</g>"
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        assert rect.get("fill") == "red"  # child CSS wins

    def test_no_style_element(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<rect fill="red" x="10" y="10" width="100" height="50"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        assert rect.get("fill") == "red"

    def test_multiple_style_blocks(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>.a { fill: red; }</style>'
            '<style>.b { stroke: blue; }</style>'
            '<rect class="a" x="10" y="10" width="100" height="50"/>'
            '<circle class="b" cx="200" cy="150" r="50"/>'
            "</svg>"
        )
        root = etree.fromstring(svg.encode())
        apply_css_blocks(root)
        rect = _find(root, "rect")
        circle = _find(root, "circle")
        assert rect is not None
        assert circle is not None
        assert rect.get("fill") == "red"
        assert circle.get("stroke") == "blue"


class TestSanitizeWithCSS:
    """Test CSS parsing integrated into the full sanitize() pipeline."""

    def test_class_based_svg_compiles(self):
        from pptx_designer.compiler import SVGCompiler
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>.box { fill: #3366FF; stroke: #000; stroke-width: 2; }</style>'
            '<rect class="box" x="50" y="50" width="300" height="200"/>'
            "</svg>"
        )
        result = SVGCompiler().compile(svg, slide, (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_llm_typical_svg(self):
        """Typical LLM-generated SVG with class-based styling."""
        from pptx_designer.compiler import SVGCompiler
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            '<style>'
            ".background { fill: #1a1a2e; }"
            ".card { fill: #16213e; stroke: #0f3460; stroke-width: 1; }"
            ".title { fill: #e94560; font-size: 18px; font-weight: bold; }"
            ".text { fill: #ffffff; font-size: 12px; }"
            "</style>"
            '<rect class="background" x="0" y="0" width="400" height="300"/>'
            '<rect class="card" x="20" y="20" width="360" height="260" rx="8"/>'
            '<text class="title" x="200" y="60" text-anchor="middle">Dashboard</text>'
            '<text class="text" x="40" y="100">Revenue: $1.2M</text>'
            '<text class="text" x="40" y="130">Users: 12,400</text>'
            "</svg>"
        )
        result = SVGCompiler().compile(svg, slide, (1, 1, 8, 6))
        assert result.shape_count >= 4  # 2 rects + 2+ texts

    def test_empty_style(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">'
            "<style></style>"
            '<rect fill="red" x="10" y="10" width="100" height="50"/>'
            "</svg>"
        )
        root = sanitize(svg)
        rect = _find(root, "rect")
        assert rect is not None
        assert rect.get("fill") == "red"
