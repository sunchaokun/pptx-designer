"""Contract tests for Build-owned atomic composition under VI constraints."""

from copy import deepcopy

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_designer.enterprise.vi_adapter import VITemplateAdapter
from pptx_designer.qa.structural import StructuralQA


def _context() -> dict:
    return {
        "colors": {"background": "#FFFFFF"},
        "atom_styles": {
            "type.claim": {
                "kind": "text",
                "font_size": 28,
                "bold": True,
                "color": "text_dark",
                "font_name": "Microsoft YaHei",
            },
            "surface.comparison": {"kind": "rect", "fill": "#E8F0EA"},
        },
        "visual_grammar": {
            "allowed_atom_kinds": ["text", "rect"],
            "safe_area": {"left": 0.5, "top": 0.4, "width": 12.3, "height": 6.7},
            "forbidden_zones": [{"left": 0, "top": 7.1, "width": 13.333, "height": 0.4}],
            "min_font_size": 16,
        },
        "fixed_bases": {
            "base.content": {
                "safe_to_clone": True,
                "fixed_objects": ["footer.logo"],
                "removed_objects": [],
                "unsupported_objects": [],
            }
        },
    }


def _atomic_plan() -> dict:
    return {
        "base_id": "base.content",
        "content_model": {
            "claim": "Choose the operating model deliberately.",
            "relations": [{"id": "contrast:centralized:distributed", "type": "contrasts"}],
        },
        "atoms": [
            {
                "id": "claim",
                "kind": "text",
                "geometry": {"left": 0.7, "top": 0.6, "width": 8.2, "height": 0.7},
                "style_ref": "type.claim",
                "data": "Choose the operating model deliberately.",
                "z_index": 20,
            },
            {
                "id": "left-panel",
                "kind": "rect",
                "geometry": {"left": 0.7, "top": 1.7, "width": 5.6, "height": 4.6},
                "style_ref": "surface.comparison",
                "data": None,
                "z_index": 1,
            },
            {
                "id": "right-panel",
                "kind": "rect",
                "geometry": {"left": 6.7, "top": 1.7, "width": 5.6, "height": 4.6},
                "style_ref": "surface.comparison",
                "data": None,
                "z_index": 1,
            },
        ],
        "relation_bindings": [
            {"relation_id": "contrast:centralized:distributed", "atom_ids": ["left-panel", "right-panel"]}
        ],
    }


def test_adapter_preserves_build_owned_atoms_and_exact_geometry():
    context = _context()
    original = deepcopy(context)
    plan = _atomic_plan()

    spec = VITemplateAdapter(context).compile_atomic(page_role="content", atomic_build_plan=plan)

    assert spec["atomic_build"] is True
    assert spec["content_model"] == plan["content_model"]
    assert spec["relation_bindings"] == plan["relation_bindings"]
    assert [item["atom_id"] for item in spec["components"]] == ["left-panel", "right-panel", "claim"]
    claim = next(item for item in spec["components"] if item["atom_id"] == "claim")
    assert claim["recipe"]["bounds"] == plan["atoms"][0]["geometry"]
    assert claim["recipe"]["font_size"] == 28
    assert spec["fixed_base"]["id"] == "base.content"
    assert context == original


def test_atomic_build_spec_renders_build_geometry_not_template_component_bounds():
    adapter = VITemplateAdapter(_context())
    spec = adapter.compile_atomic(page_role="content", atomic_build_plan=_atomic_plan())
    presentation = Presentation()

    adapter.render(spec, presentation)

    claim = next(shape for shape in presentation.slides[0].shapes if shape.has_text_frame and shape.text)
    assert round(claim.left / Inches(1), 2) == 0.7
    assert round(claim.top / Inches(1), 2) == 0.6
    assert claim.text_frame.paragraphs[0].runs[0].font.name == "Microsoft YaHei"


def test_atomic_build_passes_content_relations_through_without_interpreting_them():
    plan = _atomic_plan()
    plan["relation_bindings"] = []

    spec = VITemplateAdapter(_context()).compile_atomic(page_role="content", atomic_build_plan=plan)

    assert spec["content_model"] == plan["content_model"]
    assert spec["relation_bindings"] == []


def test_atomic_build_rejects_vi_constraint_instead_of_rearranging_atoms():
    plan = _atomic_plan()
    plan["atoms"][0]["geometry"]["top"] = 6.6

    with pytest.raises(ValueError, match="atom_outside_safe_area:claim"):
        VITemplateAdapter(_context()).compile_atomic(page_role="content", atomic_build_plan=plan)


def test_atomic_build_does_not_route_content_type_or_template_archetype():
    context = _context()
    context["archetypes"] = [
        {
            "id": "legacy-gallery",
            "page_role": "content",
            "content_types": ["gallery"],
            "component_slots": ["legacy-photo"],
        }
    ]

    spec = VITemplateAdapter(context).compile_atomic(page_role="content", atomic_build_plan=_atomic_plan())

    assert "archetype_id" not in spec
    assert "variant_id" not in spec
    assert all(item["atom_id"] != "legacy-photo" for item in spec["components"])


def test_qa_rejects_an_atomic_plan_with_an_uncovered_relation(tmp_path):
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    output = tmp_path / "atomic-qa.pptx"
    presentation.save(output)
    malformed_plan = {
        "atomic_build": True,
        "page_role": "content",
        "render_strategy": "components",
        "content_model": {"relations": [{"id": "supports:proof:claim"}]},
        "components": [{"atom_id": "claim"}],
        "relation_bindings": [],
    }

    report = StructuralQA().check(output, vi_plans=[malformed_plan])

    assert "relation_uncovered" in {issue.kind for issue in report.fatal}
