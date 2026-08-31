from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_designer.qa import run_structural_qa
from pptx_designer.qa.visual_baseline import compare, create


def _pptx(path: Path, *, out_of_bounds: bool = False) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(12.9 if out_of_bounds else 1)
    box = slide.shapes.add_textbox(left, Inches(1), Inches(1), Inches(0.5))
    box.text = "Readable text"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    prs.save(path)


def test_structural_qa_passes_reopenable_editable_deck(tmp_path):
    path = tmp_path / "deck.pptx"
    _pptx(path)
    report = run_structural_qa(path, expected_slides=1)
    assert report.status == "pass"
    assert report.shape_count == 1
    assert report.editable_ratio == 1


def test_structural_qa_reports_slide_count_and_bounds(tmp_path):
    path = tmp_path / "deck.pptx"
    _pptx(path, out_of_bounds=True)
    report = run_structural_qa(path, expected_slides=2)
    assert report.status == "fail"
    assert any(issue.kind == "slide_count" for issue in report.fatal)
    assert any(issue.kind == "out_of_bounds" for issue in report.fatal)


def test_named_background_decoration_may_be_intentionally_clipped(tmp_path):
    path = tmp_path / "clipped-decoration.pptx"
    _pptx(path, out_of_bounds=True)
    presentation = Presentation(path)
    presentation.slides[0].shapes[0].name = "Background Decoration 1"
    presentation.save(path)

    report = run_structural_qa(path, expected_slides=1)

    assert report.status == "pass"


def test_visual_baseline_create_and_compare(tmp_path):
    rendered = tmp_path / "rendered"
    baseline = tmp_path / "baseline"
    rendered.mkdir()
    Image.new("RGBA", (4, 4), "#123456").save(rendered / "slide01.png")
    assert create(rendered, baseline)["status"] == "created"
    assert compare(rendered, baseline)["status"] == "pass"
    Image.new("RGBA", (4, 4), "#654321").save(rendered / "slide01.png")
    assert compare(rendered, baseline, threshold=0.1)["status"] == "fail"


def test_structural_qa_can_report_text_media_overlap(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(2), Inches(1))
    box.text = "Text over image"
    # A later picture layer can cover the text and must be reported.
    slide.shapes.add_picture(str(_image(tmp_path, "#123456")), Inches(1), Inches(1), Inches(4), Inches(3))
    path = tmp_path / "overlap.pptx"
    prs.save(path)

    report = run_structural_qa(path, check_overlaps=True)

    assert any(issue.kind == "text_media_overlap" for issue in report.warnings)


def _image(tmp_path, color):
    path = tmp_path / f"{color[1:]}.png"
    Image.new("RGB", (80, 60), color).save(path)
    return path
