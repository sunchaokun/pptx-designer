"""Acceptance tests for end-to-end theme application.

These tests describe the target behavior from the theme integration design:
theme changes must be visible in the editable PPTX, rather than only existing
as an intermediate color dictionary.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation as PptxPresentation

from pptx_designer import Presentation, set_slide_theme
from pptx_designer.core.pipeline import generate_ppt
from pptx_designer.diagrams.diagram_style import DiagramStyle
from pptx_designer.renderer.theme import ThemeComposer
from pptx_designer.tools.charts import bar_chart
from pptx_designer.tools.layout import page_header
from pptx_designer.tools.shapes import rect
from pptx_designer.tools.text import text

PAGES = [
    {
        "goal": "hook",
        "title": "Theme application test",
        "subtitle": "The same content rendered with different visual systems",
        "bullets": ["A measurable result"],
    },
    {
        "goal": "data",
        "title": "Evidence",
        "subtitle": "Theme differences should survive into editable objects",
        "bullets": ["42%: Adoption", "3.2x: Growth"],
    },
]


def _render_theme(tmp_path: Path, style: str) -> tuple[dict, PptxPresentation]:
    output = tmp_path / f"{style}.pptx"
    result = generate_ppt(
        content={"pages": PAGES},
        style=style,
        style_seed=17,
        output=str(output),
    )
    return result, PptxPresentation(str(output))


def _shape_fill_hex(shape) -> str | None:
    """Return a solid RGB fill when the shape has one."""
    try:
        if not shape.fill.fore_color.rgb:
            return None
        return str(shape.fill.fore_color.rgb)
    except (AttributeError, TypeError, ValueError):
        return None


def _text_font_names(prs: PptxPresentation) -> set[str]:
    names: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        names.add(run.font.name)
    return names


def test_same_content_with_two_themes_changes_editable_visual_output(tmp_path):
    """A theme must change more than the intermediate theme metadata."""
    dark_result, dark_prs = _render_theme(tmp_path, "dark-tech")
    elegant_result, elegant_prs = _render_theme(tmp_path, "warm-elegant")

    dark_fills = {_shape_fill_hex(shape) for slide in dark_prs.slides for shape in slide.shapes}
    elegant_fills = {_shape_fill_hex(shape) for slide in elegant_prs.slides for shape in slide.shapes}

    assert dark_result["theme_atoms"] != elegant_result["theme_atoms"]
    assert dark_fills != elegant_fills
    assert dark_result["slide_count"] == elegant_result["slide_count"] == len(PAGES)


def test_resolved_theme_typography_reaches_editable_text(tmp_path):
    """The selected heading/body fonts must be present in the generated PPTX."""
    theme = ThemeComposer().compose(style="warm-elegant", seed=17)
    _, prs = _render_theme(tmp_path, "warm-elegant")

    font_names = _text_font_names(prs)
    assert theme["typography"]["heading"] in font_names
    assert theme["typography"]["body"] in font_names


def test_theme_result_exposes_complete_traceable_application_context():
    theme = ThemeComposer().compose(style="dark-tech", seed=17)

    assert {"colors", "typography", "decoration", "layout_variant"} <= theme.keys()
    assert "source" in theme
    assert "semantic_roles" in theme
    assert theme["source"]["seed"] == 17


def test_theme_composition_is_reproducible_with_explicit_atoms():
    composer = ThemeComposer()
    kwargs = {
        "style": "warm-elegant",
        "palette": "golden-luxury",
        "fonts": "serif-editorial",
        "decoration": "gold-trim",
        "layout": "centered",
        "seed": 17,
    }

    assert composer.compose(**kwargs) == composer.compose(**kwargs)


def test_explicit_theme_atoms_override_preset_atoms():
    theme = ThemeComposer().compose(
        style="dark-tech",
        palette="golden-luxury",
        fonts="serif-editorial",
        decoration="gold-trim",
        layout="centered",
        seed=17,
    )

    assert theme["atoms"]["palette"] == "golden-luxury"
    assert theme["atoms"]["fonts"] == "serif-editorial"
    assert theme["atoms"]["decoration"] == "gold-trim"
    assert theme["atoms"]["layout"] == "centered"


def test_semantic_solution_color_does_not_revert_to_fixed_green(tmp_path):
    output = tmp_path / "solution-theme.pptx"
    generate_ppt(
        content={
            "pages": [
                {
                    "goal": "solution",
                    "title": "Solution",
                    "subtitle": "Theme semantic color test",
                    "bullets": ["A solution item"],
                }
            ]
        },
        style="warm-elegant",
        style_seed=17,
        output=str(output),
    )
    prs = PptxPresentation(str(output))
    fills = {_shape_fill_hex(shape) for shape in prs.slides[0].shapes}

    assert "22C55E" not in fills


def test_generate_ppt_uses_a_locked_resolved_theme_without_discovery(tmp_path, monkeypatch):
    locked_theme = ThemeComposer().compose(style="warm-elegant", seed=17)

    def fail_if_called(*args, **kwargs):
        raise AssertionError("Theme discovery must not run for a locked theme")

    monkeypatch.setattr(ThemeComposer, "compose", fail_if_called)
    output = tmp_path / "locked-theme.pptx"
    result = generate_ppt(content={"pages": PAGES}, theme=locked_theme, output=str(output))

    assert output.exists()
    assert result["theme_context"] == locked_theme


def test_theme_application_reports_not_yet_consumed_theme_fields(tmp_path):
    result, _ = _render_theme(tmp_path, "dark-tech")
    not_applied = {item["field"] for item in result["theme_application"]["not_applied"]}

    assert {"decoration", "layout_variant", "text_effect_preset", "image_effect"} <= not_applied


def test_build_mode_helpers_inherit_presentation_theme(tmp_path):
    theme = ThemeComposer().compose(style="warm-elegant", seed=17)
    prs = Presentation(theme=theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rect(slide, 0, 0, 13.333, 7.5, "background")
    page_header(slide, "Inherited heading", "Inherited body")
    text(slide, 1, 2, 8, 0.5, "Inherited text")
    bar_chart(slide, 2, 3, [("Adoption", 0.7, "70%")])

    output = tmp_path / "build-theme.pptx"
    prs.save(output)
    generated = PptxPresentation(str(output))

    assert _shape_fill_hex(generated.slides[0].shapes[0]) == theme["semantic_roles"]["background"].lstrip("#")
    assert theme["typography"]["heading"] in _text_font_names(generated)
    assert theme["typography"]["body"] in _text_font_names(generated)


def test_slide_theme_and_explicit_values_override_presentation_defaults():
    warm = ThemeComposer().compose(style="warm-elegant", seed=17)
    dark = ThemeComposer().compose(style="dark-tech", seed=17)
    prs = Presentation(theme=warm)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_theme(slide, dark)

    inherited = text(slide, 1, 1, 5, 0.5, "Slide override")
    explicit = text(slide, 1, 2, 5, 0.5, "Element override", font_name="Arial", C={"text_body": "#123456"})

    inherited_run = inherited.text_frame.paragraphs[0].runs[0]
    explicit_run = explicit.text_frame.paragraphs[0].runs[0]
    assert inherited_run.font.name == dark["typography"]["body"]
    assert str(inherited_run.font.color.rgb) == dark["semantic_roles"]["ink"].lstrip("#")
    assert explicit_run.font.name == "Arial"
    assert str(explicit_run.font.color.rgb) == "123456"


def test_diagram_style_uses_resolved_semantic_roles():
    theme = ThemeComposer().compose(style="warm-elegant", seed=17)
    style = DiagramStyle.from_theme(theme)

    assert style.resolve_color("primary") == theme["semantic_roles"]["data-series-1"]
    assert style.resolve_color("foreground") == theme["semantic_roles"]["ink"]
    assert style.resolve_color("muted") == theme["semantic_roles"]["surface"]
