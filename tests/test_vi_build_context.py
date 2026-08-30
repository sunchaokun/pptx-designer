"""Contract tests for template-derived VI Build contexts.

These tests intentionally describe Build's responsibilities before the
extractor is expanded.  A context is only useful when Build either applies a
rule or reports why it cannot do so.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from pptx_designer import Presentation
from pptx_designer.enterprise.brand import BrandSpec
from pptx_designer.enterprise.design_dna_extractor import extract_design_context, extract_design_dna
from pptx_designer.enterprise.vi_context import (
    VIBuildSession,
    design_context_from_brand_spec,
    normalize_design_context,
)


def _photo_archetype_context() -> dict:
    return {
        "source": {"kind": "template", "confidence": 1.0},
        "semantic_roles": {"background": "#FFFFFF", "ink": "#111111"},
        "typography": {"heading": "Aptos Display", "body": "Aptos"},
        "assets": {
            "image_grammar": {
                "required": True,
                "subjects": ["botanical"],
                "crop": {"mode": "cover"},
                "min_area_ratio": 0.35,
            }
        },
        "components": {
            "forest_photo_panel": {
                "kind": "photo_panel",
                "image_mode": "cover",
                "bounds": {"left": 7, "top": 0, "width": 6.333, "height": 7.5},
            }
        },
        "archetypes": [
            {
                "id": "text-photo-right",
                "reference_slide": 2,
                "required_assets": ["supporting_photo"],
                "permitted_components": ["forest_photo_panel"],
            }
        ],
        "content_slots": [
            {
                "id": "page_title",
                "max_chars": 24,
                "bounds": {"left": 0.7, "top": 1.2, "width": 5.5, "height": 1},
                "font_size": 32,
                "bold": True,
            }
        ],
        "locks": [{"field": "footer.text", "mode": "template-locked"}],
        "acceptance": {"must_coverage": ["image_present", "component_applied"]},
    }


def test_normalization_preserves_theme_contract_and_vi_fields():
    context = normalize_design_context(_photo_archetype_context())

    assert context["schema_version"] == "1.0"
    assert context["semantic_roles"]["background"] == "#FFFFFF"
    assert context["assets"]["image_grammar"]["required"] is True
    assert context["archetypes"][0]["id"] == "text-photo-right"
    assert context["diagnostics"]["warnings"] == []


def test_brand_spec_adapts_to_the_same_design_context():
    context = design_context_from_brand_spec(
        BrandSpec(
            source="brand_json",
            colors={"primary": "#115E32", "foreground": "#0B0C11", "background": "#FFFFFF"},
            fonts={"heading": "Proxima Nova", "body": "Aptos"},
            dark_mode=False,
        )
    )

    assert context["source"]["kind"] == "merged"
    assert context["semantic_roles"]["primary"] == "#115E32"
    assert context["semantic_roles"]["ink"] == "#0B0C11"
    assert context["typography"]["heading"] == "Proxima Nova"


def test_build_requires_asset_when_archetype_demands_photo():
    session = VIBuildSession(_photo_archetype_context())

    result = session.plan_page(
        "text-photo-right",
        components=["forest_photo_panel"],
        slot_values={"page_title": "春季系列发布"},
    )

    assert result["status"] == "NEEDS_ASSET"
    assert result["asset_plan"]["missing"] == ["supporting_photo"]
    assert "image_present" in result["acceptance"]["blocked"]


def test_build_applies_asset_component_slot_and_acceptance_rules(tmp_path: Path):
    photo = tmp_path / "botanical.png"
    photo.write_bytes(b"placeholder image path is sufficient for planning")
    session = VIBuildSession(_photo_archetype_context(), assets={"supporting_photo": str(photo)})

    result = session.plan_page(
        "text-photo-right",
        components=["forest_photo_panel"],
        slot_values={"page_title": "春季系列发布"},
    )

    assert result["status"] == "READY"
    assert result["asset_plan"]["resolved"]["supporting_photo"] == str(photo)
    assert result["slot_bindings"] == [{"id": "page_title", "status": "bound"}]
    assert result["acceptance"]["passed"] == ["image_present", "component_applied"]


def test_build_rejects_unknown_or_overflowing_content_slot():
    session = VIBuildSession(_photo_archetype_context(), assets={"supporting_photo": "photo.png"})

    with pytest.raises(ValueError, match="unknown content slot"):
        session.plan_page("text-photo-right", slot_values={"unknown": "x"})
    with pytest.raises(ValueError, match="exceeds max_chars"):
        session.plan_page("text-photo-right", slot_values={"page_title": "x" * 25})


def test_build_rejects_bound_slot_without_render_geometry(tmp_path: Path):
    context = _photo_archetype_context()
    context["content_slots"][0].pop("bounds")
    photo = tmp_path / "botanical.png"
    Image.new("RGB", (800, 500), "#115E32").save(photo)
    session = VIBuildSession(context, assets={"supporting_photo": str(photo)})
    presentation = Presentation(theme=context)

    with pytest.raises(ValueError, match="has no render bounds"):
        session.render_page(
            presentation,
            "text-photo-right",
            components=["forest_photo_panel"],
            slot_values={"page_title": "春季系列发布"},
        )


def test_template_lock_blocks_unapproved_override():
    session = VIBuildSession(_photo_archetype_context())

    with pytest.raises(PermissionError, match="template-locked"):
        session.validate_overrides({"footer.text": "new footer"})

    approved = session.validate_overrides(
        {"footer.text": "new footer"}, allow_template_override=True
    )
    assert approved["overrides"][0]["status"] == "approved_override"


def test_vi_build_renders_required_photo_component_or_blocks_without_asset(tmp_path: Path):
    photo = tmp_path / "botanical.png"
    Image.new("RGB", (800, 500), "#115E32").save(photo)
    context = _photo_archetype_context()
    presentation = Presentation(theme=context)

    blocked = VIBuildSession(context).render_page(
        presentation,
        "text-photo-right",
        components=["forest_photo_panel"],
        slot_values={"page_title": "春季系列发布"},
    )
    rendered = VIBuildSession(context, assets={"supporting_photo": str(photo)}).render_page(
        presentation,
        "text-photo-right",
        components=["forest_photo_panel"],
        slot_values={"page_title": "春季系列发布"},
    )

    assert blocked["status"] == "NEEDS_ASSET"
    assert blocked["slide"] is None
    assert rendered["status"] == "READY"
    assert rendered["slide"] is presentation.slides[0]
    assert rendered["design_application"]["applied_to"] == ["forest_photo_panel", "page_title"]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in rendered["slide"].shapes)
    assert any(shape.has_text_frame and shape.text == "春季系列发布" for shape in rendered["slide"].shapes)


def _make_template_with_deterministic_visual_evidence(tmp_path: Path) -> Path:
    image_path = tmp_path / "botanical.png"
    Image.new("RGB", (800, 500), "#164A3B").save(image_path)

    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(0), Inches(6.333), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(22, 74, 59)
    panel.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(5.5), Inches(1))
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "自然的系列"
    run.font.name = "Aptos Display"
    run.font.size = Pt(32)
    run.font.color.rgb = RGBColor(18, 18, 18)
    rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(2.4), Inches(6.0), Inches(2.4))
    rule.line.color.rgb = RGBColor(255, 255, 255)
    slide.shapes.add_picture(str(image_path), Inches(7.2), Inches(0.5), Inches(5.6), Inches(6.5))
    template_path = tmp_path / "template.pptx"
    prs.save(template_path)
    return template_path


def test_extractor_emits_stable_unified_context_and_keeps_legacy_projection(tmp_path: Path):
    template_path = _make_template_with_deterministic_visual_evidence(tmp_path)

    first = extract_design_context(str(template_path))
    second = extract_design_context(str(template_path))
    legacy = extract_design_dna(str(template_path))

    assert first["source"]["kind"] == "template"
    assert first["source"]["template_fingerprint"] == second["source"]["template_fingerprint"]
    assert first["typography"]["heading"] == "Aptos Display"
    assert first["semantic_roles"]["primary"] == "#164A3B"
    assert first["assets"]["references"][0]["kind"] == "image"
    assert first["archetypes"][0]["required_assets"] == ["supporting_photo"]
    assert any(component["kind"] == "color_panel" for component in first["components"].values())
    assert any(component["kind"] == "rule" for component in first["components"].values())
    assert "#164A3B" in legacy["colors"].values()
    assert legacy["slides"][0]["text_zones"][0]["font_name"] == "Aptos Display"
    assert legacy["slides"][0]["images"]
    assert legacy["slides"][0]["shapes"]
