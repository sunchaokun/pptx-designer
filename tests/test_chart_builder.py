from __future__ import annotations

import pytest
from pptx import Presentation

from pptx_designer.renderer.chart_builder import ChartBuilder


def _slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _position():
    return {"x": 0, "y": 0, "width": 2, "height": 2}


@pytest.mark.parametrize("chart_type", ["radar", "radar_markers"])
def test_radar_chart_types_are_native(chart_type):
    chart = ChartBuilder().build(
        _slide(),
        {"type": chart_type, "categories": ["A", "B"], "series": [{"name": "S", "values": [1, 2]}]},
        _position(),
    )
    assert chart.chart_type.name == chart_type.upper()


@pytest.mark.parametrize("chart_type", ["scatter", "scatter_lines", "scatter_smooth"])
def test_scatter_chart_types_use_xy_data(chart_type):
    chart = ChartBuilder().build(
        _slide(),
        {"type": chart_type, "series": [{"name": "S", "values": [[1, 2], [2, 4]]}]},
        _position(),
    )
    assert chart.chart_type.name == {"scatter": "XY_SCATTER", "scatter_lines": "XY_SCATTER_LINES", "scatter_smooth": "XY_SCATTER_SMOOTH"}[chart_type]


def test_brand_fonts_is_accepted_by_precision_renderer_call_shape():
    chart = ChartBuilder().build(
        _slide(),
        {"type": "bar", "categories": ["A"], "series": [{"name": "S", "values": [1]}]},
        _position(),
        brand_fonts={"heading": "Arial", "body": "Arial"},
    )
    assert chart is not None


@pytest.mark.parametrize("chart_type", ["bar_3d", "pie_3d"])
def test_unsupported_3d_chart_types_fail_clearly(chart_type):
    with pytest.raises(ValueError, match="Unsupported native chart type"):
        ChartBuilder().build(
            _slide(),
            {"type": chart_type, "categories": ["A"], "series": [{"name": "S", "values": [1]}]},
            _position(),
        )
