"""Contract tests for the narrowed VI Adapter boundary."""

from pptx import Presentation

from pptx_designer.enterprise.vi_adapter import VITemplateAdapter


def _framework_context():
    return {
        "schema_version": "1.1",
        "framework_pages": [{"id": "cover", "role": "cover", "reference_slide": 1}],
        "content_slots": [
            {"id": "cover.title", "page_role": "cover", "target": {"object_id": "title-1"}},
        ],
        "archetypes": [],
    }


def test_content_pages_have_no_archetype_compiler_path():
    try:
        VITemplateAdapter(_framework_context()).compile(page_role="content", content={})
    except ValueError as exc:
        assert "compile_atomic" in str(exc)
    else:
        raise AssertionError("content pages must be Build-owned atomic plans")


def test_framework_pages_compile_only_confirmed_slots():
    spec = VITemplateAdapter(_framework_context()).compile(
        page_role="cover", page_goal="open", content={"slots": {"cover.title": "Approved title"}}
    )

    assert spec["render_strategy"] == "prototype"
    assert spec["reference_slide"] == 1
    assert spec["delivery_origin"] == "framework_rebound"
    assert spec["slot_instances"] == [{"slot_id": "cover.title", "value": "Approved title"}]


def test_framework_renderer_rebinds_editable_text(tmp_path):
    presentation = Presentation()
    source = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = source.shapes.add_textbox(1000000, 1000000, 4000000, 600000)
    title.name = "title-1"
    title.text = "Template sample"
    adapter = VITemplateAdapter(_framework_context())
    spec = adapter.compile(page_role="cover", content={"slots": {"cover.title": "Approved title"}})

    adapter.render(spec, presentation)
    output = tmp_path / "framework-output.pptx"
    presentation.save(output)
    reopened = Presentation(output)

    assert len(reopened.slides) == 2
    assert reopened.slides[1].shapes[0].text == "Approved title"
