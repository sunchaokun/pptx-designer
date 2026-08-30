"""End-to-end, offline regression tests for the public CLI and pipeline."""

from __future__ import annotations

import importlib
import sys

from pptx import Presentation as PptxPresentation

from pptx_designer import recommend_styles
from pptx_designer.core.pipeline import generate_ppt
from pptx_designer.renderer.theme import ThemeComposer

cli_module = importlib.import_module("pptx_designer.cli.main")


def test_cli_generates_presentation_with_forwarded_options(monkeypatch, capsys, tmp_path):
    captured = {}
    output = tmp_path / "cli-output.pptx"

    def fake_generate_ppt(**kwargs):
        captured.update(kwargs)
        return {"output_path": str(output), "page_count": 2, "strategy": "pitch"}

    monkeypatch.setattr(cli_module, "generate_ppt", fake_generate_ppt)
    monkeypatch.setattr(cli_module, "_load_dotenv", lambda: None)
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "pptx-designer",
            "AI roadmap",
            "--slides",
            "2",
            "--style",
            "editorial",
            "--decoration",
            "corner-accent",
            "--layout-variant",
            "centered",
            "--mood",
            "elegant",
            "--style-seed",
            "17",
            "--image-mode",
            "placeholder",
            "--no-auto-detect",
            "--output",
            str(output),
        ],
    )

    cli_module.main()

    assert captured["query"] == "AI roadmap"
    assert captured["slides"] == 2
    assert captured["style"] == "editorial"
    assert captured["decoration"] == "corner-accent"
    assert captured["layout_variant"] == "centered"
    assert captured["mood"] == "elegant"
    assert captured["style_seed"] == 17
    assert captured["auto_detect"] is False
    assert "Generated:" in capsys.readouterr().out


def test_cli_image_subcommand_prefers_search_shortcut(monkeypatch, capsys):
    captured = {}

    def fake_fetch_image(**kwargs):
        captured.update(kwargs)
        return {"path": "generated.png", "provider": None}

    monkeypatch.setattr(cli_module, "fetch_image", fake_fetch_image)
    monkeypatch.setattr(cli_module, "_load_dotenv", lambda: None)
    monkeypatch.setattr(sys, "argv", ["pptx-designer", "image", "abstract architecture", "--fetch-images"])

    cli_module.main()

    assert captured["mode"] == "search"
    assert captured["keywords"] == "abstract architecture"
    assert capsys.readouterr().out.strip() == "generated.png"


def test_cli_analyze_subcommand_outputs_json(monkeypatch, capsys):
    monkeypatch.setattr(cli_module, "extract_design_dna", lambda path: {"source": path, "slides": 3})
    monkeypatch.setattr(cli_module, "_load_dotenv", lambda: None)
    monkeypatch.setattr(sys, "argv", ["pptx-designer", "analyze", "source.pptx"])

    cli_module.main()

    output = capsys.readouterr().out
    assert '"source": "source.pptx"' in output
    assert '"slides": 3' in output


def test_pipeline_generates_editable_multi_goal_presentation(tmp_path):
    pages = [
        {"goal": goal, "title": goal.title(), "subtitle": "QA", "bullets": ["One", "Two", "Three"]}
        for goal in ("hook", "problem", "solution", "features", "data", "code", "exercise", "overview", "content")
    ]
    output = tmp_path / "multi-goal.pptx"

    result = generate_ppt(content={"pages": pages}, output=str(output))

    assert output.exists() and output.stat().st_size > 0
    assert result["slide_count"] == len(pages)
    assert result["shapes_count"] >= len(pages)

    generated = PptxPresentation(str(output))
    assert len(generated.slides) == len(pages)
    assert all(len(slide.shapes) > 0 for slide in generated.slides)


def test_pipeline_forwards_explicit_theme_atoms(tmp_path):
    output = tmp_path / "explicit-theme.pptx"

    result = generate_ppt(
        content={"pages": [{"goal": "hook", "title": "Theme QA", "bullets": ["One"]}]},
        style="dark-tech",
        decoration="corner-accent",
        layout_variant="centered",
        mood="elegant",
        style_seed=17,
        output=str(output),
    )

    assert output.exists()
    assert result["theme_atoms"] == {
        "palette": "cyber-neon",
        "fonts": "tech-mono",
        "decoration": "accent-bar",
        "layout": "centered",
        "moods": ["elegant"],
    }
    assert result["theme_context"]["source"]["fallbacks"] == [
        {"field": "decoration", "requested": "corner-accent", "used": "accent-bar"}
    ]


def test_style_recommendations_are_distinct_and_composable():
    recommendations = recommend_styles("AI technology platform")

    assert [item["style"] for item in recommendations] == ["dark-tech", "neon", "sci"]
    assert len({item["palette"] for item in recommendations}) >= 2
    assert all({"style", "palette", "fonts", "decoration", "layout"} <= item.keys() for item in recommendations)

    theme = ThemeComposer().compose(style=recommendations[0]["style"])
    assert theme["atoms"]["palette"] == recommendations[0]["palette"]


def test_style_recommendations_recognize_chinese_luxury_keywords():
    recommendations = recommend_styles("高定香水画册")

    assert [item["style"] for item in recommendations] == ["warm-elegant", "zen", "professional"]
