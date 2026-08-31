"""Regression coverage for theme propagation across every Diagram type."""

from __future__ import annotations

import pytest
from pptx import Presentation as PptxPresentation

from pptx_designer import Presentation
from pptx_designer.diagrams import (
    CycleDiagram,
    DiagramStyle,
    FlowchartDiagram,
    FunnelDiagram,
    HierarchyDiagram,
    MatrixDiagram,
    PyramidDiagram,
    Region,
    SwotDiagram,
    TableDiagram,
    TimelineDiagram,
    VennDiagram,
)
from pptx_designer.renderer.theme import ThemeComposer


def _solid_fills(slide) -> set[str]:
    fills: set[str] = set()
    for shape in slide.shapes:
        try:
            if shape.fill.fore_color.rgb:
                fills.add(str(shape.fill.fore_color.rgb))
        except (AttributeError, TypeError, ValueError):
            continue
    return fills


DIAGRAMS = [
    (FlowchartDiagram, {"direction": "horizontal", "nodes": [{"label": "A"}, {"label": "B"}]}),
    (TimelineDiagram, {"events": [{"date": "2025", "label": "A"}, {"date": "2026", "label": "B"}]}),
    (SwotDiagram, {"quadrants": [{"label": "S"}, {"label": "W"}, {"label": "O"}, {"label": "T"}]}),
    (MatrixDiagram, {"rows": [{"label": "R1"}], "cols": [{"label": "C1"}], "cells": [["X"]]}),
    (TableDiagram, {"headers": ["A", "B"], "rows": [["1", "2"]]}),
    (HierarchyDiagram, {"nodes": [{"id": "root", "label": "Root", "level": 0}, {"id": "child", "label": "Child", "level": 1, "parent": "root"}]}),
    (VennDiagram, {"sets": [{"label": "A"}, {"label": "B"}]}),
    (CycleDiagram, {"nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}),
    (FunnelDiagram, {"stages": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}),
    (PyramidDiagram, {"levels": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}),
]


@pytest.mark.parametrize("diagram_cls,data", DIAGRAMS)
@pytest.mark.parametrize("theme_style", ["warm-elegant", "dark-tech"])
def test_every_diagram_uses_theme_colors(diagram_cls, data, theme_style):
    theme = ThemeComposer().compose(style=theme_style, seed=42)
    prs = Presentation(theme=theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    region = Region(left=1, top=1, width=11, height=5)

    style = DiagramStyle.from_theme(theme)
    # Disable the visual gradient so the regression assertion can inspect the
    # resolved solid fill through python-pptx. Gradient colors are covered by
    # the existing visual/integration tests.
    style.node_gradient = False
    diagram_cls(data=data, style=style, region=region).render(slide)

    fills = _solid_fills(slide)
    assert fills
    assert "2563EB" not in fills
    assert style.resolve_color("primary") == theme["semantic_roles"]["data-series-1"]
    assert "2563EB" not in fills


def test_all_diagram_theme_render_can_be_reopened(tmp_path):
    theme = ThemeComposer().compose(style="warm-elegant", seed=42)
    prs = Presentation(theme=theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = DiagramStyle.from_theme(theme)
    style.node_gradient = False
    FlowchartDiagram(
        data={"direction": "horizontal", "nodes": [{"label": "输入"}, {"label": "输出"}]},
        style=style,
        region=Region(left=1, top=1, width=11, height=5),
    ).render(slide)

    output = tmp_path / "diagram-theme-roundtrip.pptx"
    prs.save(output)
    reopened = PptxPresentation(output)

    assert len(reopened.slides) == 1
    assert len(reopened.slides[0].shapes) >= 2
