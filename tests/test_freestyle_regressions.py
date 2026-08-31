"""Regression tests for the editable FreeStyle renderer."""

from pptx import Presentation as PptxPresentation

from pptx_designer.core.pipeline import generate_ppt
from pptx_designer.core.professional_renderer import _contrast_ratio
from pptx_designer.renderer.theme import ThemeComposer


def _texts(prs):
    return [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text
    ]


def test_freestyle_preserves_subtitles_for_all_structured_page_types(tmp_path):
    goals = ["problem", "solution", "features", "data", "code", "content"]
    pages = [
        {"goal": goal, "title": goal, "subtitle": f"subtitle-{goal}", "bullets": ["One"]}
        for goal in goals
    ]
    output = tmp_path / "subtitles.pptx"

    generate_ppt(content={"pages": pages}, output=str(output))
    texts = _texts(PptxPresentation(str(output)))

    assert all(f"subtitle-{goal}" in texts for goal in goals)


def test_freestyle_cta_renders_subtitle_and_action(tmp_path):
    output = tmp_path / "cta.pptx"

    generate_ppt(
        content={"pages": [{"goal": "cta", "title": "Next step", "subtitle": "Book a pilot"}]},
        output=str(output),
    )
    texts = _texts(PptxPresentation(str(output)))

    assert "Book a pilot" in texts
    assert "开始试点  →" in texts


def test_freestyle_long_metric_value_uses_compact_type_and_separate_label(tmp_path):
    output = tmp_path / "metrics.pptx"

    generate_ppt(
        content={
            "pages": [
                {
                    "goal": "data",
                    "title": "Metrics",
                    "bullets": ["检索耗时: 18min → 3min", "一次解决率: 62% → 86%"],
                }
            ]
        },
        output=str(output),
    )
    prs = PptxPresentation(str(output))
    slide = prs.slides[0]
    value = next(shape for shape in slide.shapes if shape.text == "18min → 3min")
    label = next(shape for shape in slide.shapes if shape.text == "检索耗时")

    assert value.text_frame.paragraphs[0].runs[0].font.size.pt == 24
    assert value.top + value.height <= label.top


def test_freestyle_metric_values_are_contrast_safe_for_every_preset(tmp_path):
    page = {"goal": "data", "title": "Metrics", "bullets": ["Latency: 18min → 3min"]}

    for style in ThemeComposer.available_presets():
        output = tmp_path / f"{style}.pptx"
        generate_ppt(content={"pages": [page]}, style=style, style_seed=17, output=str(output))
        prs = PptxPresentation(str(output))
        value = next(shape for shape in prs.slides[0].shapes if shape.text == "18min → 3min")
        value_color = f"#{value.text_frame.paragraphs[0].runs[0].font.color.rgb}"
        surface = ThemeComposer().compose(style=style, seed=17)["semantic_roles"]["surface"]
        assert _contrast_ratio(value_color, surface) >= 3.0
