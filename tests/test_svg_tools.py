"""SVG tools tests — validates the high-level svg_chart helper."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_designer import Presentation as DesignerPresentation
from pptx_designer import set_slide_theme
from pptx_designer import svg_chart as root_svg_chart
from pptx_designer.compiler import SVGCompileError
from pptx_designer.renderer.theme import ThemeComposer
from pptx_designer.tools import svg_chart as tools_svg_chart
from pptx_designer.tools.svg import svg_chart


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestSvgChart:
    def test_public_import_paths_export_same_function(self):
        assert root_svg_chart is svg_chart
        assert tools_svg_chart is svg_chart

    def test_basic_circle(self):
        svg = '<svg viewBox="0 0 400 300"><circle cx="200" cy="150" r="100" fill="#4472C4"/></svg>'
        result = svg_chart(_slide(), svg, x=1, y=1, w=8, h=6)
        assert result.shape_count >= 1

    def test_empty_svg_returns_result(self):
        result = svg_chart(_slide(), "")
        assert result.shape_count == 0
        assert "empty svg_text" in result.warnings

    def test_with_color_context(self):
        svg = '<svg viewBox="0 0 400 300"><rect x="50" y="50" width="300" height="200" fill="#FFF"/></svg>'
        colors = {"primary": "#4472C4", "text_dark": "#333"}
        result = svg_chart(_slide(), svg, C=colors)
        assert result.shape_count >= 1

    def test_inherits_presentation_theme_color_context(self):
        theme = ThemeComposer().compose(style="warm-elegant", seed=17)
        prs = DesignerPresentation(theme=theme)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        svg = '<svg viewBox="0 0 100 100"><rect width="100" height="100" fill="primary"/></svg>'

        result = svg_chart(slide, svg)

        assert result.shape_count == 1
        assert str(slide.shapes[0].fill.fore_color.rgb) == theme["semantic_roles"]["data-series-1"].lstrip("#")

    def test_explicit_color_context_overrides_inherited_theme(self):
        theme = ThemeComposer().compose(style="warm-elegant", seed=17)
        prs = DesignerPresentation(theme=theme)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        svg = '<svg viewBox="0 0 100 100"><rect width="100" height="100" fill="primary"/></svg>'

        svg_chart(slide, svg, C={"primary": "#123456"})

        assert str(slide.shapes[0].fill.fore_color.rgb) == "123456"

    @pytest.mark.parametrize("token", ["background", "surface"])
    def test_inherited_semantic_background_tokens_are_resolved(self, token):
        theme = ThemeComposer().compose(style="glassmorphism", seed=42)
        prs = DesignerPresentation(theme=theme)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        svg = f'<svg viewBox="0 0 100 100"><rect width="100" height="100" fill="{token}"/></svg>'

        result = svg_chart(slide, svg)

        assert result.shape_count == 1
        expected = theme["semantic_roles"]["background" if token == "background" else "surface"]
        assert str(slide.shapes[0].fill.fore_color.rgb) == expected.lstrip("#")

    def test_slide_theme_overrides_presentation_theme(self):
        presentation_theme = ThemeComposer().compose(style="warm-elegant", seed=42)
        slide_theme = ThemeComposer().compose(style="dark-tech", seed=42)
        prs = DesignerPresentation(theme=presentation_theme)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_theme(slide, slide_theme)
        svg = '<svg viewBox="0 0 100 100"><rect width="100" height="100" fill="primary"/></svg>'

        svg_chart(slide, svg)

        assert str(slide.shapes[0].fill.fore_color.rgb) == slide_theme["semantic_roles"]["data-series-1"].lstrip("#")

    def test_css_variable_uses_passed_theme_context(self):
        theme = ThemeComposer().compose(style="warm-elegant", seed=42)
        prs = DesignerPresentation(theme=theme)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        svg = (
            '<svg viewBox="0 0 100 100"><style>:root { --brand: primary; }</style>'
            '<rect width="100" height="100" fill="var(--brand)"/></svg>'
        )

        result = svg_chart(slide, svg)

        assert result.warnings == []
        assert str(slide.shapes[0].fill.fore_color.rgb) == theme["semantic_roles"]["data-series-1"].lstrip("#")

    def test_gradient_stop_semantic_tokens_use_theme_context(self):
        theme = ThemeComposer().compose(style="warm-elegant", seed=42)
        prs = DesignerPresentation(theme=theme)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        svg = (
            '<svg viewBox="0 0 100 100"><defs><linearGradient id="g">'
            '<stop offset="0" stop-color="primary"/><stop offset="1" stop-color="accent"/>'
            '</linearGradient></defs><rect width="100" height="100" fill="url(#g)"/></svg>'
        )

        result = svg_chart(slide, svg)

        assert result.warnings == []
        assert "gradFill" in slide.shapes[0]._element.xml

    def test_gradient_defs_work_without_xml_namespace(self):
        svg = (
            '<svg viewBox="0 0 100 100"><defs><linearGradient id="g">'
            '<stop offset="0" stop-color="#112233"/><stop offset="1" stop-color="#445566"/>'
            '</linearGradient></defs><rect width="100" height="100" fill="url(#g)"/></svg>'
        )

        slide = _slide()
        result = svg_chart(slide, svg)

        assert result.warnings == []
        assert "gradFill" in slide.shapes[0]._element.xml

    def test_invalid_svg_raises(self):
        with pytest.raises(SVGCompileError):
            svg_chart(_slide(), "not svg")

    def test_svg_with_text(self):
        svg = '<svg viewBox="0 0 400 300"><text x="200" y="150" font-size="24" fill="#333">Title</text></svg>'
        result = svg_chart(_slide(), svg)
        assert result.shape_count >= 1

    def test_text_style_applies_explicit_role_sizes(self):
        slide = _slide()
        svg = (
            '<svg viewBox="0 0 400 300">'
            '<rect x="0" y="0" width="400" height="300" fill="#132238"/>'
            '<text class="title" x="200" y="100" text-anchor="middle">Title</text>'
            '<text class="body" x="200" y="150" text-anchor="middle">Body</text>'
            '</svg>'
        )
        result = svg_chart(
            slide,
            svg,
            x=1,
            y=1,
            w=8,
            h=6,
            text_style={
                "title": {"font_size": 20, "color": "#FFFFFF"},
                "body": {"font_size": 12, "color": "#D7E2EF"},
            },
        )
        text_shapes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text_frame.paragraphs[0].runs
        ]
        sizes = [shape.text_frame.paragraphs[0].runs[0].font.size.pt for shape in text_shapes]
        assert 20 in sizes
        assert 12 in sizes
        assert result.warnings == []

    def test_layout_contract_warns_and_can_fail_on_text_zone_violations(self):
        svg = '<svg viewBox="0 0 400 300"><text class="title" x="390" y="20">Title</text></svg>'
        layout = {"safe_margin": 20, "zones": {"title": (20, 20, 250, 80)}}
        result = svg_chart(_slide(), svg, layout=layout)
        assert any("unsafe SVG margin" in warning for warning in result.warnings)
        assert any("outside zone" in warning for warning in result.warnings)

        with pytest.raises(SVGCompileError, match="layout contract failed"):
            svg_chart(_slide(), svg, layout={**layout, "text_collision": "error"})

    def test_layout_contract_applies_nested_transforms(self):
        svg = (
            '<svg viewBox="0 0 400 300">'
            '<g transform="translate(100 50)"><g transform="scale(2)">'
            '<text class="title" x="50" y="50">Title</text>'
            "</g></g></svg>"
        )
        result = svg_chart(
            _slide(),
            svg,
            layout={"safe_margin": 20, "zones": {"title": (190, 140, 120, 80)}},
        )
        assert result.warnings == []

    def test_invalid_layout_mode_is_rejected(self):
        with pytest.raises(ValueError, match="layout.text_collision"):
            svg_chart(_slide(), '<svg viewBox="0 0 10 10"/>', layout={"text_collision": "fail"})

    def test_svg_with_gradient(self):
        svg = """<svg viewBox="0 0 400 300">
            <defs><linearGradient id="g1" x1="0" y1="0" x2="1" y2="0">
                <stop offset="0" stop-color="#4472C4"/><stop offset="1" stop-color="#2E75B6"/>
            </linearGradient></defs>
            <rect x="0" y="0" width="400" height="300" fill="url(#g1)"/>
        </svg>"""
        result = svg_chart(_slide(), svg)
        assert result.shape_count >= 1
