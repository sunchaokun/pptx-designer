"""SVG compiler integration tests — covers the main compilation path."""

from __future__ import annotations

import time

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_designer.compiler import SVGCompileError, SVGCompiler, SVGRenderReport, SVGResult


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


class TestBasicShapes:
    def test_circle(self):
        svg = '<svg viewBox="0 0 400 300"><circle cx="200" cy="150" r="100" fill="#4472C4"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert isinstance(result, SVGResult)
        assert result.shape_count >= 1

    def test_rect(self):
        svg = '<svg viewBox="0 0 400 300"><rect x="50" y="50" width="300" height="200" fill="#E74C3C"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_ellipse(self):
        svg = '<svg viewBox="0 0 400 300"><ellipse cx="200" cy="150" rx="150" ry="100" fill="#2ECC71"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_polygon(self):
        svg = '<svg viewBox="0 0 400 300"><polygon points="200,30 370,280 30,280" fill="#F39C12"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_line(self):
        svg = (
            '<svg viewBox="0 0 400 300"><line x1="50" y1="50" x2="350" y2="250" stroke="#333" stroke-width="3"/></svg>'
        )
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_result_contains_shapes_metrics_and_source_mapping(self):
        svg = (
            '<svg viewBox="0 0 100 100">'
            '<rect id="background" x="0" y="0" width="100" height="100" fill="#4472C4"/>'
            "</svg>"
        )
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert isinstance(result, SVGRenderReport)
        assert result.shapes == result.native_shapes
        assert len(result.shapes) == result.shape_count
        assert result.source_to_output["background"]
        assert result.metrics["node_count"] >= 2
        assert result.metrics["ir_node_count"] == result.metrics["node_count"]
        assert result.metrics["ir_build_ms"] >= 0
        assert result.ir_document is not None
        assert result.ir_document.nodes_for_id("background")[0].tag == "rect"
        assert result.metrics["total_ms"] >= 0
        assert result.feature_levels["rect"] == "NATIVE"

    def test_total_metric_includes_post_processing(self, monkeypatch):
        def slow_overlap(slide, pre_count, result):
            time.sleep(0.02)

        monkeypatch.setattr(SVGCompiler, "_detect_text_overlaps", staticmethod(slow_overlap))
        svg = '<svg viewBox="0 0 10 10"><rect width="10" height="10" fill="#4472C4"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.metrics["total_ms"] >= 20


class TestPath:
    def test_simple_path(self):
        svg = '<svg viewBox="0 0 400 300"><path d="M100,100 L300,100 L200,250 Z" fill="#9B59B6"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_curved_path(self):
        svg = '<svg viewBox="0 0 400 300"><path d="M100,200 C150,50 250,50 300,200" fill="none" stroke="#1ABC9C" stroke-width="4"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1


class TestText:
    def test_simple_text(self):
        svg = '<svg viewBox="0 0 400 300"><text x="200" y="150" text-anchor="middle" font-size="24" fill="#333">Hello</text></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_tspan(self):
        svg = '<svg viewBox="0 0 400 300"><text x="100" y="100"><tspan font-size="18" fill="#E74C3C">Line 1</tspan><tspan x="100" dy="30" font-size="14" fill="#333">Line 2</tspan></text></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_tspan_font_size_uses_svg_to_ppt_scale(self):
        svg = '<svg viewBox="0 0 1280 720"><text x="100" y="100" font-size="20"><tspan>Scaled</tspan></text></svg>'
        result = SVGCompiler().compile(svg, _slide(), (0, 0, 13.333, 7.5))
        run = result.shapes[0].text_frame.paragraphs[0].runs[0]
        assert run.font.size.pt == pytest.approx(15.0, abs=0.1)


class TestGradients:
    def test_linear_gradient(self):
        svg = """<svg viewBox="0 0 400 300">
            <defs><linearGradient id="g1" x1="0" y1="0" x2="1" y2="1">
                <stop offset="0" stop-color="#4472C4"/><stop offset="1" stop-color="#2E75B6"/>
            </linearGradient></defs>
            <rect x="50" y="50" width="300" height="200" fill="url(#g1)"/>
        </svg>"""
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_radial_gradient(self):
        svg = """<svg viewBox="0 0 400 300">
            <defs><radialGradient id="g2" cx="50%" cy="50%" r="50%">
                <stop offset="0" stop-color="#FFF"/><stop offset="1" stop-color="#4472C4"/>
            </radialGradient></defs>
            <circle cx="200" cy="150" r="120" fill="url(#g2)"/>
        </svg>"""
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1


class TestTransform:
    def test_translate(self):
        svg = '<svg viewBox="0 0 400 300"><g transform="translate(100,50)"><rect x="0" y="0" width="200" height="150" fill="#3498DB"/></g></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1

    def test_scale(self):
        svg = '<svg viewBox="0 0 400 300"><g transform="scale(2)"><rect x="10" y="10" width="50" height="40" fill="#E74C3C"/></g></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1


class TestClipPath:
    def test_clip_path(self):
        svg = """<svg viewBox="0 0 400 300">
            <defs><clipPath id="cp"><rect x="50" y="50" width="200" height="200"/></clipPath></defs>
            <rect x="0" y="0" width="400" height="300" fill="#4472C4" clip-path="url(#cp)"/>
        </svg>"""
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count >= 1


class TestUseAndVisibility:
    def test_use_maps_referenced_and_use_ids(self):
        svg = """<svg viewBox="0 0 100 100">
            <defs><rect id="source" x="0" y="0" width="10" height="10" fill="#4472C4"/></defs>
            <use id="copy" href="#source" x="20" y="20"/>
        </svg>"""
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.source_to_output["source"]
        assert result.source_to_output["copy"]

    def test_symbol_use_preserves_references_and_translation(self):
        svg = """<svg viewBox="0 0 100 100">
            <defs><symbol id="badge"><rect id="badge-rect" width="10" height="10" fill="#4472C4"/></symbol></defs>
            <use id="placed-badge" href="#badge" x="20" y="30"/>
        </svg>"""
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))

        shape = result.source_to_output["placed-badge"][0]
        assert result.source_to_output["badge"]
        assert result.source_to_output["badge-rect"]
        # Default `contain` scaling centers the 6×6-inch square viewBox area
        # inside the 8×6-inch target rectangle, adding one inch on the x axis.
        assert shape.left.inches == pytest.approx(3.2, abs=0.01)
        assert shape.top.inches == pytest.approx(2.8, abs=0.01)

    @pytest.mark.parametrize("attr", ['display="none"', 'visibility="hidden"'])
    def test_hidden_elements_are_not_rendered(self, attr):
        svg = f'<svg viewBox="0 0 10 10"><rect {attr} width="10" height="10" fill="#4472C4"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert result.shape_count == 0
        assert result.shapes == []


class TestGroupOpacity:
    def test_group_opacity_fails_before_creating_incorrect_shapes(self):
        svg = '<svg viewBox="0 0 10 10"><g id="faded" opacity="0.5"><rect width="10" height="10" fill="#112233"/></g></svg>'
        slide = _slide()
        with pytest.raises(SVGCompileError, match="requires raster fallback"):
            SVGCompiler().compile(svg, slide, (0, 0, 1, 1))
        assert len(slide.shapes) == 0


class TestUnsupported:
    def test_image_warning(self):
        svg = '<svg viewBox="0 0 400 300"><image href="photo.jpg" width="400" height="300"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        assert any("image" in w.lower() or "unsupported" in w.lower() for w in result.warnings)
        assert result.feature_levels["image"] == "RASTER_FALLBACK_CANDIDATE"

    @pytest.mark.parametrize("feature", ["filter", "mask"])
    def test_raster_only_feature_is_exposed_in_render_report(self, feature):
        definition = (
            '<filter id="fx"><feGaussianBlur stdDeviation="2"/></filter>'
            if feature == "filter"
            else '<mask id="fx"><rect width="100" height="100" fill="white"/></mask>'
        )
        svg = (
            '<svg viewBox="0 0 100 100"><defs>'
            f"{definition}"
            f'</defs><rect width="100" height="100" fill="#123456" {feature}="url(#fx)"/></svg>'
        )

        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))

        assert result.shape_count == 1
        assert result.feature_levels[feature] == "RASTER_FALLBACK_CANDIDATE"
        assert "raster_fallback_candidate" in result.feature_levels


class TestRoundTrip:
    def test_complex_svg_can_be_saved_and_reopened_as_pptx(self, tmp_path):
        svg = """<svg viewBox="0 0 400 300">
            <style>.panel { fill: #123456; } .label { fill: #FFFFFF; font-size: 28px; }</style>
            <defs>
                <linearGradient id="accent" x1="0" y1="0" x2="1" y2="1">
                    <stop offset="0" stop-color="#00B8D9"/><stop offset="1" stop-color="#6554C0"/>
                </linearGradient>
                <clipPath id="window"><rect x="40" y="40" width="320" height="220" rx="16"/></clipPath>
            </defs>
            <rect class="panel" width="400" height="300"/>
            <g transform="translate(40,40)" clip-path="url(#window)">
                <rect width="320" height="220" fill="url(#accent)"/>
                <text class="label" x="30" y="80">SVG QA</text>
            </g>
        </svg>"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        result = SVGCompiler().compile(svg, slide, (1, 1, 8, 6))
        output = tmp_path / "svg-roundtrip.pptx"
        prs.save(output)
        reopened = Presentation(output)

        assert result.shape_count >= 3
        assert len(reopened.slides) == 1
        assert len(reopened.slides[0].shapes) == result.shape_count


class TestTextOpacity:
    def test_text_opacity_writes_alpha(self):
        svg = '<svg viewBox="0 0 100 100"><text x="10" y="50" opacity="0.5" fill="#112233">Hello</text></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6))
        run_xml = result.shapes[0].text_frame.paragraphs[0].runs[0]._r.xml
        assert 'a:alpha val="50000"' in run_xml


class TestScalingModes:
    def test_contain(self):
        svg = '<svg viewBox="0 0 400 300"><rect x="0" y="0" width="400" height="300" fill="#4472C4"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6), scaling="contain")
        assert result.shape_count >= 1

    def test_cover(self):
        svg = '<svg viewBox="0 0 400 300"><rect x="0" y="0" width="400" height="300" fill="#4472C4"/></svg>'
        result = SVGCompiler().compile(svg, _slide(), (1, 1, 8, 6), scaling="cover")
        assert result.shape_count >= 1


class TestErrors:
    def test_empty_svg(self):
        with pytest.raises(SVGCompileError):
            SVGCompiler().compile("", _slide(), (1, 1, 8, 6))

    def test_invalid_svg(self):
        with pytest.raises(SVGCompileError):
            SVGCompiler().compile("not svg at all", _slide(), (1, 1, 8, 6))

    def test_svg_size_limit(self):
        compiler = SVGCompiler(limits={"max_svg_bytes": 10})
        with pytest.raises(SVGCompileError, match="size"):
            compiler.compile('<svg viewBox="0 0 1 1"/>', _slide(), (1, 1, 8, 6))

    def test_node_limit(self):
        compiler = SVGCompiler(limits={"max_nodes": 2})
        svg = '<svg viewBox="0 0 1 1"><g><rect width="1" height="1"/></g></svg>'
        with pytest.raises(SVGCompileError, match="node count"):
            compiler.compile(svg, _slide(), (1, 1, 8, 6))

    def test_node_limit_is_checked_before_ir_is_built(self, monkeypatch):
        import pptx_designer.compiler._compiler as compiler_module

        called = False

        def fail_if_called(_root):
            nonlocal called
            called = True
            raise AssertionError("IR must not be built after a node-limit failure")

        monkeypatch.setattr(compiler_module, "build_svg_ir", fail_if_called)
        compiler = SVGCompiler(limits={"max_nodes": 2})
        svg = '<svg viewBox="0 0 1 1"><g><rect width="1" height="1"/></g></svg>'
        with pytest.raises(SVGCompileError, match="node count"):
            compiler.compile(svg, _slide(), (0, 0, 1, 1))
        assert not called

    def test_tree_depth_limit(self):
        compiler = SVGCompiler(limits={"max_tree_depth": 3})
        svg = '<svg viewBox="0 0 1 1"><g><g><rect width="1" height="1"/></g></g></svg>'
        with pytest.raises(SVGCompileError, match="tree depth"):
            compiler.compile(svg, _slide(), (0, 0, 1, 1))

    def test_path_command_limit(self):
        compiler = SVGCompiler(limits={"max_path_commands": 2})
        svg = '<svg viewBox="0 0 1 1"><path d="M0 0 L1 0 L1 1 Z"/></svg>'
        with pytest.raises(SVGCompileError, match="path command"):
            compiler.compile(svg, _slide(), (1, 1, 8, 6))


class TestCompatShim:
    """Verify the old import path still works via the compatibility shim."""

    def test_import_from_renderer_svg_compiler(self):
        from pptx_designer.renderer.svg_compiler import SVGCompiler as OldSVGCompiler
        from pptx_designer.renderer.svg_compiler import SVGResult as OldSVGResult

        assert OldSVGCompiler is SVGCompiler
        assert OldSVGResult is SVGResult
