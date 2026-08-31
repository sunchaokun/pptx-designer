"""Regression coverage for template-source isolation in VI delivery."""

from pptx import Presentation

from pptx_designer.enterprise import VIBuildDelivery, VITemplateAdapter


def _strict_framework_context():
    return {
        "schema_version": "1.1",
        "framework_pages": [
            {
                "id": "cover",
                "role": "cover",
                "reference_slide": 1,
                "text_contract": {"strict": True},
            }
        ],
        "content_slots": [
            {
                "id": "cover.title",
                "page_role": "cover",
                "target": {"shape_index": 0},
                "text_style": {"font_size": 28},
            }
        ],
        "archetypes": [],
    }


def test_strict_framework_contract_rejects_unowned_template_text():
    presentation = Presentation()
    source = presentation.slides.add_slide(presentation.slide_layouts[6])
    source.shapes.add_textbox(0, 0, 1000000, 400000).text = "Template title"
    source.shapes.add_textbox(0, 600000, 1000000, 400000).text = "Template placeholder"
    adapter = VITemplateAdapter(_strict_framework_context())
    spec = adapter.compile(page_role="cover", content={"slots": {"cover.title": "Approved"}})

    try:
        adapter.render(spec, presentation)
    except ValueError as exc:
        assert "unowned framework text shapes: 1" in str(exc)
    else:
        raise AssertionError("strict framework pages must reject unowned template text")


def test_delivery_removes_template_source_pages_and_checks_provenance(tmp_path):
    presentation = Presentation()
    source = presentation.slides.add_slide(presentation.slide_layouts[6])
    source.shapes.add_textbox(0, 0, 4000000, 600000).text = "Template sample title"
    adapter = VITemplateAdapter(_strict_framework_context())
    delivery = VIBuildDelivery(presentation, adapter)
    spec = adapter.compile(page_role="cover", content={"slots": {"cover.title": "Approved delivery title"}})

    delivery.add(spec)
    output = tmp_path / "delivery.pptx"
    report = delivery.finalize(output, sample_texts=["Template sample title"], check_overlaps=False)
    reopened = Presentation(output)

    assert len(reopened.slides) == 1
    assert reopened.slides[0].shapes[0].text == "Approved delivery title"
    assert report.status == "pass"
    assert report.vi_consumption["sample_text_leaks"] == []
    assert report.vi_consumption["delivery_origin_issues"] == []
