from PIL import Image
from pptx import Presentation

from pptx_designer.core.build_spec import render_build_spec


def test_build_spec_renderer_supports_native_editable_component_families(tmp_path):
    image = tmp_path / "photo.png"
    Image.new("RGB", (120, 80), "#164A3B").save(image)
    context = {
        "colors": {"background": "#FFFFFF", "primary": "#164A3B"},
        "components": {
            "title": {"kind": "text", "bounds": {"left": 0.5, "top": 0.4, "width": 5, "height": 0.5}},
            "kpi": {"kind": "kpi_card", "bounds": {"left": 0.5, "top": 1.5, "width": 2.4, "height": 1.5}},
            "chart": {"kind": "bar_chart", "bounds": {"left": 3.5, "top": 1.5, "width": 5, "height": 2}},
            "photo": {"kind": "photo_panel", "bounds": {"left": 9, "top": 0.5, "width": 3.5, "height": 2.5}},
        },
    }
    spec = {
        "kind": "BuildSpec",
        "render_strategy": "components",
        "components": [
            {"component_id": "title", "data": "Editable title"},
            {"component_id": "kpi", "data": {"number": "42", "label": "Users", "trend": "+8%"}},
            {"component_id": "chart", "data": {"data": [("A", 0.6, "60"), ("B", 0.3, "30")]}},
            {"component_id": "photo", "data": str(image)},
        ],
    }
    prs = Presentation()
    render_build_spec(spec, prs, context)
    output = tmp_path / "build-spec.pptx"
    prs.save(output)
    reopened = Presentation(output)
    texts = [shape.text for shape in reopened.slides[0].shapes if getattr(shape, "has_text_frame", False)]
    assert "Editable title" in texts
    assert "42" in texts
    assert any(shape.shape_type == 13 for shape in reopened.slides[0].shapes)
