from pptx import Presentation

from pptx_designer.qa.structural import StructuralQA, run_structural_qa


def test_structural_qa_rejects_vi_regressions(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(100000, 100000, 1000000, 300000).text = "Sample title"
    path = tmp_path / "qa.pptx"
    prs.save(path)
    plans = [
        {"page_role": "content", "render_strategy": "components", "variant_id": "split-a"},
        {"page_role": "content", "render_strategy": "components", "variant_id": "split-a"},
        {"page_role": "content", "render_strategy": "prototype", "variant_id": "split-b"},
    ]
    report = StructuralQA().check(path, vi_plans=plans, sample_texts=["Sample title"])
    kinds = {issue.kind for issue in report.fatal}
    assert report.status == "fail"
    assert {"content_prototype", "adjacent_variant_repeat", "sample_text_leak"} <= kinds


def test_run_structural_qa_forwards_vi_plans_to_the_check(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "qa-wrapper.pptx"
    prs.save(path)

    report = run_structural_qa(
        path,
        vi_plans=[{"page_role": "content", "render_strategy": "prototype"}],
    )

    assert "content_prototype" in {issue.kind for issue in report.fatal}
