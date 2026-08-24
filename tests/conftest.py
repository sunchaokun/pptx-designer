"""Shared test fixtures for pptx-designer tests."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def slide():
    """Create a blank 16:9 slide for testing."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)
