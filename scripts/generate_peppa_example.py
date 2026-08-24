"""Compile the editable one-page cartoon SVG example to a PPTX file."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from pptx_designer.compiler import SVGCompileError, SVGCompiler

SVG_PATH = Path("output/peppa-pig-cartoon-example.svg")
PPTX_PATH = Path("output/peppa-pig-cartoon-example.pptx")


def main() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = SVGCompiler().compile(SVG_PATH.read_text(encoding="utf-8"), slide, (0, 0, 13.333333, 7.5))
    if result.shape_count < 25:
        raise RuntimeError(f"cartoon compilation produced too few shapes: {result.shape_count}")
    if result.warnings:
        raise RuntimeError(f"cartoon compilation emitted warnings: {result.warnings}")
    if result.ir_document is None:
        raise RuntimeError("P2b SVG IR report is missing")
    if result.metrics["ir_node_count"] != result.metrics["node_count"]:
        raise RuntimeError("IR node count does not match the sanitized SVG tree")
    if result.ir_document.nodes_for_id("scene-title")[0].get("fill") != "#6B2A4A":
        raise RuntimeError("CSS class was not materialized into the SVG IR")
    for source_id in ("peppa", "peppa-dress", "flower-source", "flower-left", "flower-right"):
        if not result.source_to_output.get(source_id):
            raise RuntimeError(f"missing source-to-output mapping for {source_id}")
    if not {"text", "use", "geometry"} <= result.features:
        raise RuntimeError(f"incomplete feature report: {result.features}")
    if result.feature_levels.get("use") != "NATIVE":
        raise RuntimeError(f"unexpected use feature level: {result.feature_levels.get('use')}")

    presentation.save(PPTX_PATH)
    reopened = Presentation(PPTX_PATH)
    if len(reopened.slides) != 1 or len(reopened.slides[0].shapes) != len(result.shapes):
        raise RuntimeError("saved PPTX did not preserve the expected editable shapes")
    _assert_group_opacity_is_rejected()
    print(f"PPTX: {PPTX_PATH.resolve()}")
    print(f"editable shapes: {result.shape_count}")
    print("P2b report, mapping, CSS materialization and group-opacity safety: verified")


def _assert_group_opacity_is_rejected() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    svg = '<svg viewBox="0 0 10 10"><g opacity="0.5"><rect width="10" height="10"/></g></svg>'
    try:
        SVGCompiler().compile(svg, slide, (0, 0, 1, 1))
    except SVGCompileError as exc:
        if "requires raster fallback" in str(exc) and not slide.shapes:
            return
        raise RuntimeError(f"invalid group-opacity failure mode: {exc}") from exc
    raise RuntimeError("group opacity was silently rendered")


if __name__ == "__main__":
    main()
