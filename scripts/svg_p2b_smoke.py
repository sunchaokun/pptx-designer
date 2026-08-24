"""Generate and structurally verify a real PPTX through the SVG compiler.

This is a manual smoke harness, intentionally separate from unit tests.  It
exercises the current native renderer together with the incremental SVG IR and
writes a reviewable presentation under ``output/``.
"""
from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.util import Inches

from pptx_designer.compiler import SVGCompiler

OUTPUT = Path("output/svg-p2b-real-smoke.pptx")


def _add_slide(presentation: Presentation):
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def main() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    compiler = SVGCompiler()

    architecture_svg = """<svg viewBox="0 0 1280 720">
      <style>
        :root { --ink: #102A43; --accent: #1D72F3; }
        .title { fill: var(--ink); font-family: Arial; font-size: 34px; font-weight: bold; }
        .label { fill: #102A43; font-family: Arial; font-size: 20px; text-anchor: middle; }
        .node { fill: url(#surface); stroke: #A9C5F7; stroke-width: 3; }
      </style>
      <defs>
        <linearGradient id="surface" x1="0" y1="0" x2="1" y2="1">
          <stop offset="0" stop-color="#EAF2FF"/><stop offset="1" stop-color="#C7DDFF"/>
        </linearGradient>
        <clipPath id="round-window"><rect x="80" y="150" width="1120" height="470" rx="28"/></clipPath>
        <g id="service"><rect class="node" width="250" height="105" rx="16"/><text class="label" x="125" y="62">Native shape</text></g>
      </defs>
      <rect width="1280" height="720" fill="#F7FAFC"/>
      <text class="title" x="80" y="90">SVG Compiler — real PPTX smoke test</text>
      <g clip-path="url(#round-window)">
        <rect x="80" y="150" width="1120" height="470" fill="#FFFFFF"/>
        <path d="M270 345 L485 345 M795 345 L1010 345" fill="none" stroke="#1D72F3" stroke-width="8"/>
        <use id="service-left" href="#service" x="145" y="292"/>
        <use id="service-right" href="#service" x="885" y="292"/>
        <circle cx="640" cy="345" r="43" fill="#1D72F3"/>
        <text class="label" x="640" y="353">IR</text>
      </g>
    </svg>"""
    result_one = compiler.compile(architecture_svg, _add_slide(presentation), (0, 0, 13.333333, 7.5))
    assert result_one.shape_count >= 8
    assert result_one.ir_document is not None
    assert result_one.source_to_output["service-left"]
    assert result_one.source_to_output["service"]

    capability_svg = """<svg viewBox="0 0 1280 720">
      <style>.heading { fill: #102A43; font-family: Arial; font-size: 38px; font-weight: bold; }</style>
      <defs><filter id="blur"><feGaussianBlur stdDeviation="8"/></filter></defs>
      <rect width="1280" height="720" fill="#F7FAFC"/>
      <text class="heading" x="80" y="90">Capability reporting and editable output</text>
      <rect id="visible" x="80" y="160" width="500" height="280" rx="22" fill="#1D72F3" opacity="0.5"/>
      <rect id="hidden" x="650" y="160" width="500" height="280" fill="#FF0000" display="none"/>
      <text x="330" y="315" fill="#FFFFFF" font-family="Arial" font-size="30" text-anchor="middle">50% opacity</text>
      <image href="missing.png" x="80" y="520" width="50" height="50"/>
      <text x="160" y="552" fill="#52606D" font-family="Arial" font-size="24">image is safely skipped and reported</text>
    </svg>"""
    result_two = compiler.compile(capability_svg, _add_slide(presentation), (0, 0, 13.333333, 7.5))
    assert result_two.ir_document is not None
    assert "hidden" in result_two.ir_document.nodes_for_id("hidden")[0].features
    assert result_two.feature_levels["filter"] == "RASTER_FALLBACK_CANDIDATE"
    assert result_two.shape_count >= 3
    assert result_two.shapes
    assert any("image" in warning for warning in result_two.warnings)

    OUTPUT.parent.mkdir(exist_ok=True)
    presentation.save(OUTPUT)

    reopened = Presentation(OUTPUT)
    assert len(reopened.slides) == 2
    assert len(reopened.slides[0].shapes) == len(result_one.shapes)
    assert len(reopened.slides[1].shapes) == len(result_two.shapes)
    with ZipFile(OUTPUT) as package:
        slide_parts = [name for name in package.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        assert len(slide_parts) == 2
        assert all(package.read(part).startswith(b"<?xml") for part in slide_parts)

    print(f"PPTX: {OUTPUT.resolve()}")
    print(f"slide 1: {result_one.shape_count} shapes, {len(result_one.warnings)} warnings")
    print(f"slide 2: {result_two.shape_count} shapes, {len(result_two.warnings)} warnings")
    for warning in result_one.warnings + result_two.warnings:
        print(f"warning: {warning}")


if __name__ == "__main__":
    main()
