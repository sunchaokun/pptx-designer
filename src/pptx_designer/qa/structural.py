"""Offline, deterministic checks for PPTX delivery reliability.

This module deliberately avoids subjective visual scoring.  It checks whether
the file can be reopened, whether content stays inside the slide, whether
text is plausibly readable, and whether the deck remains natively editable.
"""

from __future__ import annotations

from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class QAIssue:
    slide: int | None
    kind: str
    severity: str
    message: str


@dataclass
class StructuralQAReport:
    status: str
    fatal: list[QAIssue] = field(default_factory=list)
    warnings: list[QAIssue] = field(default_factory=list)
    slide_count: int = 0
    shape_count: int = 0
    editable_ratio: float = 1.0
    vi_consumption: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


class StructuralQA:
    """Run deterministic delivery checks against a PowerPoint file."""

    def __init__(
        self,
        *,
        min_font_size_pt: float = 9.0,
        edge_tolerance_in: float = 0.02,
        check_overlaps: bool = False,
    ) -> None:
        self.min_font_size_pt = min_font_size_pt
        self.edge_tolerance_in = edge_tolerance_in
        self.check_overlaps = check_overlaps

    def check(
        self,
        path: str | Path,
        *,
        expected_slides: int | None = None,
        vi_plans: list[dict[str, Any]] | None = None,
        sample_texts: list[str] | None = None,
        min_content_variants: int = 3,
    ) -> StructuralQAReport:
        report = StructuralQAReport(status="pass")
        try:
            prs = Presentation(str(path))
        except Exception as exc:  # pragma: no cover - exact parser errors vary by python-pptx version
            report.status = "fail"
            report.fatal.append(QAIssue(None, "unreadable_file", "fatal", f"Cannot reopen PPTX: {exc}"))
            return report

        report.slide_count = len(prs.slides)
        if expected_slides is not None and report.slide_count != expected_slides:
            report.fatal.append(
                QAIssue(None, "slide_count", "fatal", f"Expected {expected_slides} slides, found {report.slide_count}")
            )

        slide_w = prs.slide_width / 914400.0
        slide_h = prs.slide_height / 914400.0
        total = editable = 0
        for slide_no, slide in enumerate(prs.slides, start=1):
            for shape_no, shape in enumerate(slide.shapes, start=1):
                total += 1
                if shape.shape_type not in {MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE}:
                    editable += 1
                self._check_bounds(report, slide_no, shape_no, shape, slide_w, slide_h)
                self._check_text(report, slide_no, shape_no, shape)
            if self.check_overlaps:
                self._check_overlaps(report, slide_no, slide.shapes)
        report.shape_count = total
        report.editable_ratio = editable / total if total else 1.0
        if total and report.editable_ratio < 0.8:
            report.warnings.append(
                QAIssue(None, "editable_ratio", "warning", f"Native editable shape ratio is {report.editable_ratio:.2f}")
            )
        if vi_plans is not None:
            self._check_vi_consumption(
                report,
                prs,
                vi_plans,
                sample_texts or [],
                min_content_variants=min_content_variants,
            )
        if report.fatal:
            report.status = "fail"
        elif report.warnings:
            report.status = "pass_with_warnings"
        return report

    def _check_vi_consumption(
        self,
        report: StructuralQAReport,
        prs: Any,
        plans: list[dict[str, Any]],
        sample_texts: list[str],
        *,
        min_content_variants: int,
    ) -> None:
        """Check that a VI plan was actually consumed by the output deck.

        The checks are deliberately plan-driven.  A plain PPTX cannot infer
        whether a page was intended to be a framework page or a Build page;
        the compiler already has that provenance and passes it here.
        """
        content = [item for item in plans if item.get("page_role") == "content"]
        prototype_content = [
            item for item in content if item.get("render_strategy") == "prototype"
        ]
        legacy_content = [item for item in content if not item.get("atomic_build")]
        variants = [str(item.get("variant_id")) for item in legacy_content if item.get("variant_id")]
        repeated_adjacent: list[int] = []
        previous = None
        for index, item in enumerate(plans, start=1):
            if item.get("page_role") != "content" or item.get("atomic_build"):
                previous = None
                continue
            variant = item.get("variant_id")
            if variant and variant == previous:
                repeated_adjacent.append(index)
            previous = variant

        leaked = []
        output_text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text
        )
        for sample in sample_texts:
            if sample and sample in output_text:
                leaked.append(sample)

        atomic_relation_issues: list[QAIssue] = []
        delivery_origin_issues: list[QAIssue] = []
        allowed_origins = {"framework_rebound", "build_atomic", "build_components"}
        if any("delivery_origin" in item for item in plans):
            for index, item in enumerate(plans, start=1):
                origin = item.get("delivery_origin")
                if origin not in allowed_origins:
                    delivery_origin_issues.append(
                        QAIssue(index, "delivery_origin_missing", "fatal", "Delivery page has no approved origin")
                    )
        for index, item in enumerate(plans, start=1):
            if item.get("page_role") != "content" or not item.get("atomic_build"):
                continue
            slide = item.get("page_number", index)
            content_model = item.get("content_model")
            if not isinstance(content_model, dict):
                atomic_relation_issues.append(
                    QAIssue(slide, "content_model_missing", "fatal", "Atomic content page has no content model")
                )
                continue
            relations = content_model.get("relations", [])
            if not isinstance(relations, list):
                atomic_relation_issues.append(
                    QAIssue(slide, "content_model_invalid", "fatal", "Content model relations must be a list")
                )
                continue
            relation_ids = {
                str(relation.get("id"))
                for relation in relations
                if isinstance(relation, dict) and relation.get("id")
            }
            atom_ids = {
                str(component.get("atom_id"))
                for component in item.get("components", [])
                if isinstance(component, dict) and component.get("atom_id")
            }
            bound: set[str] = set()
            for binding in item.get("relation_bindings", []):
                if not isinstance(binding, dict):
                    continue
                relation_id = str(binding.get("relation_id", ""))
                binding_atoms = {str(atom_id) for atom_id in binding.get("atom_ids", [])}
                if relation_id in relation_ids and binding_atoms and binding_atoms <= atom_ids:
                    bound.add(relation_id)
                else:
                    atomic_relation_issues.append(
                        QAIssue(slide, "relation_binding_invalid", "fatal", f"Invalid relation binding: {relation_id}")
                    )
            for relation_id in sorted(relation_ids - bound):
                atomic_relation_issues.append(
                    QAIssue(slide, "relation_uncovered", "fatal", f"Content relation is not represented: {relation_id}")
                )

        report.vi_consumption = {
            "content_pages": len(content),
            "content_prototype_pages": len(prototype_content),
            "distinct_variants": sorted(set(variants)),
            "adjacent_variant_repeat_pages": repeated_adjacent,
            "sample_text_leaks": leaked,
            "atomic_content_pages": len([item for item in content if item.get("atomic_build")]),
            "atomic_relation_issues": [issue.kind for issue in atomic_relation_issues],
            "delivery_origin_issues": [issue.kind for issue in delivery_origin_issues],
        }
        report.fatal.extend(atomic_relation_issues)
        report.fatal.extend(delivery_origin_issues)
        for item in prototype_content:
            report.fatal.append(
                QAIssue(item.get("page_number"), "content_prototype", "fatal",
                        "Content page was compiled with prototype rendering")
            )
        for page_number in repeated_adjacent:
            report.fatal.append(
                QAIssue(page_number, "adjacent_variant_repeat", "fatal",
                        "Adjacent content pages reuse the same layout variant")
            )
        if len(legacy_content) >= min_content_variants and len(set(variants)) < min_content_variants:
            report.fatal.append(
                QAIssue(None, "insufficient_layout_variants", "fatal",
                        f"Content pages use {len(set(variants))} variants; expected at least {min_content_variants}")
            )
        for sample in leaked:
            report.fatal.append(
                QAIssue(None, "sample_text_leak", "fatal", f"Template sample text remains: {sample[:80]}")
            )

    def _check_bounds(self, report, slide_no, shape_no, shape, slide_w, slide_h) -> None:
        # Intentional off-canvas ornaments are named explicitly by renderers.
        # Keep this narrow so normal content is never exempted accidentally.
        if shape.name.startswith("Background Decoration"):
            return
        if any(value is None for value in (shape.left, shape.top, shape.width, shape.height)):
            report.warnings.append(
                QAIssue(
                    slide_no,
                    "missing_geometry",
                    "warning",
                    f"Shape {shape_no} has no resolved geometry and was excluded from bounds checking",
                )
            )
            return
        x = shape.left / 914400.0
        y = shape.top / 914400.0
        w = shape.width / 914400.0
        h = shape.height / 914400.0
        if x < -self.edge_tolerance_in or y < -self.edge_tolerance_in or x + w > slide_w + self.edge_tolerance_in or y + h > slide_h + self.edge_tolerance_in:
            report.fatal.append(
                QAIssue(slide_no, "out_of_bounds", "fatal", f"Shape {shape_no} exceeds slide bounds: ({x:.2f},{y:.2f},{w:.2f},{h:.2f})")
            )

    def _check_text(self, report, slide_no, shape_no, shape) -> None:
        if not getattr(shape, "has_text_frame", False):
            return
        text = shape.text.strip()
        if not text:
            return
        sizes = [run.font.size.pt for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.font.size]
        if sizes and min(sizes) < self.min_font_size_pt:
            report.warnings.append(
                QAIssue(slide_no, "small_text", "warning", f"Shape {shape_no} contains {min(sizes):g}pt text")
            )
        lines = max(text.count("\n") + 1, len(shape.text_frame.paragraphs))
        if sizes and shape.height / 914400.0 < lines * min(sizes) * 1.15 / 72.0:
            report.warnings.append(QAIssue(slide_no, "text_box_height", "warning", f"Shape {shape_no} may clip {lines} text lines"))

    def _check_overlaps(self, report, slide_no, shapes) -> None:
        """Report meaningful text/media collisions without rejecting layering."""
        boxes = []
        for index, shape in enumerate(shapes, start=1):
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            boxes.append((index, "text", shape.left / 914400.0, shape.top / 914400.0,
                          shape.width / 914400.0, shape.height / 914400.0))
        pictures = [
            (index, "picture", shape.left / 914400.0, shape.top / 914400.0,
             shape.width / 914400.0, shape.height / 914400.0)
            for index, shape in enumerate(shapes, start=1)
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
        ]
        for text_box in boxes:
            for picture in pictures:
                # A picture below text is a legitimate editorial treatment;
                # only flag it when the later picture layer can cover text.
                if picture[0] < text_box[0]:
                    continue
                overlap = self._intersection_ratio(text_box[2:], picture[2:])
                if overlap >= 0.15:
                    report.warnings.append(
                        QAIssue(slide_no, "text_media_overlap", "warning",
                                f"Text shape {text_box[0]} overlaps picture {picture[0]} by {overlap:.0%}")
                    )
        for index, first in enumerate(boxes):
            for second in boxes[index + 1:]:
                overlap = self._intersection_ratio(first[2:], second[2:])
                if overlap >= 0.15:
                    report.warnings.append(
                        QAIssue(slide_no, "text_text_overlap", "warning",
                                f"Text shape {first[0]} overlaps text shape {second[0]} by {overlap:.0%}")
                    )

    @staticmethod
    def _intersection_ratio(first, second) -> float:
        ax, ay, aw, ah = first
        bx, by, bw, bh = second
        area = max(0.0, min(ax + aw, bx + bw) - max(ax, bx)) * max(0.0, min(ay + ah, by + bh) - max(ay, by))
        base = max(aw * ah, 0.0001)
        return area / base


def run_structural_qa(
    path: str | Path,
    *,
    expected_slides: int | None = None,
    vi_plans: list[dict[str, Any]] | None = None,
    sample_texts: list[str] | None = None,
    min_content_variants: int = 3,
    **kwargs,
) -> StructuralQAReport:
    """Run structural and optional VI-consumption checks through one public API."""
    return StructuralQA(**kwargs).check(
        path,
        expected_slides=expected_slides,
        vi_plans=vi_plans,
        sample_texts=sample_texts,
        min_content_variants=min_content_variants,
    )
