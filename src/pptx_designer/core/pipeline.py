"""Core pipeline — FreeStyle + Build mode PPT generation.

Usage:
    from pptx_designer import generate_ppt, Presentation

    # FreeStyle mode: one-liner generation
    result = generate_ppt("AI startup pitch deck", style="dark cyberpunk")

    # Build mode: pixel-perfect control
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # ... add shapes, text, charts ...
    prs.save("output.pptx")
"""

from __future__ import annotations

import os
import warnings
from collections.abc import Mapping
from typing import Any

from pptx import Presentation as _PptxPresentation


def Presentation(
    template_path: str | None = None,
    theme: Mapping[str, Any] | None = None,
    *,
    strict_theme: bool = False,
) -> _PptxPresentation:
    """Create a Presentation with pptx-designer defaults.

    Automatically sets widescreen 16:9 format.

    Args:
        template_path: Optional path to a .pptx template file.
        theme: Optional resolved theme inherited by Build Mode helpers.
        strict_theme: Require ``theme`` to be a complete resolved FreeStyle
            theme. Leave false for partial VI/template design contexts.

    Returns:
        A python-pptx Presentation object.
    """
    prs = _PptxPresentation(template_path) if template_path else _PptxPresentation()
    # Set widescreen 16:9
    prs.slide_width = 12192000  # 13.333 inches
    prs.slide_height = 6858000  # 7.5 inches
    if theme is not None:
        if strict_theme:
            from pptx_designer.renderer.theme import validate_resolved_theme

            validate_resolved_theme(theme)
        from pptx_designer.renderer.theme_context import set_presentation_theme

        set_presentation_theme(prs, theme)
    return prs


def generate_ppt(
    query: str = "",
    *,
    content: dict[str, Any] | None = None,
    style: str | None = None,
    palette: str | None = None,
    fonts: str | None = None,
    decoration: str | None = None,
    layout: str | None = None,
    layout_variant: str | None = None,
    mood: str | None = None,
    style_seed: int | None = None,
    theme: Mapping[str, Any] | None = None,
    template: str | None = None,
    output: str = "output.pptx",
    slides: int | None = None,
    **kwargs: Any,
) -> dict[str, Any]:
    """Generate a PowerPoint presentation.

    FreeStyle mode: provide a query, get a complete PPT.
    Build mode: provide structured content for full control.

    Args:
        query: Natural language description (FreeStyle mode).
        content: Structured FreeStyle content dict with title, pages, etc.
        style: Design style (e.g., "dark cyberpunk", "warm elegant").
        palette: Exact color palette name.
        fonts: Exact font pair name.
        decoration: Exact decoration style name.
        layout: Exact layout variant name.
        layout_variant: Backward-compatible alias for ``layout``.
        mood: Optional mood hint used when composing a theme.
        style_seed: Optional seed for reproducible theme selection.
        theme: Previously resolved theme. When supplied, it is used as-is and
            wins over theme discovery arguments.
        template: Path to a .pptx template.
        output: Output file path.
        slides: Number of slides (FreeStyle mode).

    Returns:
        Dictionary with keys: output_path, slide_count, shapes_count, and
        traceable resolved-theme information.
    """
    from pptx_designer.renderer.theme import ThemeComposer, validate_resolved_theme

    # A supplied theme is a locked, complete FreeStyle result. Validate it
    # before planning so a Theme Lock cannot fail later with an unrelated
    # renderer KeyError.
    if theme is not None:
        validate_resolved_theme(theme)

    ignored_arguments = {
        field: value
        for field, value in {
            "style": style,
            "palette": palette,
            "fonts": fonts,
            "decoration": decoration,
            "layout": layout,
            "layout_variant": layout_variant,
            "mood": mood,
            "style_seed": style_seed,
        }.items()
        if value is not None
    }
    if theme is not None and ignored_arguments:
        ignored_names = ", ".join(ignored_arguments)
        warnings.warn(
            f"theme= was supplied, so these theme-discovery arguments were ignored: {ignored_names}",
            UserWarning,
            stacklevel=2,
        )

    # ── Step 1: Get pages ──────────────────────────────────────────
    layout = layout or layout_variant
    if content:
        # Structured FreeStyle: user supplied page plan.
        pages = content.get("pages", [])
    else:
        # FreeStyle mode: use planner
        from pptx_designer.core.planner import StoryPlanner

        planner = StoryPlanner()
        story = planner.plan(query, slide_count_override=slides)
        pages = [
            {
                "goal": p.goal,
                "title": p.title,
                "subtitle": p.subtitle,
                "bullets": p.bullets,
            }
            for p in story.pages
        ]
        # The planner supplies only a style hint.  Theme resolution happens
        # once below so FreeStyle does not silently re-randomize a theme.
        if not style:
            style = story.style_hint

    # ── Step 2: Resolve one complete theme ─────────────────────────
    if theme is None:
        composer = ThemeComposer()
        theme = composer.compose(
            style=style,
            palette=palette,
            fonts=fonts,
            decoration=decoration,
            layout=layout,
            mood=mood,
            seed=style_seed,
            query=query,
        )

    # ── Step 3: Create presentation ────────────────────────────────
    prs = Presentation(template)

    # ── Step 4: Render each page ───────────────────────────────────
    from pptx_designer.core.professional_renderer import render_professional_page

    for i, page in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        goal = page.get("goal", "content")
        page_title = page.get("title", "")
        subtitle = page.get("subtitle", "")
        bullets = page.get("bullets", [])

        render_professional_page(slide, goal, page_title, subtitle, bullets, theme, i, len(pages))

    # ── Step 5: Save ───────────────────────────────────────────────
    output_dir = os.path.dirname(output)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    prs.save(output)

    shape_count = sum(len(slide.shapes) for slide in prs.slides)

    return {
        "output_path": output,
        "slide_count": len(prs.slides),
        "shapes_count": shape_count,
        "theme": theme["name"],
        "theme_atoms": theme["atoms"],
        "theme_context": theme,
        "theme_application": {
            "requested": theme.get("source", {}).get("requested", {}),
            "resolved": theme.get("source", {}).get("resolved", theme.get("atoms", {})),
            "applied_to": [
                "professional_renderer",
                "text",
                "semantic_roles",
                "typography",
                "decoration",
                "layout_variant",
            ],
            "not_applied": [
                {
                    "field": field,
                    "reason": "FreeStyle renderer does not consume this theme field yet.",
                }
                for field in ("text_effect_preset", "image_effect")
                if theme.get(field)
            ],
            "fallbacks": theme.get("source", {}).get("fallbacks", []),
            "warnings": theme.get("source", {}).get("warnings", []),
            "ignored_arguments": ignored_arguments if theme is not None else {},
        },
    }


def _render_page(
    slide: Any,
    goal: str,
    title: str,
    subtitle: str,
    bullets: list[str],
    C: dict,
    page_index: int,
    total_pages: int,
) -> None:
    """Render a single page based on its goal type."""
    from pptx_designer.tools.cards import code_block, highlight_cards, kpi_card
    from pptx_designer.tools.layout import page_header, page_number, top_bar
    from pptx_designer.tools.shapes import oval, rect, rrect
    from pptx_designer.tools.text import text

    # Common: top accent bar
    top_bar(slide, C.get("primary", "#1D78FA"))

    if goal == "hook":
        # ── Hero slide ───────────────────────────────────────────
        # Large title centered
        text(slide, 0.5, 1.8, 12.333, 1.5, title, font_size=44, bold=True, color="text_dark", C=C, anchor="middle")
        if subtitle:
            text(slide, 0.5, 3.5, 12.333, 0.8, subtitle, font_size=20, color="text_body", C=C, anchor="middle")
        # Decorative line
        rect(slide, 5.5, 4.5, 2.333, 0.04, C.get("accent", "#FF6B35"))

    elif goal == "problem":
        # ── Problem slide ────────────────────────────────────────
        page_header(slide, title, subtitle or "Why this matters", C=C)
        if bullets:
            y = 2.3
            for bullet in bullets:
                # Red accent bullet
                oval(slide, 1.0, y + 0.08, 0.16, 0.16, C.get("destructive", "#EF4444"))
                text(slide, 1.4, y, 10.5, 0.4, bullet, font_size=15, color="text_body", C=C)
                y += 0.6

    elif goal == "solution":
        # ── Solution slide ───────────────────────────────────────
        page_header(slide, title, subtitle or "How we solve it", C=C)
        if bullets:
            y = 2.3
            for bullet in bullets:
                # Green accent bullet
                oval(slide, 1.0, y + 0.08, 0.16, 0.16, "#22C55E")
                text(slide, 1.4, y, 10.5, 0.4, bullet, font_size=15, color="text_body", C=C)
                y += 0.6

    elif goal == "features":
        # ── Features with cards ──────────────────────────────────
        page_header(slide, title, subtitle or "What we offer", C=C)
        if bullets:
            cards_data = []
            accents = [C.get("primary", "#1D78FA"), C.get("accent", "#FF6B35"), "#22C55E", "#8B5CF6", "#F59E0B"]
            for j, bullet in enumerate(bullets[:5]):
                accent = accents[j % len(accents)]
                cards_data.append((bullet, "", accent))
            if cards_data:
                highlight_cards(slide, 1.0, 2.2, cards_data, total_width=11, C=C)

    elif goal == "data":
        # ── Data/Metrics with KPI cards ──────────────────────────
        page_header(slide, title, subtitle or "Key metrics", C=C)
        if bullets:
            # Try to parse as KPI cards (format: "Label: Value")
            kpi_data = []
            for bullet in bullets:
                if ": " in bullet or "：" in bullet:
                    sep = ": " if ": " in bullet else "："
                    parts = bullet.split(sep, 1)
                    kpi_data.append((parts[1] if len(parts) > 1 else "", parts[0]))
                else:
                    kpi_data.append((bullet, ""))

            if len(kpi_data) <= 4:
                # KPI card layout
                card_width = min(3.0, 10.0 / len(kpi_data))
                gap = 0.3
                total_width = len(kpi_data) * card_width + (len(kpi_data) - 1) * gap
                start_x = (13.333 - total_width) / 2
                for j, (value, label) in enumerate(kpi_data):
                    x = start_x + j * (card_width + gap)
                    kpi_card(slide, x, 2.5, card_width, 1.8, value, label, C=C)
            else:
                # Bullet list fallback
                y = 2.3
                for bullet in bullets:
                    text(slide, 1.0, y, 11.0, 0.4, f"• {bullet}", font_size=14, color="text_body", C=C)
                    y += 0.5

    elif goal == "code":
        # ── Code block ───────────────────────────────────────────
        page_header(slide, title, subtitle, C=C)
        if bullets:
            code_block(slide, 1.0, 2.2, 11.0, 4.0, bullets, language="python", C=C)

    elif goal == "exercise":
        # ── Exercise/Practice ────────────────────────────────────
        page_header(slide, title, subtitle or "Hands-on practice", C=C)
        if bullets:
            y = 2.3
            for j, bullet in enumerate(bullets):
                # Numbered steps
                rrect(slide, 1.0, y, 0.4, 0.35, C.get("primary", "#1D78FA"))
                text(
                    slide,
                    1.05,
                    y + 0.02,
                    0.3,
                    0.3,
                    str(j + 1),
                    font_size=12,
                    color="#FFFFFF",
                    bold=True,
                    align="center",
                    C=C,
                )
                text(slide, 1.6, y, 10.0, 0.4, bullet, font_size=14, color="text_body", C=C)
                y += 0.6

    elif goal == "overview":
        # ── Overview with sidebar ────────────────────────────────
        # Left sidebar
        rect(slide, 0, 0, 4.5, 7.5, C.get("primary", "#1D78FA"))
        text(slide, 0.5, 2.0, 3.5, 1.0, title, font_size=28, bold=True, color="#FFFFFF", C=C)
        if subtitle:
            text(slide, 0.5, 3.2, 3.5, 0.5, subtitle, font_size=14, color="#FFFFFF", C=C)
        # Right content
        if bullets:
            y = 2.0
            for bullet in bullets:
                text(slide, 5.2, y, 7.5, 0.4, f"• {bullet}", font_size=14, color="text_body", C=C)
                y += 0.5

    else:
        # ── Default content ──────────────────────────────────────
        page_header(slide, title, subtitle, C=C)
        if bullets:
            y = 2.3
            for bullet in bullets:
                text(slide, 1.0, y, 11.0, 0.4, f"• {bullet}", font_size=14, color="text_body", C=C)
                y += 0.5

    # ── Page number (all slides except first) ─────────────────────
    if page_index > 0:
        page_number(slide, page_index + 1, total_pages, style="simple", C=C)
