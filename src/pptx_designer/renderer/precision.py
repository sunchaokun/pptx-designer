"""PrecisionRenderer — brand-compliant + pixel-perfect rendering.

Combines Pipeline's brand constraint system (BrandSpec, template layouts,
logo/footer/watermark) with Build Script's rendering precision (run-level
fonts, Pillow cropping, per-element coordinates).

This is the third rendering path:
  - EnterpriseRenderer: template-driven, paragraph-level fonts → low quality
  - PPTRenderer (freestyle fallback): layout-registry driven → medium quality
  - PrecisionRenderer: brand-aware + run-level precision → delivery quality
"""

from __future__ import annotations

import hashlib
import logging
import os
import tempfile
from contextlib import suppress
from typing import Any

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from pptx_designer.effects.image_effects import (
    add_circle_image as _add_circle_image,
)
from pptx_designer.effects.image_effects import (
    add_image_in_shape as _add_image_in_shape,
)
from pptx_designer.effects.image_effects import (
    apply_blip_artistic,
    apply_blip_duotone,
)
from pptx_designer.effects.shape_effects import (
    GradientFill,
    GradientStop,
    apply_3d,
    apply_bevel,
    apply_frosted_glass,
    apply_glow,
    apply_gradient,
    apply_pattern_fill,
    apply_shadow,
    apply_soft_edge,
)
from pptx_designer.effects.text_effects import (
    apply_text_gradient,
    apply_text_gradient_preset,
    apply_text_outline,
    set_text_rotation,
    set_vertical_text,
)
from pptx_designer.enterprise.brand import BrandSpec
from pptx_designer.renderer.layout import SLIDE_HEIGHT, SLIDE_WIDTH
from pptx_designer.renderer.shapes import ShapeFactory

logger = logging.getLogger(__name__)

CORNER_RADIUS_SCALE = {
    "sm": 4,
    "md": 8,
    "lg": 12,
    "pill": 50,
}

# Unified 16:9 layout grid — every slide type aligns to the same frame.
LAYOUT_GRID = {
    "margin_left": 0.75,
    "margin_right": 0.75,
    "title_y": 0.6,
    "title_size": 30,
    "subtitle_y": 1.55,
    "content_y": 2.05,
    "content_bottom": 6.85,
    "footer_y": 7.0,
    "top_accent_h": 0.05,
}


class PrecisionRenderer:
    def __init__(self, brand_spec: BrandSpec | None = None, template_path: str | None = None):
        self._brand = brand_spec
        self._template_path = template_path
        self._has_template = bool(template_path and os.path.exists(template_path))
        self._crop_cache_dir = os.path.join(tempfile.gettempdir(), "ppt-precision-crops")
        os.makedirs(self._crop_cache_dir, exist_ok=True)
        self._shape_factory: ShapeFactory | None = None
        from pptx_designer.effects.decoration_renderer import DecorationRenderer
        from pptx_designer.renderer.layout import LayoutEngine

        self._layout_engine = LayoutEngine()
        self._decoration_renderer = DecorationRenderer()

    @property
    def brand(self) -> BrandSpec:
        return self._brand or BrandSpec()

    @property
    def shape_factory(self) -> ShapeFactory:
        if self._shape_factory is None:
            brand_colors = {}
            if self._brand and self._brand.colors:
                brand_colors = dict(self._brand.colors)
            self._shape_factory = ShapeFactory(brand_colors=brand_colors)
        return self._shape_factory

    def _c(self, role: str, fallback: str = "#000000") -> str:
        if self._brand and self._brand.colors:
            return self._brand.colors.get(role, fallback)
        return fallback

    def _font_h(self) -> str:
        if self._brand and self._brand.fonts:
            return self._brand.fonts.get("heading", "Inter")
        return "Inter"

    def _font_b(self) -> str:
        if self._brand and self._brand.fonts:
            return self._brand.fonts.get("body", "Inter")
        return "Inter"

    def _is_dark(self) -> bool:
        if self._brand and self._brand.dark_mode:
            return True
        bg_hex = self._c("background", "#FFFFFF").lstrip("#")
        r, g, b = int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16)
        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        return luminance < 0.5

    @staticmethod
    def _rgb(h: str) -> RGBColor:
        return RGBColor.from_string(h.lstrip("#"))

    def create_presentation(self) -> Presentation:
        if self._has_template:
            try:
                from pptx_designer.enterprise.slide_utils import remove_slide

                prs = Presentation(self._template_path)
                while len(prs.slides) > 0:
                    remove_slide(prs, 0)
                return prs
            except Exception:
                pass
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        return prs

    def add_slide(self, prs: Presentation, layout_name: str | None = None):
        if self._has_template:
            if layout_name:
                for layout in prs.slide_layouts:
                    if layout.name == layout_name:
                        return prs.slides.add_slide(layout)
            for layout in prs.slide_layouts:
                if "blank" in layout.name.lower():
                    return prs.slides.add_slide(layout)
            return prs.slides.add_slide(prs.slide_layouts[0])
        blank = None
        for layout in prs.slide_layouts:
            if "blank" in layout.name.lower():
                blank = layout
                break
        return prs.slides.add_slide(blank or prs.slide_layouts[-1])

    # ── Text helpers (run-level fonts) ──

    def add_text(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        font: str | None = None,
        size: int = 20,
        color_role: str = "foreground",
        color_hex: str | None = None,
        bold: bool = False,
        align: str = "left",
    ) -> object:
        font = font or self._font_h()
        color = color_hex or self._c(color_role)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = self._rgb(color)
        run.font.bold = bold
        from pptx_designer.renderer.theme_mapper import get_cjk_companion

        cjk = get_cjk_companion(font, "heading" if size >= 20 else "body")
        self._set_font_with_cjk(run, font, cjk)
        return tb

    def add_multiline(
        self,
        slide,
        lines: list[str],
        x: float,
        y: float,
        w: float,
        h: float,
        font: str | None = None,
        size: int = 14,
        color_role: str = "foreground",
        color_hex: str | None = None,
        bold: bool = False,
        align: str = "left",
        spacing: int = 6,
    ) -> object:
        font = font or self._font_b()
        color = color_hex or self._c(color_role)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        from pptx_designer.renderer.theme_mapper import get_cjk_companion

        cjk = get_cjk_companion(font, "heading" if size >= 20 else "body")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = self._rgb(color)
            run.font.bold = bold
            self._set_font_with_cjk(run, font, cjk)
            p.space_after = Pt(spacing)
        return tb

    # ── Text effects helpers (Phase 1) ──

    def add_text_with_gradient(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        text: str,
        gradient_preset: str = "gold-shine",
        gradient_stops: list[tuple[str, int]] | None = None,
        font_size: int = 44,
        bold: bool = False,
        font: str | None = None,
        align: str = "left",
    ) -> object:
        font = font or self._font_h()
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.bold = bold
        from pptx_designer.renderer.theme_mapper import get_cjk_companion

        cjk = get_cjk_companion(font, "heading" if font_size >= 20 else "body")
        self._set_font_with_cjk(run, font, cjk)
        if gradient_stops:
            apply_text_gradient(run, gradient_stops)
        else:
            apply_text_gradient_preset(run, gradient_preset)
        return tb

    def add_vertical_text(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        text: str,
        direction: str = "ea",
        font_name: str = "STKaiti",
        font_size: int = 24,
        color_role: str = "foreground",
        color_hex: str | None = None,
        bold: bool = False,
    ) -> object:
        color = color_hex or self._c(color_role)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        set_vertical_text(tf, direction)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = self._rgb(color)
        run.font.bold = bold
        from pptx_designer.renderer.theme_mapper import get_cjk_companion

        cjk = get_cjk_companion(font_name, "heading" if font_size >= 20 else "body")
        self._set_font_with_cjk(run, font_name, cjk)
        return tb

    def add_seal_stamp(
        self,
        slide,
        x: float,
        y: float,
        size: float,
        text: str,
        fill_hex: str = "#C41E3A",
        font_name: str = "SimSun",
        rotation: float = -15,
    ) -> object:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._rgb(fill_hex)
        border_hex = self._lighten(fill_hex.lstrip("#"), 40)
        sh.line.color.rgb = self._rgb(border_hex)
        sh.line.width = Pt(2)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(max(11, int(size * 18)))
        run.font.color.rgb = self._rgb("#FFFFFF")
        run.font.bold = True
        apply_text_outline(run, fill_hex.lstrip("#"), 1.5)
        if rotation != 0:
            set_text_rotation(sh, rotation)
        return sh

    # ── Image effects helpers (Phase 2) ──

    def add_circle_image(
        self, slide, cx: float, cy: float, radius: float, image_path: str, border_hex: str | None = None
    ) -> object:
        if not image_path or not os.path.isfile(image_path):
            return None
        border = border_hex or self._c("border", None)
        return _add_circle_image(slide, cx, cy, radius, image_path, border_hex=border)

    def add_image_with_soft_edge(
        self, slide, image_path: str, x: float, y: float, w: float, h: float, radius_pt: float = 10
    ) -> object:
        if not image_path or not os.path.isfile(image_path):
            return None
        shape = slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(w), Inches(h))
        apply_soft_edge(shape, radius_pt=radius_pt)
        return shape

    def add_image_with_duotone(
        self,
        slide,
        image_path: str,
        x: float,
        y: float,
        w: float,
        h: float,
        color1: str = "#0000FF",
        color2: str = "#FF0000",
    ) -> object:
        if not image_path or not os.path.isfile(image_path):
            return None
        shape = _add_image_in_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, image_path)
        apply_blip_duotone(shape, color1, color2)
        return shape

    def add_image_with_artistic(
        self,
        slide,
        image_path: str,
        x: float,
        y: float,
        w: float,
        h: float,
        effect: str = "watercolor_sponge",
        params: dict | None = None,
    ) -> object:
        if not image_path or not os.path.isfile(image_path):
            return None
        shape = _add_image_in_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, image_path)
        apply_blip_artistic(shape, effect, params)
        return shape

    # ── 3D / Pattern / Frosted helpers (Phase 4) ──

    def add_3d_shape(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        depth_pt: float = 10.0,
        material: str = "powder",
        extrusion_color: str = "#000000",
        shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
    ) -> object:
        sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        apply_3d(sh, depth_pt=depth_pt, material=material, extrusion_color=extrusion_color)
        return sh

    def add_bevel_shape(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        top_w: float = 4.0,
        top_h: float = 2.0,
        material: str = "powder",
        shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
    ) -> object:
        sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        apply_bevel(sh, top_w=top_w, top_h=top_h, material=material)
        return sh

    def add_pattern_fill_shape(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        pattern_type: str,
        fg_color: str,
        bg_color: str,
        fg_alpha: int | None = None,
        shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
    ) -> object:
        sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        apply_pattern_fill(sh, pattern_type, fg_color, bg_color, fg_alpha=fg_alpha)
        return sh

    def add_frosted_panel(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        tint_color: str = "#FFFFFF",
        tint_alpha: int = 50,
        soft_edge: float = 8,
    ) -> object:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        apply_frosted_glass(sh, tint_color=tint_color, tint_alpha=tint_alpha, soft_edge=soft_edge)
        return sh

    # ── Image helpers (Pillow pre-crop) ──

    def add_image(self, slide, path: str, x: float, y: float, w: float, h: float) -> None:
        if not os.path.isfile(path):
            return
        with PILImage.open(path) as img:
            iw, ih = img.size
            box_ratio = w / h
            img_ratio = iw / ih
            if img_ratio > box_ratio:
                cw, ch = int(ih * box_ratio), ih
                cl, ct = (iw - cw) // 2, 0
            else:
                cw, ch = iw, int(iw / box_ratio)
                cl, ct = 0, (ih - ch) // 2
            cropped = img.crop((cl, ct, cl + cw, ct + ch))
            cp = os.path.join(self._crop_cache_dir, hashlib.md5(f"{path}:{w}x{h}".encode()).hexdigest() + ".png")
            if not os.path.isfile(cp):
                cropped.save(cp, "PNG")
        slide.shapes.add_picture(cp, Inches(x), Inches(y), Inches(w), Inches(h))

    # ── Shape helpers ──

    def add_rect(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        fill_role: str | None = None,
        fill_hex: str | None = None,
        border_role: str | None = None,
        border_hex: str | None = None,
        gradient: bool = False,
        shadow: bool = False,
    ) -> object:
        fill = fill_hex or self._c(fill_role or "muted")
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if gradient:
            apply_gradient(sh, self._lighten(fill), fill, gradient_type="linear", angle=5400000)
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=15)
        border = border_hex or (self._c(border_role) if border_role else None)
        if border:
            sh.line.color.rgb = self._rgb(border)
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def add_rounded_rect(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        fill_role: str | None = None,
        fill_hex: str | None = None,
        border_role: str | None = None,
        border_hex: str | None = None,
        gradient: bool = False,
        shadow: bool = False,
        corner_radius: str | int = "md",
    ) -> object:
        fill = fill_hex or self._c(fill_role or "muted")
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        radius_pt = CORNER_RADIUS_SCALE.get(corner_radius, 8) if isinstance(corner_radius, str) else corner_radius
        min_dim = min(Inches(w), Inches(h))
        adj_val = int((radius_pt * 12700) / (min_dim / 2) * 100000) if min_dim > 0 else 16667
        adj_val = min(100000, max(0, adj_val))
        spPr = sh._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            for existing_gd in avLst.findall(qn("a:gd")):
                if existing_gd.get("name") == "adj":
                    existing_gd.set("fmla", f"val {adj_val}")
                    break
            else:
                gd = etree.SubElement(avLst, qn("a:gd"))
                gd.set("name", "adj")
                gd.set("fmla", f"val {adj_val}")
        if gradient:
            apply_gradient(sh, self._lighten(fill), fill, gradient_type="linear", angle=2700000)
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=15)
        border = border_hex or (self._c(border_role) if border_role else None)
        if border:
            sh.line.color.rgb = self._rgb(border)
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def add_oval(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        fill_role: str | None = None,
        fill_hex: str | None = None,
        gradient: bool = True,
        shadow: bool = True,
        label: str = "",
        font_size: int = 16,
        font_color: str | None = None,
    ) -> object:
        if fill_role is None and fill_hex is None:
            fill_role = "muted" if self._is_dark() else "primary"
        fill = fill_hex or self._c(fill_role)
        effective_font_color = font_color or (
            self._c("foreground", "#FFFFFF") if self._is_dark() else self._c("on-primary", "#FFFFFF")
        )
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        if gradient:
            apply_gradient(sh, fill, self._darken(fill), gradient_type="path")
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=25)
        sh.line.fill.background()
        if label:
            tf = sh.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = self._rgb(effective_font_color)
            run.font.bold = True
        return sh

    def add_donut(
        self,
        slide,
        x: float,
        y: float,
        size: float,
        fill_role: str | None = None,
        fill_hex: str | None = None,
        gradient: bool = True,
        shadow: bool = True,
        label: str = "",
        font_size: int = 18,
    ) -> object:
        if fill_role is None and fill_hex is None:
            fill_role = "muted" if self._is_dark() else "primary"
        fill = fill_hex or self._c(fill_role)
        font_color = self._c("foreground", "#FFFFFF") if self._is_dark() else self._c("on-primary", "#FFFFFF")
        sh = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(x), Inches(y), Inches(size), Inches(size))
        if gradient:
            apply_gradient(sh, fill, self._darken(fill), gradient_type="path")
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=6, distance_pt=3, alpha_pct=20)
        sh.line.fill.background()
        if label:
            tf = sh.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = self._rgb(font_color)
            run.font.bold = True
        return sh

    def add_hexagon(
        self,
        slide,
        x: float,
        y: float,
        size: float,
        fill_role: str | None = None,
        fill_hex: str | None = None,
        gradient: bool = True,
        shadow: bool = True,
        label: str = "",
        font_size: int = 16,
    ) -> object:
        if fill_role is None and fill_hex is None:
            fill_role = "muted" if self._is_dark() else "primary"
        fill = fill_hex or self._c(fill_role)
        font_color = self._c("on-primary", "#FFFFFF") if not self._is_dark() else self._c("foreground", "#FFFFFF")
        sh = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(size), Inches(size * 0.87))
        if gradient:
            apply_gradient(sh, fill, self._darken(fill), gradient_type="linear", angle=5400000)
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=25)
        sh.line.fill.background()
        if label:
            tf = sh.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = self._rgb(font_color)
            run.font.bold = True
        return sh

    @staticmethod
    def _lighten(hex_color: str, amount: int = 30) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = min(255, r + amount)
        g = min(255, g + amount)
        b = min(255, b + amount)
        return f"{r:02X}{g:02X}{b:02X}"

    @staticmethod
    def _darken(hex_color: str, amount: int = 30) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = max(0, r - amount)
        g = max(0, g - amount)
        b = max(0, b - amount)
        return f"{r:02X}{g:02X}{b:02X}"

    # ── Overlay helpers ──

    def add_dark_overlay(self, slide, opacity: float = 0.65) -> None:
        bg_hex = self._c("background", "#060B18" if self._is_dark() else "#000000")
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT))
        ov.fill.solid()
        ov.fill.fore_color.rgb = self._rgb(bg_hex)
        ov.line.fill.background()
        el = ov._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if el is not None:
            a = etree.SubElement(el, qn("a:alpha"))
            a.set("val", str(int(opacity * 100000)))

    def add_gradient_overlay(
        self, slide, opacity_bottom: float = 0.72, opacity_top: float = 0.0, color_role: str = "background"
    ) -> None:
        bg_hex = self._c(color_role, "#000000")
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT))
        ov.line.fill.background()
        top_alpha = int(opacity_top * 100000)
        mid_alpha = int(opacity_bottom * 0.4 * 100000)
        bot_alpha = int(opacity_bottom * 100000)
        grad = GradientFill(
            stops=[
                GradientStop(color=bg_hex, position=0, alpha=top_alpha),
                GradientStop(color=bg_hex, position=40000, alpha=mid_alpha),
                GradientStop(color=bg_hex, position=100000, alpha=bot_alpha),
            ],
            angle=5400000,
        )
        grad.apply(ov)

    def add_overlay(
        self, slide, x: float, y: float, w: float, h: float, color_hex: str = "#000000", opacity: float = 0.65
    ) -> None:
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        ov.fill.solid()
        ov.fill.fore_color.rgb = self._rgb(color_hex)
        ov.line.fill.background()
        el = ov._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if el is not None:
            a = etree.SubElement(el, qn("a:alpha"))
            a.set("val", str(int(opacity * 100000)))

    # ── Brand visual design ──

    def render_slide(
        self,
        prs: Presentation,
        page: dict[str, Any],
        layout_variant: dict | None = None,
        page_index: int = 0,
        total_pages: int = 0,
    ) -> object:
        elements = page.get("elements")
        if elements:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            self._render_elements(slide, elements)
            return slide

        goal = page.get("goal", "content")
        title = page.get("title", "")
        subtitle = page.get("subtitle")
        bullets = page.get("bullets") or []
        image_path = page.get("image")
        cards = page.get("cards") or []
        diagram_type = page.get("diagram_type")
        diagram_data = page.get("diagram_data")
        svg_diagram = page.get("svg_diagram")
        code = page.get("code")
        exercise = page.get("exercise")
        chart = page.get("chart")
        notes = page.get("notes")
        links = page.get("links") or []
        _image_grid = page.get("image_grid")
        _icons = page.get("icons")
        explicit_layout = page.get("layout")

        if explicit_layout:
            goal = self._remap_layout_to_goal(explicit_layout, goal)

        is_hero = goal in ("hook", "cta")
        is_section = goal == "section"

        variant = layout_variant or {}
        margin_left = variant.get("content_margin_left", 0.9)
        title_align = variant.get("title_alignment", "left")
        decoration_style = variant.get("decoration_style", "accent-bar")
        if isinstance(page, dict):
            deco = page.get("decoration")
            if isinstance(deco, dict):
                decoration_style = deco.get("style", decoration_style)

        layout_name = page.get("template_layout_name")
        slide = self.add_slide(prs, layout_name=layout_name)

        if is_section:
            section_num = page.get("section_number")
            if section_num is None:
                section_num = page_index + 1
            section_sub = subtitle or ""
            self.render_section_divider(slide, section_num, title, section_sub)
            return slide

        if is_hero:
            blocks = page.get("blocks")
            has_blocks = bool(blocks)
            has_image = bool(image_path and os.path.isfile(image_path))
            self.apply_hero_overlay(slide, prs, image_path=image_path if has_image else None)
            if has_blocks:
                self._render_hero_content(slide, title, subtitle, bullets, has_image=has_image, compact=True)
                self._render_blocks(slide, blocks, is_hero=True)
            else:
                self._render_hero_content(slide, title, subtitle, bullets, has_image=has_image)
            self._render_standalone_links(slide, links)
        else:
            self._draw_sidebar(slide, variant)
            deco_cfg = variant.get("decoration") or {}
            self.apply_brand_background(
                slide, prs, goal=goal, page_index=page_index, total_pages=total_pages, decoration=deco_cfg
            )

            sidebar_side = (variant or {}).get("sidebar_side")
            sidebar_width = (variant or {}).get("sidebar_width", 0)
            if sidebar_side and sidebar_width and sidebar_side == "left":
                margin_left = max(margin_left, sidebar_width + 0.6)

            cx, cy, cw, ch = self._content_rect(margin_left, has_subtitle=bool(subtitle))

            if sidebar_side == "right" and sidebar_width:
                cw = max(4.0, SLIDE_WIDTH - sidebar_width - margin_left - LAYOUT_GRID["margin_right"] - 0.6)

            has_content_image = bool(image_path and os.path.isfile(image_path))
            if has_content_image:
                img_w = min(4.0, cw * 0.34)
                content_cw = cw - img_w - 0.4
            else:
                img_w = 0
                content_cw = cw

            if has_content_image:
                self._render_content_image(slide, image_path, cx, cy, cw, ch)

            if title:
                self._render_title_band(slide, title, subtitle, margin_left, title_align, decoration_style)

            blocks = page.get("blocks")
            if blocks:
                self._render_blocks(slide, blocks)
            elif cards:
                self._render_cards(slide, cards, cx, cy, content_cw, ch)
            elif diagram_type and diagram_data:
                self._render_diagram_on_slide(slide, diagram_type, diagram_data, cx, cy, content_cw, ch)
            elif svg_diagram:
                self._render_svg_diagram_on_slide(slide, svg_diagram, cx, cy, content_cw, ch)
            elif code:
                self._render_code_on_slide(slide, code, cx, cy, content_cw, ch)
            elif exercise:
                self._render_exercise_on_slide(slide, exercise, cx, cy, content_cw, ch)
            elif bullets:
                self._render_bullets_on_slide(slide, bullets, cx, cy, content_cw, ch, links=links)

            if chart:
                self._render_chart_on_slide(slide, chart, cx, cy, content_cw if has_content_image else cw, ch)

            self._render_footer(slide, page_index, total_pages)

        if notes:
            self._render_notes_on_slide(slide, notes)

        return slide

    # ── Unified layout helpers ──

    def _content_rect(self, margin_left: float, has_subtitle: bool = False) -> tuple[float, float, float, float]:
        x = margin_left
        y = LAYOUT_GRID["content_y"] if has_subtitle else LAYOUT_GRID["content_y"] - 0.35
        w = SLIDE_WIDTH - margin_left - LAYOUT_GRID["margin_right"]
        h = LAYOUT_GRID["content_bottom"] - y
        return (x, y, w, h)

    def _draw_sidebar(self, slide, variant: dict | None) -> None:
        variant = variant or {}
        side = variant.get("sidebar_side")
        width = variant.get("sidebar_width")
        if not side or not width:
            return
        x = SLIDE_WIDTH - width if side == "right" else 0.0
        primary = self._c("primary", "#2563EB")
        self.add_rect(slide, x, 0, width, SLIDE_HEIGHT, fill_hex=primary)
        accent = self._c("accent", "#F97316")
        edge_x = x + width - 0.05 if side == "right" else x
        self.add_rect(slide, edge_x, 0, 0.05, SLIDE_HEIGHT, fill_hex=accent)

    def _render_title_band(
        self, slide, title: str, subtitle, margin_left: float, title_align: str, decoration_style: str
    ) -> None:
        align = "center" if title_align == "center" else "left"
        tw = SLIDE_WIDTH - margin_left - LAYOUT_GRID["margin_right"]
        tx = margin_left if align == "left" else (SLIDE_WIDTH - tw) / 2
        self.add_text(
            slide,
            title,
            tx,
            LAYOUT_GRID["title_y"],
            tw,
            0.9,
            size=LAYOUT_GRID["title_size"],
            color_role="foreground",
            bold=True,
            align=align,
        )

        colors = {
            "primary": self._c("primary", "#2563EB"),
            "accent": self._c("accent", "#F97316"),
            "foreground": self._c("foreground", "#1E293B"),
            "muted": self._c("muted", "#F1F5F9"),
            "border": self._c("border", "#E2E8F0"),
        }
        self._decoration_renderer.apply_title_decoration(
            slide,
            tx,
            LAYOUT_GRID["title_y"],
            tw,
            decoration_style,
            colors,
            add_rect_fn=lambda sl, x, y, w, h, **kw: self.add_rect(sl, x, y, w, h, **kw),
            add_oval_fn=lambda sl, x, y, w, h, **kw: self.add_oval(sl, x, y, w, h, **kw),
            add_text_fn=lambda sl, t, x, y, w, h, **kw: self.add_text(sl, t, x, y, w, h, **kw),
            apply_glow_fn=apply_glow,
        )
        if subtitle:
            self.add_text(
                slide,
                subtitle,
                tx,
                LAYOUT_GRID["subtitle_y"],
                tw,
                0.5,
                font=self._font_b(),
                size=14,
                color_role="muted-foreground",
                align=align,
            )

    def _render_hero_content(self, slide, title, subtitle, bullets, has_image: bool, compact: bool = False) -> None:
        main_role = "on-primary"  # white on gradient or scrim; scrim guarantees contrast
        if compact:
            if title:
                self.add_text(
                    slide,
                    title,
                    0.75,
                    0.6,
                    SLIDE_WIDTH - 1.5,
                    0.9,
                    size=30,
                    color_role=main_role,
                    bold=True,
                    align="left",
                )
            if subtitle:
                self.add_text(
                    slide,
                    subtitle,
                    0.75,
                    1.5,
                    SLIDE_WIDTH - 1.5,
                    0.5,
                    font=self._font_b(),
                    size=14,
                    color_role=main_role,
                    align="left",
                )
            return

        # ── Tier 4 — Hero/CTA adaptive vertical layout ──
        # Canvas: 13.333×7.5". Use 3-zone vertical distribution:
        #   top padding (10%), title block (28-32%), subtitle (8-12%),
        #   bullets (16-22%), bottom padding (15-25%)
        accent = self._c("accent", self._c("primary", "#2563EB"))
        SLIDE_TOP_Y = 0.55  # top safe margin (above title)
        SLIDE_BOTTOM_Y = SLIDE_HEIGHT - 0.55  # bottom safe margin
        USABLE_H = SLIDE_BOTTOM_Y - SLIDE_TOP_Y  # ≈ 6.4

        # Component heights — adapt to presence/absence
        components = []
        if title:
            components.append(("title", 1.75))  # baseline 1.75" tall
        if subtitle:
            components.append(("subtitle", 0.85))
        if bullets:
            n = min(4, len(bullets))
            components.append(("bullets", 0.42 * n))

        if not components:
            return

        # Reserve ratiometric heights, then compute y positions
        total_req = sum(h for _, h in components)
        if total_req < USABLE_H:
            extra = USABLE_H - total_req
            # Distribute 60% of extra as padding between components
            pad_between = (extra * 0.6) / max(1, len(components) - 1) if len(components) > 1 else 0
            bottom_pad = (extra * 0.4) / 2  # split 20% to top, 20% to bottom
            top_pad = bottom_pad
        else:
            # Doesn't fit — compress
            scale = USABLE_H / total_req
            components = [(name, h * scale) for name, h in components]
            pad_between = 0.15
            top_pad = 0.05
            bottom_pad = 0.0

        # Compute absolute y for each component
        y = SLIDE_TOP_Y + top_pad
        positions = []
        for name, h in components:
            positions.append((name, y, h))
            y += h + pad_between

        # Render
        for name, py, ph in positions:
            if name == "title":
                tx = 1.5
                tw = SLIDE_WIDTH - 3.0
                self.add_text(slide, title, tx, py, tw, ph, size=44, color_role=main_role, bold=True, align="center")
            elif name == "subtitle":
                sx = 2.0
                sw = SLIDE_WIDTH - 4.0
                self.add_text(
                    slide, subtitle, sx, py, sw, ph, font=self._font_b(), size=18, color_role=main_role, align="center"
                )
            elif name == "bullets":
                n = min(4, len(bullets))
                bullet_lines = [b for b in bullets[:n]]
                bx = 2.5
                bw = SLIDE_WIDTH - 5.0
                size = 14
                # If only 1-2 bullets, make them feel weighty by adding divider
                if n <= 2 and ph > 1.0:
                    # Add accent divider above bullets (≤2 bullets case)
                    self.add_rect(slide, (SLIDE_WIDTH - 1.2) / 2, py, 1.2, 0.04, fill_hex=accent)
                    py += 0.18
                self.add_multiline(
                    slide,
                    bullet_lines,
                    bx,
                    py,
                    bw,
                    ph,
                    size=size,
                    color_role="muted-foreground",
                    spacing=8,
                    align="center",
                )

    def _render_standalone_links(self, slide, links: list | None) -> None:
        """Render standalone link buttons (text + url + position) for hero/cta pages.

        Links without bullet_index are drawn as accent-colored underlined text
        at the given position (default bottom_right).
        """
        if not links:
            return
        from pptx_designer.renderer.theme_mapper import get_cjk_companion

        body_font = self._font_b()
        cjk = get_cjk_companion(body_font, "body")
        accent = self._c("accent", self._c("primary", "#2563EB"))
        for lnk in links:
            if lnk.get("bullet_index") is not None:
                continue  # handled by bullet renderer
            url = lnk.get("url")
            text = lnk.get("text") or url or "Link"
            if not url:
                continue
            position = lnk.get("position", "bottom_right")
            pos = {
                "bottom_right": (SLIDE_WIDTH - 4.6, SLIDE_HEIGHT - 1.1, 3.8, 0.5),
                "bottom_left": (0.8, SLIDE_HEIGHT - 1.1, 3.8, 0.5),
                "center": ((SLIDE_WIDTH - 4.0) / 2, 5.6, 4.0, 0.5),
            }.get(position, (SLIDE_WIDTH - 4.6, SLIDE_HEIGHT - 1.1, 3.8, 0.5))
            x, y, w, h = pos
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = text
            r.font.name = body_font
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = self._rgb(accent)
            r.font.underline = True
            self._set_font_with_cjk(r, body_font, cjk)
            r.hyperlink.address = url

    def _render_cards(self, slide, cards, cx: float, cy: float, cw: float, ch: float) -> None:
        n = len(cards)
        if n <= 0:
            return
        gap = 0.35
        card_w = min((cw - gap * (n - 1)) / n, 4.2)
        total_w = n * card_w + gap * (n - 1)
        x0 = cx + (cw - total_w) / 2
        for i, card in enumerate(cards):
            xx = x0 + i * (card_w + gap)
            card_title = card.get("title", "")
            card_body = card.get("text", card.get("body", ""))
            self.add_card(slide, xx, cy, card_w, ch, card_title, card_body, featured=(i == 0))

    def _render_content_image(self, slide, image_path, cx: float, cy: float, cw: float, ch: float) -> None:
        img_w = min(4.0, cw * 0.34)
        self.add_masked_image(slide, image_path, cx + cw - img_w, cy, img_w, ch)

    def _render_footer(self, slide, page_index: int, total_pages: int) -> None:
        if total_pages <= 0:
            return
        line_y = SLIDE_HEIGHT - 0.05
        self.add_rect(slide, 0, line_y, SLIDE_WIDTH, 0.05, fill_hex=self._c("border", "#E2E8F0"))
        accent = self._c("accent", self._c("primary", "#2563EB"))
        fill_w = SLIDE_WIDTH * ((page_index + 1) / total_pages)
        self.add_rect(slide, 0, line_y, fill_w, 0.05, fill_hex=accent)
        num = str(page_index + 1)
        self.add_text(
            slide,
            num,
            SLIDE_WIDTH - 0.9,
            LAYOUT_GRID["footer_y"],
            0.6,
            0.3,
            font=self._font_b(),
            size=11,
            color_role="muted-foreground",
            align="right",
        )

    def _render_blocks(self, slide, blocks: list[dict], is_hero: bool = False) -> None:
        from pptx_designer.enterprise.block_renderer import BlockRenderer

        br = BlockRenderer(self)
        br.render(slide, blocks, is_hero=is_hero)

    def _render_elements(self, slide, elements: list[dict]) -> None:
        for el in elements:
            self._dispatch_element(slide, el)

    def _dispatch_element(self, slide, el: dict) -> None:
        kind = el.get("type")
        x = el.get("x", 0.0)
        y = el.get("y", 0.0)
        w = el.get("w", 4.0)
        h = el.get("h", 1.0)

        if kind == "text":
            color = el.get("color")
            color_role = el.get("color_role")
            effective_color = color or self._c(color_role or "foreground", "#1A1A1A")
            self.add_text(
                slide,
                el.get("text", ""),
                x,
                y,
                w,
                h,
                font=el.get("font", self._font_h()),
                size=el.get("size", 18),
                color_hex=effective_color,
                bold=el.get("bold", False),
                align=el.get("align", "left"),
            )

        elif kind == "multiline":
            lines = el.get("lines", [])
            color = el.get("color")
            color_role = el.get("color_role")
            effective_color = color or self._c(color_role or "foreground", "#1A1A1A")
            self.add_multiline(
                slide,
                lines,
                x,
                y,
                w,
                h,
                font=el.get("font", self._font_b()),
                size=el.get("size", 14),
                color_hex=effective_color,
                bold=el.get("bold", False),
                align=el.get("align", "left"),
                spacing=el.get("spacing", 6),
            )

        elif kind == "rect":
            fill = el.get("fill")
            fill_role = el.get("fill_role")
            border = el.get("border")
            border_role = el.get("border_role")
            self.add_rect(
                slide,
                x,
                y,
                w,
                h,
                fill_hex=fill,
                fill_role=fill_role,
                border_hex=border,
                border_role=border_role,
                gradient=el.get("gradient", False),
                shadow=el.get("shadow", False),
            )

        elif kind == "rounded_rect":
            fill = el.get("fill")
            fill_role = el.get("fill_role")
            border = el.get("border")
            border_role = el.get("border_role")
            self.add_rounded_rect(
                slide,
                x,
                y,
                w,
                h,
                fill_hex=fill,
                fill_role=fill_role,
                border_hex=border,
                border_role=border_role,
                gradient=el.get("gradient", False),
                shadow=el.get("shadow", False),
                corner_radius=el.get("radius", "md"),
            )

        elif kind == "image":
            path = el.get("path", "")
            if os.path.isfile(path):
                self.add_image(slide, path, x, y, w, h)

        elif kind == "overlay":
            color = el.get("color") or self._c("background", "#000000")
            opacity = el.get("opacity", 0.65)
            self.add_overlay(slide, x, y, w, h, color, opacity)

        elif kind == "gradient_line":
            color = el.get("color") or self._c("accent", self._c("primary", "#2563EB"))
            self.add_gradient_line(slide, x, y, w, h, color)

        elif kind == "hexagon":
            size = el.get("size", min(w, h))
            fill = el.get("fill")
            fill_role = el.get("fill_role", "primary")
            self.add_hexagon(
                slide,
                x,
                y,
                size,
                fill_hex=fill,
                fill_role=fill_role,
                gradient=el.get("gradient", True),
                shadow=el.get("shadow", True),
                label=el.get("label", ""),
                font_size=el.get("font_size", 16),
            )

        elif kind == "oval":
            fill = el.get("fill")
            fill_role = el.get("fill_role", "primary")
            self.add_oval(
                slide,
                x,
                y,
                w,
                h,
                fill_hex=fill,
                fill_role=fill_role,
                gradient=el.get("gradient", True),
                shadow=el.get("shadow", True),
                label=el.get("label", ""),
                font_size=el.get("font_size", 16),
            )

        elif kind == "donut":
            size = el.get("size", min(w, h))
            fill = el.get("fill")
            fill_role = el.get("fill_role", "primary")
            self.add_donut(
                slide,
                x,
                y,
                size,
                fill_hex=fill,
                fill_role=fill_role,
                gradient=el.get("gradient", True),
                shadow=el.get("shadow", True),
                label=el.get("label", ""),
                font_size=el.get("font_size", 18),
            )

    def _render_bullets_on_slide(
        self,
        slide,
        bullets: list,
        cx: float = 0.75,
        cy: float = 2.05,
        cw: float = 11.83,
        ch: float = 4.8,
        links: list | None = None,
    ) -> None:
        accent = self._c("accent", self._c("primary", "#2563EB"))
        links = links or []

        # ── Tier 4 — Auto-stack: split content area into zones ──
        # If bullets are few (≤4), use a single column with generous sizing
        # and decorative fill below the last bullet to occupy the vertical space.
        if len(bullets) <= 4:
            # Use full width, single column
            size = self._fit_bullet_size(bullets, cw, ch, base=16, target_fill=0.75, lo=14, hi=24)
            self._render_bullet_column(
                slide, bullets[:8], cx, cy, cw, ch, accent, links=links, link_base_index=0, override_size=size
            )
            # After rendering, check if bottom of content area is near-empty
            # and fill with a subtle decorative element
            total_h = self._estimate_bullet_height(bullets, cw, size, max(11, int(size * 0.9)))
            bottom_gap = ch - total_h
            if bottom_gap > 0.6:
                # Place a subtle accent separator at ~70% of the content area
                sep_y = cy + total_h + 0.2
                self.add_rect(slide, cx + 1.0, sep_y, cw - 2.0, 0.03, fill_hex=accent)
            return

        if len(bullets) >= 6:
            mid = (len(bullets) + 1) // 2
            col_gap = 0.6
            col_w = (cw - col_gap) / 2
            self._render_bullet_column(slide, bullets[:mid], cx, cy, col_w, ch, accent, links=links, link_base_index=0)
            self._render_bullet_column(
                slide, bullets[mid:], cx + col_w + col_gap, cy, col_w, ch, accent, links=links, link_base_index=mid
            )
        else:
            self._render_bullet_column(slide, bullets[:8], cx, cy, cw, ch, accent, links=links, link_base_index=0)

    def _estimate_bullet_height(self, bullets: list, w: float, size: int, spacing: int) -> float:
        """Rough text-height estimate (CJK-aware) to drive adaptive sizing."""
        em = size / 72.0
        cpl = max(6, int(w / (em * 0.98)))
        lines = 0
        for b in bullets:
            lines += max(1, (len(b) + cpl - 1) // cpl)
        line_h = size / 72.0 * 1.28
        sp = spacing / 72.0
        return lines * line_h + max(0, len(bullets) - 1) * sp

    def _fit_bullet_size(
        self, bullets: list, w: float, h: float, base: int = 15, target_fill: float = 0.66, lo: int = 13, hi: int = 20
    ) -> int:
        size = base
        for _ in range(14):
            spacing = max(10, int(size * 0.8))
            est = self._estimate_bullet_height(bullets, w, size, spacing)
            if est > h * 0.96:
                size -= 1
            elif est < h * target_fill and size < hi:
                size += 1
            else:
                break
        return max(lo, min(hi, size))

    def _render_bullet_column(
        self,
        slide,
        bullets: list,
        x: float,
        y: float,
        w: float,
        h: float,
        accent: str,
        links: list | None = None,
        link_base_index: int = 0,
        override_size: int | None = None,
    ) -> object:
        from pptx_designer.renderer.theme_mapper import get_cjk_companion

        size = override_size if override_size else self._fit_bullet_size(bullets, w, h)
        spacing = max(11, int(size * 0.9))
        total_h = self._estimate_bullet_height(bullets, w, size, spacing)
        # ── Tier 4 — clamp textbox to available height (never overflow slide) ──
        # Older code sized the textbox to full `h` and centered it, which pushed
        # the bottom edge past SLIDE_HEIGHT when total_h was small.
        top = y + max(0.0, (h - total_h) / 2)
        box_h = max(total_h, 0.4)
        max_top = SLIDE_HEIGHT - 0.4 - box_h
        if top > max_top:
            top = max_top

        # map bullet_index (global) → url
        link_map: dict[int, str] = {}
        for lnk in links or []:
            if lnk.get("bullet_index") is not None and lnk.get("url"):
                link_map[int(lnk["bullet_index"])] = lnk["url"]

        tb = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(w), Inches(box_h))
        tf = tb.text_frame
        tf.word_wrap = True
        body_font = self._font_b()
        cjk = get_cjk_companion(body_font, "body")
        text_color = self._c("foreground", "#1E293B")
        accent_rgb = self._rgb(accent)
        marker = "\u25aa  "
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(spacing)
            r1 = p.add_run()
            r1.text = marker
            r1.font.name = body_font
            r1.font.size = Pt(size)
            r1.font.color.rgb = accent_rgb
            r1.font.bold = True
            self._set_font_with_cjk(r1, body_font, cjk)
            r2 = p.add_run()
            r2.text = b
            r2.font.name = body_font
            r2.font.size = Pt(size)
            r2.font.color.rgb = self._rgb(text_color)
            self._set_font_with_cjk(r2, body_font, cjk)
            url = link_map.get(link_base_index + i)
            if url:
                r2.hyperlink.address = url
                r2.font.color.rgb = self._rgb(accent)
                r2.font.underline = True
        return tb

    def _compute_content_area(self, page: dict[str, Any]) -> tuple[float, float, float, float]:
        goal = page.get("goal", "content")
        has_image = bool(page.get("image"))
        has_title = bool(page.get("title"))
        has_subtitle = bool(page.get("subtitle"))

        slide_w = SLIDE_WIDTH
        slide_h = SLIDE_HEIGHT

        margin_left = 0.9
        margin_right = 0.4
        margin_bottom = 0.5

        top = 0.5
        if has_title:
            top += 0.9
        if has_subtitle:
            top += 0.5

        right_edge = slide_w - margin_right
        if has_image and goal not in ("hook", "cta"):
            right_edge = 8.0

        content_w = right_edge - margin_left
        content_h = slide_h - top - margin_bottom

        return (margin_left, top, content_w, content_h)

    def _remap_layout_to_goal(self, layout_name: str, fallback_goal: str) -> str:
        layout_to_goal = {
            "title-slide": "hook",
            "cta-closing": "cta",
            "content-with-title": "content",
            "two-column": "content",
            "three-column-cards": "features",
            "four-metrics": "traction",
            "big-number": "agitation",
            "quote": "testimonials",
            "chart-focus": "data",
            "image-plus-text": "content",
            "sidebar-left": "overview",
            "exercise-layout": "exercise",
            "code-block": "code",
            "funnel": "funnel",
            "grid-2x2-cards": "features",
            "dense-bullets": "content",
            "swot-matrix": "content",
            "cycle-diagram": "content",
            "timeline-horizontal": "content",
            "table-layout": "content",
            "section-header": "content",
            "blank": "content",
        }
        return layout_to_goal.get(layout_name, fallback_goal)

    def _render_diagram_on_slide(
        self,
        slide,
        diagram_type: str,
        diagram_data: dict,
        cx: float | None = None,
        cy: float | None = None,
        cw: float | None = None,
        ch: float | None = None,
    ) -> None:
        from pptx_designer.renderer.diagram.diagram_style import DiagramStyle
        from pptx_designer.renderer.diagram.layout_engine import Region
        from pptx_designer.renderer.diagram_engine import DiagramEngine

        if cx is None:
            cx, cy, cw, ch = self._content_rect(LAYOUT_GRID["margin_left"])
        style = DiagramStyle.from_brand_spec(self.brand) if self._brand else DiagramStyle()
        region = Region(left=cx, top=cy, width=cw, height=ch)
        engine = DiagramEngine()
        try:
            engine.render(slide, diagram_type, diagram_data, style, region)
        except Exception as e:
            logger.error("Diagram engine failed for %s: %s", diagram_type, e)

    def _render_svg_diagram_on_slide(
        self,
        slide,
        svg_diagram,
        cx: float | None = None,
        cy: float | None = None,
        cw: float | None = None,
        ch: float | None = None,
    ):
        from pptx_designer.compiler import SVGCompileError, SVGCompiler

        if cx is None:
            cx, cy, cw, ch = self._content_rect(LAYOUT_GRID["margin_left"])
        svg_text = svg_diagram if isinstance(svg_diagram, str) else svg_diagram.get("svg", "")
        if not svg_text:
            return None
        C = self._build_svg_context()
        try:
            result = SVGCompiler(C=C).compile(svg_text, slide, (cx, cy, cw, ch))
            if result.warnings:
                for w in result.warnings:
                    logger.warning("SVG compiler warning: %s", w)
            return result
        except SVGCompileError as e:
            logger.error("SVG compilation failed: %s", e)
            return None

    def _build_svg_context(self) -> dict:
        C = {}
        if self._brand and self._brand.colors:
            C.update(self._brand.colors)
        elif getattr(self, "_colors", None) or getattr(self, "colors", None):
            C.update(getattr(self, "_colors", None) or getattr(self, "colors", None))
        if self._brand and self._brand.fonts:
            C["font_heading"] = self._brand.fonts.get("heading", "Inter")
            C["font_body"] = self._brand.fonts.get("body", "Inter")
        if self._brand and self._brand.dark_mode:
            C["dark_mode"] = True
        return C

    def _render_code_on_slide(
        self,
        slide,
        code_data,
        cx: float | None = None,
        cy: float | None = None,
        cw: float | None = None,
        ch: float | None = None,
    ) -> None:
        if cx is None:
            cx, cy, cw, ch = self._content_rect(LAYOUT_GRID["margin_left"])
        code_text = code_data if isinstance(code_data, str) else code_data.get("source", code_data.get("code", ""))
        language = code_data.get("language", "") if isinstance(code_data, dict) else ""
        code_bg = "#1E293B"
        self.add_rounded_rect(slide, cx, cy, cw, ch, fill_hex=code_bg, corner_radius="md")
        lines = code_text.split("\n")
        if language:
            badge_text = language.upper()
            badge_w = len(badge_text) * 0.1 + 0.3
            badge_y = cy + 0.15
            self.add_rounded_rect(slide, cx + 0.15, badge_y, badge_w, 0.3, fill_role="primary", corner_radius="sm")
            self.add_text(
                slide,
                badge_text,
                cx + 0.25,
                badge_y + 0.02,
                badge_w - 0.2,
                0.26,
                size=11,
                color_role="on-primary",
                bold=True,
            )
            code_lines = [f"  {line}" for line in lines[:30]]
            self.add_multiline(
                slide,
                code_lines,
                cx + 0.35,
                cy + 0.6,
                cw - 0.7,
                ch - 0.7,
                font="Consolas",
                size=12,
                color_role="muted-foreground",
                spacing=5,
            )
        else:
            code_lines = [f"  {line}" for line in lines[:30]]
            self.add_multiline(
                slide,
                code_lines,
                cx + 0.35,
                cy + 0.3,
                cw - 0.7,
                ch - 0.5,
                font="Consolas",
                size=12,
                color_role="muted-foreground",
                spacing=5,
            )

    def _render_exercise_on_slide(
        self,
        slide,
        exercise_data,
        cx: float | None = None,
        cy: float | None = None,
        cw: float | None = None,
        ch: float | None = None,
    ) -> None:
        if cx is None:
            cx, cy, cw, ch = self._content_rect(LAYOUT_GRID["margin_left"])
        instructions = exercise_data.get("instructions", "") if isinstance(exercise_data, dict) else str(exercise_data)
        duration = exercise_data.get("duration", "") if isinstance(exercise_data, dict) else ""
        steps = exercise_data.get("steps", []) if isinstance(exercise_data, dict) else []
        badge_text = f"Exercise {duration}" if duration else "Exercise"
        self.add_badge(slide, badge_text, cx, cy + 0.15, variant="solid")
        y = cy + 0.55
        if instructions:
            self.add_text(
                slide, instructions, cx, y, cw, 0.7, font=self._font_b(), size=15, color_role="muted-foreground"
            )
            y += 0.9
        if steps:
            step_lines = [f"{i + 1}.  {s}" for i, s in enumerate(steps)]
            steps_h = max(0.5, ch - (y - cy))
            self.add_multiline(slide, step_lines, cx, y, cw, steps_h, size=16, color_role="foreground", spacing=10)

    def _render_chart_on_slide(
        self,
        slide,
        chart_config: dict,
        cx: float | None = None,
        cy: float | None = None,
        cw: float | None = None,
        ch: float | None = None,
    ) -> None:
        from pptx_designer.renderer.chart_builder import ChartBuilder

        if cx is None:
            cx, cy, cw, ch = self._content_rect(LAYOUT_GRID["margin_left"])
        brand_colors = None
        brand_fonts = None
        if self._brand:
            if self._brand.colors:
                brand_colors = self._brand.colors
            if self._brand.fonts:
                brand_fonts = self._brand.fonts
        position = {"x": cx, "y": cy, "width": cw, "height": ch}
        builder = ChartBuilder()
        with suppress(Exception):
            builder.build(slide, chart_config, position=position, brand_colors=brand_colors, brand_fonts=brand_fonts)

    def _render_notes_on_slide(self, slide, notes_text: str) -> None:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    def apply_brand_background(
        self,
        slide,
        prs: Presentation,
        goal: str = "content",
        page_index: int = 0,
        total_pages: int = 0,
        decoration: dict | None = None,
    ) -> None:
        deco = decoration or {}
        bg_hex = self._c("background", "#FFFFFF")
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._rgb(bg_hex)
        except Exception:
            pass

        accent_hex = self._c("accent", self._c("primary", "#2563EB"))
        primary_hex = self._c("primary", "#2563EB")

        # Subtle gradient wash — adds depth over the flat background.
        wash_from = self._lighten(bg_hex, 22) if self._is_dark() else self._darken(bg_hex, 6)
        self.add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_hex=wash_from, gradient=True)

        # Style-driven background decoration.
        if deco.get("grid_background"):
            from pptx_designer.effects.decoration import add_grid_background

            with suppress(Exception):
                add_grid_background(slide, spacing=1.2, color=primary_hex, alpha=7)
        if deco.get("neon_accent"):
            self._add_bg_glow(slide, accent_hex)
        if deco.get("circle_decoration") or deco.get("left_accent"):
            self._add_bg_circle(slide, primary_hex)
        if deco.get("seal_decoration"):
            from pptx_designer.effects.decoration import add_seal_stamp

            with suppress(Exception):
                add_seal_stamp(
                    slide, SLIDE_WIDTH - 1.15, 0.35, 0.5, "智", fill_hex=self._lighten(accent_hex, 45), style="zhu"
                )
        if deco.get("brush_divider"):
            from pptx_designer.effects.decoration import add_brush_divider

            with suppress(Exception):
                add_brush_divider(
                    slide, 0, SLIDE_HEIGHT - 0.18, SLIDE_WIDTH, color=self._lighten(accent_hex, 20), thickness=0.06
                )
        if deco.get("left_accent"):
            self.add_rect(slide, 0, 0, 0.05, SLIDE_HEIGHT, fill_hex=accent_hex, gradient=True)

        # Corner depth accent on every content slide (subtle, low alpha).
        self._add_bg_corner_accent(slide, accent_hex, top_left=(page_index % 2 == 0))

        # Consistent top accent line — the "format anchor".
        self.add_rect(slide, 0, 0, SLIDE_WIDTH, LAYOUT_GRID["top_accent_h"], fill_hex=accent_hex)

    def _add_bg_circle(self, slide, color_hex: str) -> None:
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.8), Inches(SLIDE_HEIGHT - 3.2), Inches(4.0), Inches(4.0))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._rgb(color_hex)
        sh.line.fill.background()
        el = sh._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if el is not None:
            a = etree.SubElement(el, qn("a:alpha"))
            a.set("val", str(int(0.05 * 100000)))

    def _add_bg_corner_accent(self, slide, color_hex: str, top_left: bool) -> None:
        if top_left:
            x, y = 0.55, 0.35
        else:
            x, y = SLIDE_WIDTH - 1.75, 0.35
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.1), Inches(1.1))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._rgb(color_hex)
        sh.line.fill.background()
        el = sh._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if el is not None:
            a = etree.SubElement(el, qn("a:alpha"))
            a.set("val", str(int(0.12 * 100000)))

    def _add_bg_glow(self, slide, color_hex: str) -> None:
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.2), Inches(0.0), Inches(2.6), Inches(1.1))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._rgb(color_hex)
        sh.line.fill.background()
        el = sh._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if el is not None:
            a = etree.SubElement(el, qn("a:alpha"))
            a.set("val", str(int(0.10 * 100000)))
        with suppress(Exception):
            apply_glow(sh, radius_pt=14, color=color_hex, alpha_pct=28)

    def apply_hero_overlay(self, slide, prs: Presentation, image_path: str | None = None) -> None:
        if image_path and os.path.isfile(image_path):
            self.add_image(slide, image_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
            # Light bottom gradient so the image stays visible.
            self.add_gradient_overlay(slide, opacity_bottom=0.42, opacity_top=0.0)
            # Soft dark scrim behind the title band for text contrast.
            self.add_overlay(slide, 0, 2.0, SLIDE_WIDTH, 2.5, color_hex="#000000", opacity=0.30)
        else:
            primary_hex = self._c("primary", "#2563EB")
            self.add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_hex=self._darken(primary_hex, 12), gradient=True)
            accent_hex = self._c("accent", self._c("primary", "#2563EB"))
            # soft radial glow in a corner for depth
            sh = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(SLIDE_WIDTH - 4.5), Inches(-2.5), Inches(7.0), Inches(7.0)
            )
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(accent_hex)
            sh.line.fill.background()
            el = sh._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
            if el is not None:
                a = etree.SubElement(el, qn("a:alpha"))
                a.set("val", str(int(0.14 * 100000)))
            # bottom-left secondary glow
            sh2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(-3.0), Inches(SLIDE_HEIGHT - 4.0), Inches(6.0), Inches(6.0)
            )
            sh2.fill.solid()
            sh2.fill.fore_color.rgb = self._rgb(primary_hex)
            sh2.line.fill.background()
            el2 = sh2._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
            if el2 is not None:
                a2 = etree.SubElement(el2, qn("a:alpha"))
                a2.set("val", str(int(0.10 * 100000)))

    # ── Card component ──

    def add_card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        accent_role: str = "primary",
        featured: bool = False,
    ) -> None:
        accent_hex = self._c(accent_role)
        card_bg = self._c("muted", "#0D152A" if self._is_dark() else "#F8FAFC")
        card_bd = self._c("border", "#1A2A4A" if self._is_dark() else "#E2E8F0")
        pad = 0.32
        title_size = 22 if featured else 20
        self.add_rounded_rect(
            slide, x, y, w, h, fill_hex=card_bg, border_hex=card_bd, gradient=True, shadow=True, corner_radius="md"
        )
        bar_h = 0.08
        bar_w = (w - 2 * pad) if featured else (w - 2 * pad) * 0.4
        self.add_rect(slide, x + pad, y + pad - 0.08, bar_w, bar_h, fill_hex=accent_hex, gradient=True)
        self.add_text(
            slide, title, x + pad, y + pad + 0.1, w - 2 * pad, 0.55, self._font_h(), title_size, accent_role, bold=True
        )
        body_lines = body.split("\n")
        body_size = 15 if len(body_lines) <= 2 else 14
        # center body vertically within the remaining card space
        body_top = y + pad + 0.75
        body_h = h - pad - 0.75 - 0.2
        self.add_multiline(
            slide,
            body_lines,
            x + pad,
            body_top,
            w - 2 * pad,
            body_h,
            self._font_b(),
            body_size,
            "muted-foreground",
            spacing=6,
        )

    # ── Brand compliance: logo, footer, watermark ──

    def apply_logo(self, slide, logo_path: str, prs: Presentation, current_goal: str | None = None) -> None:
        if not self._brand or not self._brand.logo:
            return
        logo_spec = self._brand.logo
        skip_slides = logo_spec.get("skip_slides", [])
        if current_goal and current_goal in skip_slides:
            return
        if not os.path.isfile(logo_path):
            return

        position = logo_spec.get("position", "top_right")
        width_inches = logo_spec.get("width_inches", 1.0)
        width_emu = Inches(width_inches)

        try:
            with PILImage.open(logo_path) as img:
                aspect = img.size[1] / img.size[0] if img.size[0] > 0 else 0.5
        except Exception:
            aspect = 0.5
        height_emu = int(width_emu * aspect)

        slide_w, slide_h = prs.slide_width, prs.slide_height
        if position == "top_right":
            left, top = slide_w - width_emu - Inches(0.3), Inches(0.3)
        elif position == "top_left":
            left, top = Inches(0.3), Inches(0.3)
        elif position == "bottom_right":
            left, top = slide_w - width_emu - Inches(0.3), slide_h - height_emu - Inches(0.3)
        else:
            left, top = slide_w - width_emu - Inches(0.3), Inches(0.3)

        slide.shapes.add_picture(logo_path, left, top, width=width_emu, height=height_emu)

    def apply_footer(self, prs: Presentation) -> None:
        if not self._brand or not self._brand.footer:
            return
        config = self._brand.footer
        total = len(prs.slides)
        show_page_number = config.get("show_page_number", False)
        page_number_format = config.get("page_number_format", "{n}")
        page_number_position = config.get("page_number_position", "bottom_right")
        show_footer_text = config.get("show_footer_text", False)
        footer_text = config.get("footer_text", "")
        footer_position = config.get("footer_position", "bottom_center")
        font_size_pt = config.get("font_size_pt", 10)
        skip_pages = config.get("skip_pages", [])

        if not show_page_number and not show_footer_text:
            return

        muted_color = self._rgb(self._c("muted-foreground", "#999999"))
        position_map = {
            "bottom_left": Inches(0.9),
            "bottom_center": Inches(5.833),
            "bottom_right": Inches(11.433),
        }

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            is_cover = idx == 0
            should_skip = is_cover or slide_num in skip_pages

            if show_page_number and not should_skip:
                page_text = page_number_format.replace("{n}", str(slide_num)).replace("{total}", str(total))
                left = position_map.get(page_number_position, Inches(11.433))
                tb = slide.shapes.add_textbox(left, Inches(7.0), Inches(2.0), Inches(0.3))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                run = p.add_run()
                run.text = page_text
                run.font.size = Pt(font_size_pt)
                run.font.color.rgb = muted_color
                run.font.name = self._font_b()

            if show_footer_text and footer_text and not should_skip:
                left = position_map.get(footer_position, Inches(5.833))
                tb = slide.shapes.add_textbox(left, Inches(7.0), Inches(4.0), Inches(0.3))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = footer_text
                run.font.size = Pt(font_size_pt)
                run.font.color.rgb = muted_color
                run.font.name = self._font_b()

    def apply_watermark(self, prs: Presentation) -> None:
        if not self._brand or not self._brand.watermark:
            return
        config = self._brand.watermark
        text = config.get("text", "CONFIDENTIAL")
        opacity = config.get("opacity", 0.15)
        rotation = config.get("rotation", -45)
        font_size_pt = config.get("font_size_pt", 72)
        skip_pages = config.get("skip_pages", [1])

        muted_hex = self._c("muted-foreground", "#999999")

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            if slide_num in skip_pages:
                continue

            tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = self._rgb(muted_hex)
            run.font.bold = True
            run.font.name = self._font_h()

            sp = tb._element
            srgbClr_el = sp.find(".//" + qn("a:srgbClr"))
            if srgbClr_el is not None:
                alpha_elem = etree.SubElement(srgbClr_el, qn("a:alpha"))
                alpha_elem.set("val", str(int(opacity * 100000)))

            xfrm_el = sp.find(".//" + qn("a:xfrm"))
            if xfrm_el is not None:
                xfrm_el.set("rot", str(int(rotation * 60000)))

    # ── Save ──

    def save(self, prs: Presentation, output_path: str) -> None:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        prs.save(output_path)

    # ── §2.1 CJK font pairing ──

    def _set_font_with_cjk(self, run, latin_font: str, cjk_font: str | None = None) -> None:
        run.font.name = latin_font
        if cjk_font:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", cjk_font)
            cs = rPr.find(qn("a:cs"))
            if cs is None:
                cs = etree.SubElement(rPr, qn("a:cs"))
            cs.set("typeface", cjk_font)

    # ── §2.3 Badge/Tag styling ──

    def add_badge(self, slide, text: str, x: float, y: float, variant: str = "default", size: str = "sm") -> None:
        font_size = {"sm": 11, "md": 11, "lg": 12}.get(size, 11)
        badge_text = text.upper()
        cjk_count = sum(1 for ch in badge_text if "\u4e00" <= ch <= "\u9fff" or "\u3000" <= ch <= "\u303f")
        latin_count = len(badge_text) - cjk_count
        char_width = font_size * 0.009
        cjk_width = font_size * 0.016
        badge_w = latin_count * char_width + cjk_count * cjk_width + 0.3
        badge_h = 0.35
        if variant == "default":
            tint_hex = self._lighten(self._c("primary", "#2563EB"))
            self.add_rounded_rect(slide, x, y, badge_w, badge_h, fill_hex=tint_hex, corner_radius="sm")
            self.add_text(
                slide,
                badge_text,
                x + 0.1,
                y + 0.02,
                badge_w - 0.2,
                badge_h - 0.04,
                size=font_size,
                color_role="primary",
                bold=True,
            )
        elif variant == "solid":
            self.add_rounded_rect(slide, x, y, badge_w, badge_h, fill_role="primary", corner_radius="sm")
            self.add_text(
                slide,
                badge_text,
                x + 0.1,
                y + 0.02,
                badge_w - 0.2,
                badge_h - 0.04,
                size=font_size,
                color_role="on-primary",
                bold=True,
            )
        elif variant == "outline":
            self.add_rounded_rect(
                slide,
                x,
                y,
                badge_w,
                badge_h,
                fill_hex=self._c("background", "#FFFFFF"),
                border_role="primary",
                corner_radius="sm",
            )
            self.add_text(
                slide,
                badge_text,
                x + 0.1,
                y + 0.02,
                badge_w - 0.2,
                badge_h - 0.04,
                size=font_size,
                color_role="primary",
                bold=True,
            )

    # ── §2.4 Section divider ──

    def render_section_divider(
        self, slide_or_prs, section_number: int, section_title: str, section_subtitle: str = ""
    ) -> object:
        slide = slide_or_prs if hasattr(slide_or_prs, "background") else self.add_slide(slide_or_prs)
        primary_hex = self._c("primary", "#2563EB")
        bg = self._lighten(primary_hex, 115) if self._is_dark() else self._lighten(primary_hex, 75)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(bg)
        display_num = section_number if isinstance(section_number, str) else f"{section_number:02d}"
        self.add_text(
            slide,
            display_num,
            1.5,
            1.6,
            SLIDE_WIDTH - 3.0,
            2.0,
            size=72,
            color_role="primary",
            bold=True,
            align="center",
        )
        self.add_text(
            slide,
            section_title,
            1.5,
            3.6,
            SLIDE_WIDTH - 3.0,
            1.0,
            size=40,
            color_role="foreground",
            bold=True,
            align="center",
        )
        accent = self._c("accent", self._c("primary", "#2563EB"))
        bar_w = min(3.0, len(section_title) * 0.14 + 0.8)
        self.add_rect(slide, (SLIDE_WIDTH - bar_w) / 2, 4.75, bar_w, 0.045, fill_hex=accent, gradient=True)
        if section_subtitle:
            self.add_text(
                slide,
                section_subtitle,
                2.0,
                4.95,
                SLIDE_WIDTH - 4.0,
                0.5,
                font=self._font_b(),
                size=18,
                color_role="muted-foreground",
                align="center",
            )
        return slide

    # ── §3.2 Progress bar ──

    def add_progress_bar(self, slide, current: int, total: int) -> None:
        bar_y = SLIDE_HEIGHT - 0.04
        bar_h = 0.03
        self.add_rect(slide, 0, bar_y, SLIDE_WIDTH, bar_h, fill_hex=self._c("border", "#E2E8F0"))
        fill_w = SLIDE_WIDTH * (current / total)
        self.add_rect(slide, 0, bar_y, fill_w, bar_h, fill_hex=self._c("primary", "#2563EB"))

    # ── §3.4 Gradient line ──

    def add_gradient_line(
        self, slide, x: float, y: float, width: float, height: float, from_color: str, to_color: str = "transparent"
    ) -> object:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        rect.line.fill.background()
        grad = GradientFill(
            stops=[
                GradientStop(color=from_color, position=0, alpha=100000),
                GradientStop(color=from_color if to_color == "transparent" else to_color, position=100000, alpha=0),
            ],
            angle=0,
        )
        grad.apply(rect)
        return rect

    # ── §3.5 Image masking ──

    def add_masked_image(
        self,
        slide,
        image_path: str,
        x: float,
        y: float,
        w: float,
        h: float,
        padding: float = 0.15,
        corner_radius: str = "md",
    ) -> None:
        frame_bg = "#FFFFFF" if not self._is_dark() else "#1E293B"
        self.add_rounded_rect(slide, x, y, w, h, fill_hex=frame_bg, shadow=True, corner_radius=corner_radius)
        img_x = x + padding
        img_y = y + padding
        img_w = w - 2 * padding
        img_h = h - 2 * padding
        self.add_image(slide, image_path, img_x, img_y, img_w, img_h)

    # ── §3.7 Hero pattern selection ──

    def _select_hero_pattern(self, page: dict, mood: str = "professional") -> str:
        subtitle = page.get("subtitle", "")
        has_image = bool(page.get("image"))
        if not has_image:
            return "gradient"
        if len(subtitle) > 60:
            return "split-left"
        if mood in ("creative", "bold", "vibrant", "startup"):
            return "asymmetric"
        return "bottom-fade"

    _MOOD_CATEGORY_MAP: dict[str, str] = {
        "mckinsey": "hierarchy",
        "consulting": "process",
        "tech": "cycle",
        "creative": "radial",
        "dark": "chart",
        "fintech": "funnel",
        "education": "pyramid",
        "health": "cycle",
        "government": "hierarchy",
        "industrial": "process",
        "startup": "infographic",
        "luxury": "pyramid",
        "nature": "cycle",
        "minimal": "infographic",
        "bold": "comparison",
        "international": "matrix",
        "cream": "infographic",
        "frosted": "matrix",
        "pastel": "infographic",
        "retro": "timeline",
        "legal": "hierarchy",
        "pharma": "funnel",
        "realestate": "pyramid",
        "automotive": "process",
        "aviation": "process",
        "energy": "funnel",
        "telecom": "cycle",
        "logistics": "process",
    }

    def _mood_to_preferred_category(self, mood: str | None) -> str | None:
        if not mood:
            return None
        return self._MOOD_CATEGORY_MAP.get(mood)

        return None
