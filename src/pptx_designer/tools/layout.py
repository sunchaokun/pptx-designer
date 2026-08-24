"""Layout tools — slide-level layout and decoration functions for build scripts.

Provides slide management, theme control, page decorations, and background
utilities. These are the functions that operate at the slide/presentation
level rather than individual shapes.

Usage:
    from pptx_designer.tools.layout import *
    C = {'primary': '#2E6504', 'accent': '#7DA92F', ...}
    t = TYPOGRAPHY['mckinsey']
    sp = SPACING['mckinsey']
    prs = Presentation(template_path)
    s = add_slide(prs)
    page_header(s, 'Title', 'Subtitle', C, typo=t, spacing=sp)
    prs.save('output.pptx')
"""
from __future__ import annotations

import copy

from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from pptx_designer.tools.shapes import (
    _resolve_color,
    _rgb,
    _set_cjk_font,
    _strip_style,
    _add_shape,
    _set_run,
    _lighten,
    Typography,
    Spacing,
    TYPOGRAPHY,
    SPACING,
    rect,
    rrect,
    oval,
)


# ---------------------------------------------------------------------------
# Presentation-level helpers
# ---------------------------------------------------------------------------

def set_widescreen(prs):
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sldSz = prs._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz'
    )
    if sldSz is not None and 'type' in sldSz.attrib:
        del sldSz.attrib['type']


def set_dark_theme(prs, C=None):
    C = C or {}
    from lxml import etree as _et
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bg = C.get('background', '#0B1020')
    fg = C.get('text_dark', '#E2E8F0')
    theme_part = None
    for rel in prs.part.rels.values():
        if 'theme' in rel.reltype:
            theme_part = rel.target_part
            break
    if theme_part is None:
        return
    theme_el = _et.fromstring(theme_part.blob)
    clrScheme = theme_el.find(f'{{{ns}}}themeElements/{{{ns}}}clrScheme')
    if clrScheme is None:
        return
    for tag, val in [('dk1', fg), ('lt1', bg)]:
        el = clrScheme.find(f'{{{ns}}}{tag}')
        if el is not None:
            for child in list(el):
                el.remove(child)
            srgb = _et.SubElement(el, f'{{{ns}}}srgbClr')
            srgb.set('val', val.lstrip('#'))
    theme_part._blob = _et.tostring(theme_el, xml_declaration=True, encoding='UTF-8', standalone=True)


def clean_save(prs, path):
    import os, zipfile, shutil
    from lxml import etree as _et
    os.makedirs(os.path.dirname(path) if os.path.dirname(path) else '.', exist_ok=True)
    prs.save(path)
    tmp = path + '.tmp'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_ct = 'http://schemas.openxmlformats.org/package/2006/content-types'
    with zipfile.ZipFile(path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            items = {}
            for item in zin.namelist():
                if 'printerSettings' in item:
                    continue
                items[item] = zin.read(item)
    rid_map = None
    for name, data in items.items():
        if name == 'ppt/_rels/presentation.xml.rels':
            root = _et.fromstring(data)
            to_remove = [r for r in root
                         if 'printerSettings' in r.get('Target', '')]
            if to_remove:
                for r in to_remove:
                    root.remove(r)
                rid_map = {}
                for r in root:
                    old = r.get('Id', '')
                    if old.startswith('rId'):
                        rid_map[old] = 'rId%d' % (len(rid_map) + 1)
                for r in root:
                    old = r.get('Id', '')
                    if old in rid_map:
                        r.set('Id', rid_map[old])
            data = _et.tostring(root, xml_declaration=True,
                                encoding='UTF-8', standalone=True)
            items[name] = data
    if rid_map:
        for name in list(items.keys()):
            if name == 'ppt/presentation.xml':
                root = _et.fromstring(items[name])
                for el in root.iter():
                    rid = el.get(f'{{{ns_r}}}id', '')
                    if rid in rid_map:
                        el.set(f'{{{ns_r}}}id', rid_map[rid])
                items[name] = _et.tostring(root, xml_declaration=True,
                                           encoding='UTF-8', standalone=True)
    ct_data = items.get('[Content_Types].xml')
    if ct_data is not None:
        ct_root = _et.fromstring(ct_data)
        to_remove = []
        for el in ct_root:
            if el.tag == f'{{{ns_ct}}}Default':
                ct_val = el.get('ContentType', '')
                if 'printerSettings' in ct_val:
                    to_remove.append(el)
        for el in to_remove:
            ct_root.remove(el)
        if to_remove:
            items['[Content_Types].xml'] = _et.tostring(
                ct_root, xml_declaration=True, encoding='UTF-8', standalone=True)
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in items.items():
            if name.endswith('.xml') and name.startswith('ppt/slides/slide'):
                try:
                    root = _et.fromstring(data)
                    changed = False
                    for ln in root.iter(f'{{{ns_a}}}ln'):
                        if len(ln) == 0:
                            _et.SubElement(ln, f'{{{ns_a}}}noFill')
                            changed = True
                    if changed:
                        data = _et.tostring(root, xml_declaration=True,
                                            encoding='UTF-8', standalone=True)
                except Exception:
                    pass
            zout.writestr(name, data)
    shutil.move(tmp, path)


def add_slide(prs, layout_index=None):
    if layout_index is not None:
        return prs.slides.add_slide(prs.slide_layouts[layout_index])
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[-1])


def set_theme_colors(prs, C=None):
    """Write C dict colors into the PowerPoint theme clrScheme.

    Maps: primary→accent1, secondary→accent2, tertiary→accent3,
    muted→accent4, light→accent5, text_dark→dk2, text_body→lt2.
    Makes theme colors recognizable by PowerPoint (fills default to C palette).
    """
    C = C or {}
    _ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    theme_part = None
    for rel in prs.part.rels.values():
        if 'theme' in rel.reltype:
            theme_part = rel.target_part
            break
    if theme_part is None:
        return
    theme_el = etree.fromstring(theme_part.blob)
    clrScheme = theme_el.find(f'{{{_ns}}}themeElements/{{{_ns}}}clrScheme')
    if clrScheme is None:
        return

    mapping = {
        'dk1': C.get('text_dark'),
        'lt1': C.get('background'),
        'dk2': C.get('text_dark'),
        'lt2': C.get('card_bg'),
        'accent1': C.get('primary'),
        'accent2': C.get('secondary', C.get('text_body')),
        'accent3': C.get('tertiary', C.get('accent')),
        'accent4': C.get('muted'),
        'accent5': C.get('light'),
        'accent6': C.get('divider'),
    }
    for tag, val in mapping.items():
        if not val:
            continue
        el = clrScheme.find(f'{{{_ns}}}{tag}')
        if el is None:
            continue
        for child in list(el):
            el.remove(child)
        srgb = etree.SubElement(el, f'{{{_ns}}}srgbClr')
        srgb.set('val', val.lstrip('#'))

    theme_part._blob = etree.tostring(theme_el, xml_declaration=True,
                                      encoding='UTF-8', standalone=True)


# ---------------------------------------------------------------------------
# Background image
# ---------------------------------------------------------------------------

def set_slide_bg_image(slide, image_path):
    """Set slide background to an image via OOXML.

    Args:
        slide: Slide object
        image_path: Path to background image

    Returns:
        True if successful, False otherwise
    """
    import os
    from pptx.oxml.ns import qn

    if not os.path.isfile(image_path):
        return False

    # Add image to presentation and get relationship ID
    slide_part = slide.part
    image_part = slide_part.partPackage.get_or_add_image_part(image_path)
    rId = slide_part.relate_to(image_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')

    # Build background XML
    bg = slide.background._element
    # Clear existing background
    for child in list(bg):
        bg.remove(child)

    # Create bgPr with blipFill
    bgPr = etree.SubElement(bg, qn('p:bgPr'))
    blipFill = etree.SubElement(bgPr, qn('a:blipFill'))
    blipFill.set('dpi', '96')
    blipFill.set('rotWithShape', '1')

    blip = etree.SubElement(blipFill, qn('a:blip'))
    blip.set(qn('r:embed'), rId)

    # Add stretch
    stretch = etree.SubElement(blipFill, qn('a:stretch'))
    etree.SubElement(stretch, qn('a:fillRect'))

    # Add effectLst (required by OOXML)
    etree.SubElement(bgPr, qn('a:effectLst'))

    return True


# ---------------------------------------------------------------------------
# Page decorations
# ---------------------------------------------------------------------------

def page_header(slide, title, subtitle='', C=None, left=0.65, width=None,
                typo=None, spacing=None):
    from pptx_designer.tools.text import text
    C = C or {}
    cw = width or (13.333 - 2 * left)
    t = typo or TYPOGRAPHY.get('mckinsey')
    sp = spacing or SPACING.get('mckinsey')
    title_y = sp.page_margin * 0.7
    text(slide, left, title_y, cw, 0.5, title,
         font_size=t.h1, color=C.get('text_dark', '#000000'), bold=True,
         font_name=C.get('font_heading'), C=C)
    if subtitle:
        subtitle_y = title_y + 0.5
        text(slide, left, subtitle_y, cw, 0.25, subtitle,
             font_size=t.caption, color=C.get('text_muted', '#666666'),
             font_name=C.get('font_body'), C=C)
        divider_y = subtitle_y + 0.3
    else:
        divider_y = title_y + 0.5
    divider_h = 0.015
    rect(slide, left, divider_y, cw, divider_h, C.get('divider', '#CCCCCC'))


def top_bar(slide, color, width=13.333, height=0.08, C=None):
    color_val = _resolve_color(color, C)
    return rect(slide, 0, 0, width, height, color_val)


def page_number(slide, current, total, style='simple', C=None, typo=None):
    """Add page number decoration to slide.

    Args:
        slide: Slide object
        current: Current page number
        total: Total number of pages
        style: 'simple' (just number), 'gold' (gold decoration),
               'progress' (with progress bar)
        C: Color dictionary
        typo: Typography settings

    Returns:
        List of shape objects
    """
    from pptx_designer.tools.text import text
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    shapes = []

    if style == 'simple':
        page_text = f"{current} / {total}"
        s = text(slide, 11.5, 7.0, 1.5, 0.3, page_text,
                font_size=t.caption, color=C.get('text_muted', '#666666'),
                align='right', C=C)
        shapes.append(s)

    elif style == 'gold':
        num_str = str(current).zfill(2)

        num_box = text(slide, 10.5, 6.0, 2.0, 1.0, num_str,
                      font_size=48, color=C.get('accent', '#C99A4E'),
                      bold=True, align='right', C=C)
        shapes.append(num_box)

        line = rect(slide, 10.5, 7.0, 2.0, 0.02, C.get('accent', '#C99A4E'))
        shapes.append(line)

        total_box = text(slide, 10.5, 7.1, 2.0, 0.3, f"/ {total}",
                        font_size=t.caption, color=C.get('text_muted', '#666666'),
                        align='right', C=C)
        shapes.append(total_box)

    elif style == 'progress':
        num_str = str(current).zfill(2)

        num_box = text(slide, 0.5, 7.0, 1.5, 0.3, f"{num_str} / {total}",
                      font_size=t.caption, color=C.get('text_dark', '#1A1A1A'),
                      bold=True, C=C)
        shapes.append(num_box)

        bar_bg = rect(slide, 2.5, 7.1, 8.0, 0.08, C.get('divider', '#E0E0E0'))
        shapes.append(bar_bg)

        progress = current / total if total > 0 else 0
        bar_fill = rect(slide, 2.5, 7.1, 8.0 * progress, 0.08,
                       C.get('accent', '#2E6504'))
        shapes.append(bar_fill)

    return shapes


def copy_decorations(slide, template_slide, skip_long_text=True, skip_image=True):
    for shape in template_slide.shapes:
        if skip_image and shape.shape_type == 13:
            continue
        if skip_long_text and shape.has_text_frame:
            if len(shape.text_frame.text) > 50:
                continue
        el = copy.deepcopy(shape._element)
        slide.shapes._spTree.append(el)


def copy_logo(slide, template_slide, color_hints=None):
    for shape in template_slide.shapes:
        if shape.shape_type != 6:
            continue
        if color_hints:
            sp = shape._element
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            for hint in color_hints:
                if sp.find(f'.//{{{ns}}}srgbClr[@val="{hint.lstrip("#")}"]') is not None:
                    el = copy.deepcopy(sp)
                    slide.shapes._spTree.append(el)
                    return
        else:
            el = copy.deepcopy(shape._element)
            slide.shapes._spTree.append(el)
            return
