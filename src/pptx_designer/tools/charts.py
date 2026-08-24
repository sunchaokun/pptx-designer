"""charts — Chart atoms for build-mode PPT generation.

Provides:
  - bar_chart()          Horizontal bar charts
  - comparison_bars()    Side-by-side old vs new metrics
  - donut_chart()        Pie/donut charts (native or shape fallback)
  - native_chart()       Full PowerPoint native chart (bar/line/pie/...)

Usage:
    from pptx_designer.tools.charts import bar_chart, native_chart
    bar_chart(slide, 1.0, 1.5, [("Revenue", 0.75, "$12M"), ...], C=C)
"""
from __future__ import annotations

from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Import helpers from shapes module to avoid duplication
from pptx_designer.tools.shapes import (
    _resolve_color,
    _rgb,
    _set_cjk_font,
    _strip_style,
    _add_shape,
    _set_run,
    _lighten,
    Typography,
    Spacing,
    TYPOGRAPHY,
    SPACING,
    rrect,
    oval,
)


def _lighten(hex_color, amount=30):
    h = hex_color.lstrip('#')
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, r + amount)
    g = min(255, g + amount)
    b = min(255, b + amount)
    return f'{r:02X}{g:02X}{b:02X}'


# ---------------------------------------------------------------------------
# Chart functions
# ---------------------------------------------------------------------------


# ---------------------------------------------------------------------------
# Shape atoms used by chart functions
# ---------------------------------------------------------------------------


# ---------------------------------------------------------------------------
# Chart functions
# ---------------------------------------------------------------------------

def bar_chart(slide, left, top, data, max_width=5.0, bar_height=0.3, C=None,
              typo=None, spacing=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    sp = spacing or SPACING.get('mckinsey')
    bar_colors = [C.get('primary', '#1B5E20'), C.get('accent', '#4CAF50'),
                  C.get('muted', '#81C784'), C.get('light', '#C8E6C9')]
    gap = sp.bar_gap

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        for i, (label, pct, val) in enumerate(data):
            y = top + i * (bar_height + gap)

            bg = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(left), Inches(y),
                            Inches(max_width), Inches(bar_height))
            bg.fill.solid()
            bg.fill.fore_color.rgb = _rgb(C.get('bg_tint', '#F5F5F5'))
            bg.line.fill.background()

            bar_w = max_width * pct
            if bar_w > 0:
                bar = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(y),
                                 Inches(bar_w), Inches(bar_height))
                bar.fill.solid()
                bar.fill.fore_color.rgb = _rgb(bar_colors[i % len(bar_colors)])
                bar.line.fill.background()

            lbl_box = gs.add_textbox(Inches(left - 0.9), Inches(y - 0.03),
                                     Inches(0.85), Inches(bar_height))
            p = lbl_box.text_frame.paragraphs[0]
            _set_run(p, label, font_size=t.caption, color=C.get('text_body', '#333333'),
                     font_name=C.get('font_body'), C=C)
            p.alignment = PP_ALIGN.RIGHT

            val_box = gs.add_textbox(Inches(left + max_width + 0.08), Inches(y - 0.03),
                                     Inches(0.6), Inches(bar_height))
            p2 = val_box.text_frame.paragraphs[0]
            _set_run(p2, val, font_size=t.caption, color=C.get('text_dark', '#000000'),
                     bold=True, font_name=C.get('font_body'), C=C)

        return group
    else:
        for i, (label, pct, val) in enumerate(data):
            y = top + i * (bar_height + gap)
            rrect(slide, left, y, max_width, bar_height, C.get('bg_tint', '#F5F5F5'), C=C)
            bar_w = max_width * pct
            if bar_w > 0:
                rrect(slide, left, y, bar_w, bar_height,
                      bar_colors[i % len(bar_colors)], C=C)
            text(slide, left - 0.9, y - 0.03, 0.85, bar_height, label,
                 font_size=t.caption, color=C.get('text_body', '#333333'), align='right',
                 font_name=C.get('font_body'), C=C)
            text(slide, left + max_width + 0.08, y - 0.03, 0.6, bar_height, val,
                 font_size=t.caption, color=C.get('text_dark', '#000000'), bold=True,
                 font_name=C.get('font_body'), C=C)


def comparison_bars(slide, left, top, metrics, max_width=4.0, C=None,
                    typo=None, spacing=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    sp = spacing or SPACING.get('mckinsey')
    row_h = 0.55

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        for idx, (label, v_old, v_new, pct_old, pct_new) in enumerate(metrics):
            y = top + idx * row_h

            lbl_box = gs.add_textbox(Inches(left - 1.1), Inches(y - 0.02),
                                     Inches(1.0), Inches(0.2))
            p = lbl_box.text_frame.paragraphs[0]
            _set_run(p, label, font_size=t.caption, color=C.get('text_body', '#333333'),
                     bold=True, font_name=C.get('font_body'), C=C)
            p.alignment = PP_ALIGN.RIGHT

            bg1 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(left), Inches(y),
                             Inches(max_width), Inches(0.18))
            bg1.fill.solid()
            bg1.fill.fore_color.rgb = _rgb(C.get('bg_tint', '#F5F5F5'))
            bg1.line.fill.background()

            bar_old = max_width * pct_old
            if bar_old > 0:
                b1 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(y),
                                Inches(bar_old), Inches(0.18))
                b1.fill.solid()
                b1.fill.fore_color.rgb = _rgb(C.get('muted', '#81C784'))
                b1.line.fill.background()

            bg2 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(left), Inches(y + 0.22),
                             Inches(max_width), Inches(0.18))
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = _rgb(C.get('bg_tint', '#F5F5F5'))
            bg2.line.fill.background()

            bar_new = max_width * pct_new
            if bar_new > 0:
                b2 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(y + 0.22),
                                Inches(bar_new), Inches(0.18))
                b2.fill.solid()
                b2.fill.fore_color.rgb = _rgb(C.get('primary', '#1B5E20'))
                b2.line.fill.background()

            old_box = gs.add_textbox(Inches(left + max_width + 0.1), Inches(y - 0.03),
                                     Inches(0.8), Inches(0.2))
            p2 = old_box.text_frame.paragraphs[0]
            _set_run(p2, v_old, font_size=t.micro, color=C.get('text_muted', '#666666'),
                     font_name=C.get('font_body'), C=C)

            new_box = gs.add_textbox(Inches(left + max_width + 0.1), Inches(y + 0.19),
                                     Inches(0.8), Inches(0.2))
            p3 = new_box.text_frame.paragraphs[0]
            _set_run(p3, v_new, font_size=t.micro, color=C.get('text_dark', '#000000'),
                     bold=True, font_name=C.get('font_body'), C=C)

        return group
    else:
        for label, v_old, v_new, pct_old, pct_new in metrics:
            text(slide, left - 1.1, top - 0.02, 1.0, 0.2, label,
                 font_size=t.caption, color=C.get('text_body', '#333333'), bold=True,
                 align='right', font_name=C.get('font_body'), C=C)
            rrect(slide, left, top, max_width, 0.18, C.get('bg_tint', '#F5F5F5'), C=C)
            bar_old = max_width * pct_old
            if bar_old > 0:
                rrect(slide, left, top, bar_old, 0.18, C.get('muted', '#81C784'), C=C)
            rrect(slide, left, top + 0.22, max_width, 0.18, C.get('bg_tint', '#F5F5F5'), C=C)
            bar_new = max_width * pct_new
            if bar_new > 0:
                rrect(slide, left, top + 0.22, bar_new, 0.18, C.get('primary', '#1B5E20'), C=C)
            text(slide, left + max_width + 0.1, top - 0.03, 0.8, 0.2, v_old,
                 font_size=t.micro, color=C.get('text_muted', '#666666'),
                 font_name=C.get('font_body'), C=C)
            text(slide, left + max_width + 0.1, top + 0.19, 0.8, 0.2, v_new,
                 font_size=t.micro, color=C.get('text_dark', '#000000'), bold=True,
                 font_name=C.get('font_body'), C=C)
            top += row_h
        return top


def donut_chart(slide, cx, cy, radius, inner_radius, sectors, C=None,
                typo=None, grouped=True, native=True):
    """Donut/pie chart. When native=True and sectors>1, uses PowerPoint native
    doughnut chart for accurate sector angles. When native=False or sectors==1,
    falls back to Shape-based rendering for maximum visual customization."""
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    if native and len(sectors) > 1:
        chart_w = radius * 2 + 2.5
        chart_h = radius * 2 + 0.6
        chart_left = cx - radius - 0.2
        chart_top = cy - radius - 0.3

        categories = [s[0] for s in sectors]
        pct_values = []
        for s in sectors:
            pct_str = s[1].replace('%', '').strip()
            try:
                pct_values.append(float(pct_str))
            except (ValueError, AttributeError):
                pct_values.append(0)
        sector_colors = [s[2] for s in sectors]

        series = [{"name": "Share", "values": pct_values}]
        chart_style = {
            "show_legend": True,
            "legend_position": "right",
            "show_labels": True,
            "show_percentage": True,
            "show_value": False,
            "label_position": "best_fit",
            "color_scheme": sector_colors,
        }

        result = native_chart(slide, chart_left, chart_top, chart_w, chart_h,
                              "doughnut", categories=categories, series=series,
                              style=chart_style, C=C)
        if result is not None:
            return result

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        for name, pct_str, clr in sectors:
            outer = _add_shape(gs, MSO_SHAPE.OVAL,
                               Inches(cx - radius), Inches(cy - radius),
                               Inches(radius * 2), Inches(radius * 2))
            outer.fill.solid()
            outer.fill.fore_color.rgb = _rgb(clr)
            outer.line.fill.background()

        inner = _add_shape(gs, MSO_SHAPE.OVAL,
                           Inches(cx - inner_radius), Inches(cy - inner_radius),
                           Inches(inner_radius * 2), Inches(inner_radius * 2))
        inner.fill.solid()
        inner.fill.fore_color.rgb = _rgb(C.get('background', '#FFFFFF'))
        inner.line.fill.background()

        center_box = gs.add_textbox(
            Inches(cx - 0.5), Inches(cy - 0.2),
            Inches(1.0), Inches(0.4))
        p = center_box.text_frame.paragraphs[0]
        _set_run(p, '100%', font_size=t.h2, color=C.get('primary', '#1B5E20'),
                 bold=True, font_name=C.get('font_heading'), C=C)
        p.alignment = PP_ALIGN.CENTER

        ly = cy - radius
        lx = cx + radius + 0.5
        for name, pct_str, clr in sectors:
            dot = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(lx), Inches(ly),
                             Inches(0.2), Inches(0.2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = _rgb(clr)
            dot.line.fill.background()

            lbl = gs.add_textbox(Inches(lx + 0.3), Inches(ly - 0.02),
                                 Inches(1.5), Inches(0.25))
            p2 = lbl.text_frame.paragraphs[0]
            _set_run(p2, f'{name}  {pct_str}', font_size=t.caption,
                     color=C.get('text_body', '#333333'),
                     font_name=C.get('font_body'), C=C)
            ly += 0.35

        return group
    else:
        for name, pct_str, clr in sectors:
            oval(slide, cx - radius, cy - radius, radius * 2, radius * 2, clr, C=C)
        oval(slide, cx - inner_radius, cy - inner_radius,
             inner_radius * 2, inner_radius * 2,
             C.get('background', '#FFFFFF'), C=C)
        text(slide, cx - 0.5, cy - 0.2, 1.0, 0.4, '100%',
             font_size=t.h2, color=C.get('primary', '#1B5E20'), bold=True,
             align='center', font_name=C.get('font_heading'), C=C)
        ly = cy - radius
        lx = cx + radius + 0.5
        for name, pct_str, clr in sectors:
            rrect(slide, lx, ly, 0.2, 0.2, clr, C=C)
            text(slide, lx + 0.3, ly - 0.02, 1.5, 0.25, f'{name}  {pct_str}',
                 font_size=t.caption, color=C.get('text_body', '#333333'),
                 font_name=C.get('font_body'), C=C)
            ly += 0.35


def native_chart(slide, left, top, width, height, chart_type,
                 categories=None, series=None, style=None, C=None):
    """Native PowerPoint chart — editable data, axes, gridlines, legend.

    chart_type: 'bar'|'bar_stacked'|'bar_100'|'bar_3d'|
                'bar_horizontal'|'bar_horizontal_stacked'|'bar_horizontal_100'|
                'line'|'line_markers'|'line_stacked'|'line_stacked_100'|
                'pie'|'pie_3d'|'pie_exploded'|
                'doughnut'|'doughnut_exploded'|
                'area'|'area_stacked'|'area_stacked_100'|
                'scatter'|'scatter_lines'|'scatter_smooth'|
                'radar'|'radar_markers'|'bubble'
    categories: ['Q1','Q2','Q3','Q4']  (not used for scatter/bubble)
    series: [{'name':'Revenue','values':[30,45,60,75]}, ...]
            For scatter: values = [[x1,y1],[x2,y2],...]
            For bubble:  values = [[x1,y1,size1],[x2,y2,size2],...]
    style: {
        'show_legend': True,
        'legend_position': 'bottom',  # bottom|top|left|right
        'show_labels': False,
        'show_value': True,
        'show_percentage': False,
        'show_category_name': False,
        'label_font_size': 9,
        'label_position': 'outside_end',  # center|inside_end|outside_end|best_fit
        'number_format': '#,##0',
        'color_scheme': 'brand',  # 'brand'|'auto'|['#hex',...]
        'title': 'Chart Title',
        'value_axis_title': 'Revenue ($M)',
        'category_axis_title': 'Quarter',
        'gridlines': 'major_y',  # 'none'|'major_y'|'major_x'|'major_xy'
        'tick_number_format': '#,##0',
        'chart_style': 2,  # 1-48 built-in PowerPoint chart style
    }
    C: color dictionary (used for 'brand' color_scheme)
    """
    from pptx_designer.renderer.chart_builder import ChartBuilder

    if categories is None:
        categories = ["Q1", "Q2", "Q3", "Q4"]
    if series is None:
        series = [{"name": "Data", "values": [10, 25, 45, 80]}]
    if style is None:
        style = {}

    chart_config = {
        "type": chart_type,
        "categories": categories,
        "series": series,
        "style": style,
    }

    brand_colors = None
    if C:
        brand_colors = {
            "primary": C.get("primary", "#2563EB"),
            "secondary": C.get("secondary", "#64748B"),
            "accent": C.get("accent", "#F97316"),
        }

    position = {"x": left, "y": top, "width": width, "height": height}
    builder = ChartBuilder()
    return builder.build(slide, chart_config, position=position, brand_colors=brand_colors)
