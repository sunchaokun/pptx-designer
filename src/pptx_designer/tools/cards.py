"""Card components — KPI cards, highlight cards, code blocks, etc."""

from __future__ import annotations

from typing import Any

from pptx_designer.tools.shapes import (
    _resolve_color,
    _rgb,
    _set_cjk_font,
    _set_run,
    rect,
    rrect,
    Typography,
    Spacing,
    TYPOGRAPHY,
    SPACING,
)


def kpi_card(slide, left, top, width, height, number, label, trend='',
             trend_up=True, C=None, typo=None, spacing=None, grouped=True):
    """KPI card: bg + accent bar + number + label + trend.

    Args:
        slide: Slide object
        left, top, width, height: Position and size in inches
        number: Main number (e.g., "$12.8M")
        label: Label text (e.g., "Revenue")
        trend: Trend text (e.g., "+23%")
        trend_up: Whether trend is positive
        C: Color dictionary
        typo: Typography settings
        grouped: Whether to group shapes

    Returns:
        List of shapes or GroupShape
    """
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    # Background
    bg = rrect(slide, left, top, width, height, C.get('card', '#FFFFFF'),
               line=C.get('border', '#E5E5E5'))

    # Accent bar at top
    accent_color = C.get('accent', '#1D78FA')
    bar = rect(slide, left, top, width, 0.06, accent_color)

    # Number
    num_box = _add_shape(slide, left + 0.15, top + 0.2, width - 0.3, 0.6,
                         number, font_size=t.h1 if t else 28,
                         color=C.get('text_dark', '#000000'), bold=True, C=C)

    # Label
    label_box = _add_shape(slide, left + 0.15, top + 0.8, width - 0.3, 0.3,
                           label, font_size=t.caption if t else 10,
                           color=C.get('text_muted', '#666666'), C=C)

    # Trend
    shapes = [bg, bar, num_box, label_box]
    if trend:
        trend_color = '#22C55E' if trend_up else '#EF4444'
        trend_box = _add_shape(slide, left + 0.15, top + 1.1, width - 0.3, 0.25,
                               trend, font_size=t.caption if t else 10,
                               color=trend_color, bold=True, C=C)
        shapes.append(trend_box)

    return shapes


def highlight_cards(slide, left, top, cards, total_width=12.0,
                    C=None, typo=None, spacing=None, grouped=True):
    """Row of accent-topped cards.

    Args:
        slide: Slide object
        left, top: Starting position
        cards: List of (title, description, accent_color) tuples
        total_width: Total width for all cards
        C: Color dictionary
        typo: Typography settings

    Returns:
        List of shapes
    """
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    n = len(cards)
    if n == 0:
        return []
    card_width = (total_width - 0.3 * (n - 1)) / n
    shapes = []
    for i, (title, desc, accent) in enumerate(cards):
        x = left + i * (card_width + 0.3)
        # Card background
        bg = rrect(slide, x, top, card_width, 1.8, C.get('card', '#FFFFFF'),
                   line=C.get('border', '#E5E5E5'))
        shapes.append(bg)
        # Accent bar
        bar = rect(slide, x, top, card_width, 0.06, accent)
        shapes.append(bar)
        # Title
        title_box = _add_shape(slide, x + 0.15, top + 0.2, card_width - 0.3, 0.4,
                               title, font_size=t.h3 if t else 16,
                               color=C.get('text_dark', '#000000'), bold=True, C=C)
        shapes.append(title_box)
        # Description
        desc_box = _add_shape(slide, x + 0.15, top + 0.7, card_width - 0.3, 0.8,
                              desc, font_size=t.body if t else 12,
                              color=C.get('text_body', '#4A4A4A'), C=C)
        shapes.append(desc_box)
    return shapes


def code_block(slide, left, top, width, height, lines, language='python',
               C=None, typo=None, grouped=True):
    """Dark-bg code block with language badge + monospace text.

    Args:
        slide: Slide object
        left, top, width, height: Position and size
        lines: List of code lines
        language: Programming language name
        C: Color dictionary
        typo: Typography settings

    Returns:
        List of shapes
    """
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    # Dark background
    bg = rrect(slide, left, top, width, height, '#1E1E1E')

    # Language badge
    badge = rect(slide, left + 0.1, top + 0.1, 1.2, 0.3, '#3B82F6')
    badge_text = _add_shape(slide, left + 0.15, top + 0.12, 1.1, 0.25,
                            language, font_size=10, color='#FFFFFF', bold=True, C=C)

    # Code lines
    shapes = [bg, badge, badge_text]
    y_offset = top + 0.5
    for line in lines:
        line_box = _add_shape(slide, left + 0.2, y_offset, width - 0.4, 0.25,
                              line, font_size=11, color='#D4D4D4',
                              font_name='Consolas', C=C)
        shapes.append(line_box)
        y_offset += 0.28

    return shapes


def section_divider(slide, number, title, C=None, typo=None, grouped=True):
    """Full-slide section divider with large number + title.

    Args:
        slide: Slide object
        number: Section number
        title: Section title
        C: Color dictionary
        typo: Typography settings

    Returns:
        List of shapes
    """
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    # Large number
    num_box = _add_shape(slide, 0.5, 2.0, 3.0, 2.0,
                         str(number).zfill(2), font_size=96,
                         color=C.get('accent', '#1D78FA'), bold=True, C=C)

    # Title
    title_box = _add_shape(slide, 0.5, 4.2, 10.0, 1.0,
                           title, font_size=t.h1 if t else 28,
                           color=C.get('text_dark', '#000000'), bold=True, C=C)

    # Divider line
    line = rect(slide, 0.5, 5.5, 2.0, 0.04, C.get('accent', '#1D78FA'))

    return [num_box, title_box, line]


def hero_slide(slide, title, subtitle='', C=None, typo=None, grouped=True):
    """Full-slide hero with primary background + white title.

    Args:
        slide: Slide object
        title: Hero title
        subtitle: Hero subtitle
        C: Color dictionary
        typo: Typography settings

    Returns:
        List of shapes
    """
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    # Background
    bg = rect(slide, 0, 0, 13.333, 7.5, C.get('primary', '#1D78FA'))

    # Title
    title_box = _add_shape(slide, 0.5, 2.5, 12.333, 1.5,
                           title, font_size=44, color='#FFFFFF', bold=True, C=C)

    shapes = [bg, title_box]
    if subtitle:
        sub_box = _add_shape(slide, 0.5, 4.2, 12.333, 0.8,
                             subtitle, font_size=18, color='#FFFFFF', C=C)
        shapes.append(sub_box)

    return shapes


def cta_slide(slide, title, subtitle='', C=None, typo=None, grouped=True):
    """Full-slide call-to-action.

    Args:
        slide: Slide object
        title: CTA title
        subtitle: CTA subtitle
        C: Color dictionary
        typo: Typography settings

    Returns:
        List of shapes
    """
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    # Background
    bg = rect(slide, 0, 0, 13.333, 7.5, C.get('background', '#FFFFFF'))

    # Title
    title_box = _add_shape(slide, 0.5, 2.5, 12.333, 1.5,
                           title, font_size=44,
                           color=C.get('text_dark', '#000000'), bold=True, C=C)

    shapes = [bg, title_box]
    if subtitle:
        sub_box = _add_shape(slide, 0.5, 4.2, 12.333, 0.8,
                             subtitle, font_size=18,
                             color=C.get('text_body', '#4A4A4A'), C=C)
        shapes.append(sub_box)

    return shapes


def _add_shape(slide, left, top, width, height, txt, font_size=12,
               color='#000000', bold=False, font_name=None, C=None):
    """Add a text shape."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    color_val = _resolve_color(color, C)
    shape = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(color_val.lstrip('#'))
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    return shape
