"""Text tools — text creation, styling, and effect functions for build scripts.

Provides text box creation with alignment, color, CJK support, gradients,
outlines, shadows, glows, and dramatic typography.

Usage:
    from pptx_designer.tools.text import *
    C = {'primary': '#2E6504', 'accent': '#7DA92F', ...}
    prs = Presentation(template_path)
    s = add_slide(prs)
    text(s, 1.0, 1.0, 10.0, 0.5, 'Hello World', font_size=24, color='text_dark', C=C)
    prs.save('output.pptx')
"""

from __future__ import annotations

from contextlib import suppress

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx_designer.effects.text_effects import (
    TEXT_GRADIENT_PRESETS,
    set_vertical_text,
)
from pptx_designer.renderer.theme_context import resolve_color_context
from pptx_designer.tools.shapes import (
    _resolve_color,
    _rgb,
    _set_cjk_font,
    _set_run,
)

# ── Text box creation ──────────────────────────────────────────────


def text(
    slide,
    left,
    top,
    width,
    height,
    txt,
    font_size=12,
    color="text_body",
    bold=False,
    align="left",
    font_name=None,
    C=None,
    anchor="top",
):
    C = resolve_color_context(slide, C)
    font_name = font_name or C.get("font_body")
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor == "middle":
        with suppress(Exception):
            txBox.text_frame._txBody.bodyPr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    _set_run(p, txt, font_size=font_size, color=color, bold=bold, font_name=font_name, C=C)
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    return txBox


def multiline(
    slide,
    left,
    top,
    width,
    height,
    lines,
    font_size=12,
    color="text_body",
    bold=False,
    align="left",
    font_name=None,
    C=None,
    line_spacing=None,
):
    C = resolve_color_context(slide, C)
    font_name = font_name or C.get("font_body")
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_run(p, line, font_size=font_size, color=color, bold=bold, font_name=font_name, C=C)
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        if line_spacing:
            p.space_before = Pt(line_spacing)
            p.space_after = Pt(line_spacing)
        else:
            p.space_before = Pt(2)
            p.space_after = Pt(2)
    return txBox


def dramatic_text(
    slide,
    left,
    top,
    width,
    height,
    big_text,
    small_text,
    big_size=80,
    small_size=10,
    big_color="text_dark",
    small_color="text_muted",
    big_bold=True,
    small_bold=False,
    align="left",
    font_name=None,
    C=None,
    vertical=False,
):
    C = C or {}
    shapes = []

    if vertical:
        big_h = height * 0.7
        small_h = height * 0.3

        big_shape = text(
            slide,
            left,
            top,
            width,
            big_h,
            big_text,
            font_size=big_size,
            color=big_color,
            bold=big_bold,
            align=align,
            font_name=font_name,
            C=C,
        )
        shapes.append(big_shape)

        small_shape = text(
            slide,
            left,
            top + big_h,
            width,
            small_h,
            small_text,
            font_size=small_size,
            color=small_color,
            bold=small_bold,
            align=align,
            font_name=font_name,
            C=C,
        )
        shapes.append(small_shape)
    else:
        big_w = width * 0.7
        small_w = width * 0.3

        big_shape = text(
            slide,
            left,
            top,
            big_w,
            height,
            big_text,
            font_size=big_size,
            color=big_color,
            bold=big_bold,
            align=align,
            font_name=font_name,
            C=C,
        )
        shapes.append(big_shape)

        small_shape = text(
            slide,
            left + big_w,
            top,
            small_w,
            height,
            small_text,
            font_size=small_size,
            color=small_color,
            bold=small_bold,
            align=align,
            font_name=font_name,
            C=C,
        )
        shapes.append(small_shape)

    return shapes


# ── Vertical text ──────────────────────────────────────────────────


def vertical_text(
    slide,
    left,
    top,
    width,
    height,
    txt,
    direction="ea",
    font_name="STKaiti",
    font_size=24,
    color="#000000",
    bold=False,
    align="center",
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_vertical_text(tf, direction)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(_resolve_color(color, None))
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    return txBox


# ── Gradient text ──────────────────────────────────────────────────


def gradient_text(
    slide,
    left,
    top,
    width,
    height,
    txt,
    preset="gold-shine",
    stops=None,
    font_size=44,
    bold=False,
    font_name=None,
    cjk_font=None,
    align="left",
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    gradient_stops = stops or TEXT_GRADIENT_PRESETS.get(preset)
    if not gradient_stops:
        raise KeyError(f"Unknown text gradient preset: {preset!r}")

    # Per-run solid colors are deliberately used instead of OOXML gradFill.
    # PowerPoint and LibreOffice disagree on text gradients, while solid runs
    # remain editable and preserve contrast through the export pipeline.
    normalized = [(color.lstrip("#"), pos / (100000 if pos > 100 else 100)) for color, pos in gradient_stops]
    for char_index, char in enumerate(txt):
        run = p.add_run()
        run.text = char
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if font_name:
            run.font.name = font_name
        if cjk_font:
            _set_cjk_font(run, cjk_font)
        pos = char_index / max(1, len(txt) - 1)
        for stop_index in range(1, len(normalized)):
            if pos <= normalized[stop_index][1]:
                c1, p1 = normalized[stop_index - 1]
                c2, p2 = normalized[stop_index]
                ratio = (pos - p1) / max(0.0001, p2 - p1)
                rgb = tuple(round(int(c1[i : i + 2], 16) * (1 - ratio) + int(c2[i : i + 2], 16) * ratio) for i in (0, 2, 4))
                run.font.color.rgb = RGBColor(*rgb)
                break
        else:
            run.font.color.rgb = RGBColor.from_string(normalized[-1][0])
    return txBox


# ── Text outline (dual-layer overlay) ──────────────────────────────


def text_outline(
    slide,
    left,
    top,
    width,
    height,
    txt,
    color="#FFFFFF",
    width_pt=1.5,
    font_size=44,
    bold=False,
    font_name=None,
    C=None,
    align="left",
    fill_color=None,
):
    offset = width_pt * 0.015  # tiny offset for outline effect

    # Bottom layer (outline color)
    txBox_bg = slide.shapes.add_textbox(
        Inches(left - offset), Inches(top - offset), Inches(width + offset * 2), Inches(height + offset * 2)
    )
    tf_bg = txBox_bg.text_frame
    tf_bg.word_wrap = True
    p_bg = tf_bg.paragraphs[0]
    p_bg.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run_bg = p_bg.add_run()
    run_bg.text = txt
    run_bg.font.size = Pt(font_size)
    run_bg.font.bold = bold
    if font_name:
        run_bg.font.name = font_name
    run_bg.font.color.rgb = _rgb(_resolve_color(color, C)) if color else RGBColor(0xFF, 0xFF, 0xFF)
    cjk_font = (C or {}).get("font_cjk") or (C or {}).get("font_body")
    if cjk_font:
        _set_cjk_font(run_bg, cjk_font)

    # Top layer (fill color)
    txBox_fg = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf_fg = txBox_fg.text_frame
    tf_fg.word_wrap = True
    p_fg = tf_fg.paragraphs[0]
    p_fg.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run_fg = p_fg.add_run()
    run_fg.text = txt
    run_fg.font.size = Pt(font_size)
    run_fg.font.bold = bold
    if font_name:
        run_fg.font.name = font_name
    if fill_color:
        run_fg.font.color.rgb = _rgb(_resolve_color(fill_color, C))
    else:
        run_fg.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if cjk_font:
        _set_cjk_font(run_fg, cjk_font)

    return txBox_fg


# ── Text shadow ────────────────────────────────────────────────────


def text_shadow(
    slide,
    left,
    top,
    width,
    height,
    txt,
    blur_pt=8,
    distance_pt=3,
    direction_deg=90,
    color="#000000",
    alpha_pct=25,
    font_size=44,
    bold=False,
    font_name=None,
    C=None,
    align="left",
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(_resolve_color(color, C))
    if font_name:
        run.font.name = font_name
    cjk_font = (C or {}).get("font_cjk") or (C or {}).get("font_body")
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    from pptx_designer.renderer.visual_effects import apply_shadow as _apply_shadow

    _apply_shadow(
        txBox, blur_pt=blur_pt, distance_pt=distance_pt, direction_deg=direction_deg, color=color, alpha_pct=alpha_pct
    )
    return txBox


# ── Text glow ──────────────────────────────────────────────────────


def text_glow(
    slide,
    left,
    top,
    width,
    height,
    txt,
    color="#00FFFF",
    size_pt=8,
    alpha_pct=40,
    font_size=44,
    bold=False,
    font_name=None,
    C=None,
    align="left",
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(_resolve_color(color, C))
    if font_name:
        run.font.name = font_name
    cjk_font = (C or {}).get("font_cjk") or (C or {}).get("font_body")
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    from pptx_designer.renderer.visual_effects import apply_glow as _apply_glow

    _apply_glow(txBox, radius_pt=size_pt, color=color, alpha_pct=alpha_pct)
    return txBox
