"""Image tools — image placement, cropping, AI fetch, and shape-filling.

Provides high-level image functions for build scripts: cover-fit cropping,
AI image generation/fetching, shape-masked images (circle, hex, star, etc.),
soft-edge, duotone, artistic effects, gradient masks, and slide backgrounds.

Usage:
    from pptx_designer.tools.images import *
    cover_image(slide, 0, 0, 10, 7.5, 'photo.jpg')
    circle_image(slide, 5, 3.75, 1.5, 'avatar.png')
    ai_image(slide, 0, 0, 10, 7.5, 'futuristic city')
"""

from __future__ import annotations

import hashlib
import os
import tempfile

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from pptx_designer.effects.image_effects import (
    add_circle_image as _add_circle_image,
)
from pptx_designer.effects.image_effects import (
    add_image_in_shape as _add_image_in_shape,
)
from pptx_designer.effects.image_effects import (
    apply_blip_artistic,
    apply_blip_duotone,
)
from pptx_designer.effects.shape_effects import (
    GradientFill,
    GradientStop,
    apply_soft_edge,
)
from pptx_designer.tools.shapes import _add_shape


def cover_image(slide, left, top, width, height, image_path):
    """Add image with cover-fit (crop to fill, no stretch).

    Uses Pillow to pre-crop the image to the exact aspect ratio,
    then places it at the specified position. This is the correct
    way to add images to PPT — never use add_picture with stretch.
    """
    if not os.path.isfile(image_path):
        return None
    from PIL import Image as PILImage

    img = PILImage.open(image_path)
    img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        crop_w = int(img_h * box_ratio)
        crop_h = img_h
        cleft = (img_w - crop_w) // 2
        ctop = 0
    else:
        crop_w = img_w
        crop_h = int(img_w / box_ratio)
        cleft = 0
        ctop = (img_h - crop_h) // 2
    cropped = img.crop((cleft, ctop, cleft + crop_w, ctop + crop_h))
    cache_dir = os.path.join(tempfile.gettempdir(), "ppt-cropped")
    os.makedirs(cache_dir, exist_ok=True)
    crop_key = f"crop:{image_path}:{width}x{height}"
    crop_hash = hashlib.md5(crop_key.encode()).hexdigest()
    cropped_path = os.path.join(cache_dir, f"{crop_hash}.png")
    if not os.path.exists(cropped_path):
        cropped.save(cropped_path, "PNG")
    return slide.shapes.add_picture(
        cropped_path,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )


def gradient_mask_image(
    slide, left, top, width, height, bg_color="#FFFFFF", direction="bottom", alpha_start=100, alpha_end=0
):
    """Add a gradient mask overlay for seamless image-background fusion.

    Creates a rectangle with gradient fill that fades from opaque to transparent,
    useful for blending images into backgrounds (e.g., cityscape fading into sky).

    Args:
        slide: Slide object
        left, top, width, height: Position and size in inches
        bg_color: Background color for the gradient (should match slide background)
        direction: 'bottom' (fade from bottom), 'top' (fade from top),
                   'left' (fade from left), 'right' (fade from right)
        alpha_start: Starting opacity (0-100, where 100 = fully opaque)
        alpha_end: Ending opacity (0-100, where 0 = fully transparent)

    Returns:
        Shape object (the gradient mask rectangle)
    """
    shape = _add_shape(slide.shapes, MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))

    angle_map = {
        "bottom": 5400000,
        "top": 16200000,
        "left": 10800000,
        "right": 0,
    }
    angle = angle_map.get(direction, 5400000)

    gf = GradientFill(
        angle=angle,
        stops=[
            GradientStop(color=bg_color, position=0, alpha=alpha_start * 1000),
            GradientStop(color=bg_color, position=100000, alpha=alpha_end * 1000),
        ],
    )
    gf.apply(shape)

    shape.line.fill.background()

    return shape


def ai_image(
    slide,
    left,
    top,
    width,
    height,
    keywords,
    *,
    mode="auto",
    emotion="",
    goal="",
    llm_provider=None,
    llm_api_key=None,
    llm_base_url=None,
    llm_model=None,
    unsplash_access_key=None,
    pexels_api_key=None,
    image_cache_dir=None,
    auto_detect=True,
    fallback_placeholder=True,
):
    """Generate (or fetch) an image and place it cover-fit in ONE call.

    Wraps `fetch_image()` + `cover_image()` so build.py never needs to call
    image APIs directly. Preferred over hand-rolled urllib/requests scripts —
    handles cache-first, retry, multi-engine fallback, and cover-fit cropping.

    Args:
        slide: target slide.
        left/top/width/height: placement box (inches).
        keywords: image search / generation prompt.
        mode: 'auto' (default) | 'generate' | 'search' | 'enhance'.
        emotion/goal: optional design hints forwarded to image engines.
        llm_provider: e.g. 'seedream' | 'gpt-image' | 'wanx' | 'kimi' (enhance).
        llm_api_key / llm_base_url / llm_model: generation credentials.
        unsplash_access_key / pexels_api_key: search credentials.
        image_cache_dir: override image cache location.
        auto_detect: auto-detect provider/keys from environment.
        fallback_placeholder: draw a neutral placeholder box when fetch fails
            (keeps the build script from crashing mid-layout).

    Returns:
        The added picture shape, or None if fetch failed and no fallback.
    """
    from pptx_designer.tools.images import _fetch_image

    result = _fetch_image(
        keywords,
        mode=mode,
        emotion=emotion,
        goal=goal,
        width=int(width * 96),
        height=int(height * 96),
        llm_provider=llm_provider,
        llm_api_key=llm_api_key,
        llm_base_url=llm_base_url,
        llm_model=llm_model,
        unsplash_access_key=unsplash_access_key,
        pexels_api_key=pexels_api_key,
        image_cache_dir=image_cache_dir,
        auto_detect=auto_detect,
    )
    path = result.get("path") if result else None
    if path and os.path.isfile(path):
        return cover_image(slide, left, top, width, height, path)
    if fallback_placeholder:
        from pptx_designer.tools.shapes import rect
        from pptx_designer.tools.text import text

        rect(slide, left, top, width, height, "#E8ECF1", C=None)
        text(slide, left, top, width, height * 0.3, keywords, font_size=10, color="#9AA5B1", align="center", C=None)
        return None
    return None


def circle_image(slide, cx, cy, radius, image_path, border_color=None):
    return _add_circle_image(slide, cx, cy, radius, image_path, border_hex=border_color)


def hex_image(slide, cx, cy, size, image_path, border_color=None):
    x = cx - size / 2
    y = cy - size * 0.87 / 2
    return _add_image_in_shape(slide, MSO_SHAPE.HEXAGON, x, y, size, size * 0.87, image_path, border_hex=border_color)


def star_image(slide, cx, cy, size, image_path, points=5, border_color=None):
    _STAR_MAP = {
        5: MSO_SHAPE.STAR_5_POINT,
        6: MSO_SHAPE.STAR_6_POINT,
        8: MSO_SHAPE.STAR_8_POINT,
        10: MSO_SHAPE.STAR_10_POINT,
        12: MSO_SHAPE.STAR_12_POINT,
    }
    mso = _STAR_MAP.get(points, MSO_SHAPE.STAR_5_POINT)
    x = cx - size / 2
    y = cy - size / 2
    return _add_image_in_shape(slide, mso, x, y, size, size, image_path, border_hex=border_color)


def diamond_image(slide, cx, cy, size, image_path, border_color=None):
    x = cx - size / 2
    y = cy - size / 2
    return _add_image_in_shape(slide, MSO_SHAPE.DIAMOND, x, y, size, size, image_path, border_hex=border_color)


def heart_image(slide, cx, cy, size, image_path, border_color=None):
    x = cx - size / 2
    y = cy - size / 2
    return _add_image_in_shape(slide, MSO_SHAPE.HEART, x, y, size, size, image_path, border_hex=border_color)


def shape_image(slide, shape_type, left, top, width, height, image_path, border_color=None):
    _type = shape_type
    if isinstance(_type, str):
        _type = getattr(MSO_SHAPE, _type.upper(), MSO_SHAPE.OVAL)
    return _add_image_in_shape(slide, _type, left, top, width, height, image_path, border_hex=border_color)


def soft_edge_image(slide, left, top, width, height, image_path, soft_radius=10):
    if not os.path.isfile(image_path):
        return None
    shape = slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    apply_soft_edge(shape, radius_pt=soft_radius)
    return shape


def duotone_image(slide, left, top, width, height, image_path, color1="#0000FF", color2="#FF0000"):
    shape = _add_image_in_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, image_path)
    apply_blip_duotone(shape, color1, color2)
    return shape


def artistic_image(slide, left, top, width, height, image_path, effect="watercolor_sponge", params=None):
    shape = _add_image_in_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, image_path)
    apply_blip_artistic(shape, effect, params)
    return shape


def set_slide_bg_image(slide, image_path):
    """Set slide background to an image via OOXML.

    Args:
        slide: Slide object
        image_path: Path to background image

    Returns:
        True if successful, False otherwise
    """
    if not os.path.isfile(image_path):
        return False

    slide_part = slide.part
    image_part = slide_part.partPackage.get_or_add_image_part(image_path)
    rId = slide_part.relate_to(image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")

    bg = slide.background._element
    for child in list(bg):
        bg.remove(child)

    bgPr = etree.SubElement(bg, qn("p:bgPr"))
    blipFill = etree.SubElement(bgPr, qn("a:blipFill"))
    blipFill.set("dpi", "96")
    blipFill.set("rotWithShape", "1")

    blip = etree.SubElement(blipFill, qn("a:blip"))
    blip.set(qn("r:embed"), rId)

    stretch = etree.SubElement(blipFill, qn("a:stretch"))
    etree.SubElement(stretch, qn("a:fillRect"))

    etree.SubElement(bgPr, qn("a:effectLst"))

    return True


def _fetch_image(*args, **kwargs):
    """Lazy-import fetch_image — avoids circular imports at module level."""
    from pptx_designer.ai import fetch_image

    return fetch_image(*args, **kwargs)


__all__ = [
    "cover_image",
    "gradient_mask_image",
    "ai_image",
    "circle_image",
    "hex_image",
    "star_image",
    "diamond_image",
    "heart_image",
    "shape_image",
    "soft_edge_image",
    "duotone_image",
    "artistic_image",
    "set_slide_bg_image",
]
