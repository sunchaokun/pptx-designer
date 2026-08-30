"""Shape tools — low-level shape creation functions for build scripts.

Provides a comprehensive set of shape functions (rect, oval, hexagon, etc.)
with solid fills, dashed borders, and boolean operations.

Usage:
    from pptx_designer.tools.shapes import *
    C = {'primary': '#2E6504', 'accent': '#7DA92F', ...}
    prs = Presentation(template_path)
    s = add_slide(prs)
    rect(s, 1.0, 1.0, 3.0, 1.5, '#1D78FA', C=C)
    prs.save('output.pptx')
"""

from __future__ import annotations

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from pptx_designer.renderer.theme_context import resolve_color_context


def _resolve_color(val, C):
    if val is None:
        return "#000000"
    if val.startswith("#"):
        return val
    return (C or {}).get(val, "#000000")


def _rgb(hex_str):
    return RGBColor.from_string(hex_str.lstrip("#"))


def _set_cjk_font(run, font_name):
    if not font_name:
        return
    rPr = run._r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
    if rPr is None:
        return
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ea = rPr.find(f"{{{ns}}}ea")
    if ea is None:
        ea = etree.SubElement(rPr, f"{{{ns}}}ea")
    ea.set("typeface", font_name)
    cs = rPr.find(f"{{{ns}}}cs")
    if cs is None:
        cs = etree.SubElement(rPr, f"{{{ns}}}cs")
    cs.set("typeface", font_name)


def _strip_style(shape):
    sp = shape._element
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    style_el = sp.find(f"{{{ns}}}style")
    if style_el is not None:
        sp.remove(style_el)


def _add_shape(shapes, mso_type, left, top, width, height):
    sh = shapes.add_shape(mso_type, left, top, width, height)
    _strip_style(sh)
    return sh


def _set_run(paragraph, txt, font_size=12, color="text_body", bold=False, font_name=None, C=None):
    run = paragraph.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(_resolve_color(color, C))
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    cjk_font = (C or {}).get("font_cjk") or (C or {}).get("font_body")
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    return run


def _lighten(hex_color, amount=30):
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, r + amount)
    g = min(255, g + amount)
    b = min(255, b + amount)
    return f"{r:02X}{g:02X}{b:02X}"


def _centered_shape(slide, mso_type, cx, cy, width, height, fill, line=None, C=None):
    C = resolve_color_context(slide, C)
    sh = _add_shape(
        slide.shapes, mso_type, Inches(cx - width / 2), Inches(cy - height / 2), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        sh.line.color.rgb = _rgb(_resolve_color(line, C))
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def _bool_import():
    from pptx_designer.renderer.boolean_shapes import HAS_SHAPELY

    return HAS_SHAPELY


def neon_border(slide, left, top, width, height, color="#8B5CF6", radius=0.1):
    from pptx_designer.effects.decoration import add_neon_border

    return add_neon_border(slide, left, top, width, height, color=color, radius=radius)


def add_slide(prs, layout_index=None):
    if layout_index is not None:
        return prs.slides.add_slide(prs.slide_layouts[layout_index])
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[-1])


def rect(slide, left, top, width, height, fill, line=None, C=None):
    C = resolve_color_context(slide, C)
    shape = _add_shape(slide.shapes, MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        shape.line.color.rgb = _rgb(_resolve_color(line, C))
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def dashed_rect(slide, left, top, width, height, fill=None, line_color="#000000", line_width=1, dash="dash", C=None):
    """Create a rectangle with dashed border.

    Args:
        slide: Slide object
        left, top, width, height: Position and size in inches
        fill: Fill color (None for transparent)
        line_color: Border color
        line_width: Border width in points
        dash: Dash style - 'dash', 'lgDash', 'dot', 'dashDot', 'lgDashDot'
        C: Color dictionary for role-based colors

    Returns:
        Shape object
    """
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    C = resolve_color_context(slide, C)

    shape = _add_shape(slide.shapes, MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))

    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    else:
        shape.fill.background()

    shape.line.color.rgb = _rgb(_resolve_color(line_color, C))
    shape.line.width = Pt(line_width)

    dash_map = {
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "lgDash": MSO_LINE_DASH_STYLE.LONG_DASH,
        "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dashDot": MSO_LINE_DASH_STYLE.DASH_DOT,
        "lgDashDot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
    }
    if dash in dash_map:
        shape.line.dash_style = dash_map[dash]

    return shape


def rrect(slide, left, top, width, height, fill, line=None, C=None):
    C = resolve_color_context(slide, C)
    shape = _add_shape(
        slide.shapes, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        shape.line.color.rgb = _rgb(_resolve_color(line, C))
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def oval(slide, left, top, width, height, fill, line=None, C=None):
    C = resolve_color_context(slide, C)
    shape = _add_shape(slide.shapes, MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        shape.line.color.rgb = _rgb(_resolve_color(line, C))
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def shape(slide, shape_type, left, top, width, height, fill, line=None, C=None):
    C = resolve_color_context(slide, C)
    _type = shape_type
    if isinstance(_type, str):
        _type = getattr(MSO_SHAPE, _type.upper(), MSO_SHAPE.RECTANGLE)
    if width < 0 or height < 0 or width == 0 or height == 0:
        width = max(width, 0.1)
        height = max(height, 0.1)
    sh = _add_shape(slide.shapes, _type, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        sh.line.color.rgb = _rgb(_resolve_color(line, C))
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def hexagon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.HEXAGON, cx, cy, size, size * 0.87, fill, line, C)


def pentagon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.REGULAR_PENTAGON, cx, cy, size, size, fill, line, C)


def octagon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.OCTAGON, cx, cy, size, size, fill, line, C)


def diamond(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.DIAMOND, cx, cy, size, size, fill, line, C)


def triangle(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height, fill, line, C)


def right_triangle(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height, fill, line, C)


def parallelogram(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.PARALLELOGRAM, left, top, width, height, fill, line, C)


def trapezoid(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.TRAPEZOID, left, top, width, height, fill, line, C)


def star5(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_5_POINT, cx, cy, size, size, fill, line, C)


def star6(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_6_POINT, cx, cy, size, size, fill, line, C)


def star8(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_8_POINT, cx, cy, size, size, fill, line, C)


def star10(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_10_POINT, cx, cy, size, size, fill, line, C)


def star12(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_12_POINT, cx, cy, size, size, fill, line, C)


def donut(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.DONUT, cx, cy, size, size, fill, line, C)


def heart(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.HEART, cx, cy, size, size, fill, line, C)


def cross(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.CROSS, cx, cy, size, size, fill, line, C)


def arrow(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, height, fill, line, C)


def chevron(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.CHEVRON, left, top, width, height, fill, line, C)


def cloud(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.CLOUD, left, top, width, height, fill, line, C)


def lightning(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.LIGHTNING_BOLT, left, top, width, height, fill, line, C)


def gear(slide, cx, cy, size, fill, line=None, C=None, teeth=6):
    mso = MSO_SHAPE.GEAR_9 if teeth >= 9 else MSO_SHAPE.GEAR_6
    return _centered_shape(slide, mso, cx, cy, size, size, fill, line, C)


def funnel(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FUNNEL, left, top, width, height, fill, line, C)


def moon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.MOON, cx, cy, size, size, fill, line, C)


def sun(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.SUN, cx, cy, size, size, fill, line, C)


def wave(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.WAVE, left, top, width, height, fill, line, C)


def block_arc(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.BLOCK_ARC, cx, cy, size, size, fill, line, C)


def callout(slide, left, top, width, height, fill, line=None, C=None, style="rect"):
    _MAP = {
        "rect": MSO_SHAPE.RECTANGULAR_CALLOUT,
        "round": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        "oval": MSO_SHAPE.OVAL_CALLOUT,
        "cloud": MSO_SHAPE.CLOUD_CALLOUT,
    }
    mso = _MAP.get(style, MSO_SHAPE.RECTANGULAR_CALLOUT)
    return shape(slide, mso, left, top, width, height, fill, line, C)


def flow_process(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FLOWCHART_PROCESS, left, top, width, height, fill, line, C)


def flow_decision(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.FLOWCHART_DECISION, cx, cy, size, size, fill, line, C)


def flow_data(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FLOWCHART_DATA, left, top, width, height, fill, line, C)


def flow_document(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FLOWCHART_DOCUMENT, left, top, width, height, fill, line, C)


def flow_connector(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.FLOWCHART_CONNECTOR, cx, cy, size, size, fill, line, C)


def no_symbol(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.NO_SYMBOL, cx, cy, size, size, fill, line, C)


def plaque(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.PLAQUE, left, top, width, height, fill, line, C)


def frame(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FRAME, left, top, width, height, fill, line, C)


def cube(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.CUBE, left, top, width, height, fill, line, C)


def bevel(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.BEVEL, left, top, width, height, fill, line, C)


def folded_corner(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FOLDED_CORNER, left, top, width, height, fill, line, C)


def tear(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.TEAR, cx, cy, size, size, fill, line, C)


def math_plus(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.MATH_PLUS, cx, cy, size, size, fill, line, C)


def math_multiply(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.MATH_MULTIPLY, cx, cy, size, size, fill, line, C)


def spotlight(slide, cx, cy, radius, alpha=70, color="#000000"):
    # Boolean cutouts become full-page fills in PowerPoint/LibreOffice export.
    # A layered translucent halo is renderer-safe and keeps the effect local.
    group = slide.shapes.add_group_shape()
    gs = group.shapes
    from pptx_designer.effects.shape_effects import set_solid_fill_with_alpha

    for scale, strength in ((1.0, 0.12), (0.78, 0.18), (0.56, 0.26), (0.34, 0.38)):
        r = radius * scale
        ring = _add_shape(gs, MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2))
        set_solid_fill_with_alpha(ring, color, max(1, min(100, round(alpha * strength))))
        ring.line.fill.background()
    return group


def bool_donut(slide, cx, cy, outer_r, inner_r, fill="#1D78FA", line=None, C=None):
    if not _bool_import():
        return donut(slide, cx, cy, outer_r * 2, fill, line, C)
    from pptx_designer.renderer.boolean_shapes import (
        bool_shape,
        bool_subtract,
        poly_circle,
    )

    outer = poly_circle(cx, cy, outer_r)
    inner = poly_circle(cx, cy, inner_r)
    geom = bool_subtract(outer, inner)
    size = outer_r * 2
    return bool_shape(geom, slide, cx - outer_r, cy - outer_r, size, size, fill=fill, line=line, C=C)


def bool_frame(slide, x, y, w, h, border, fill=None, line=None, C=None):
    if not _bool_import():
        rrect(slide, x, y, w, h, fill or "#1D78FA", line, C)
        return None
    from pptx_designer.renderer.boolean_shapes import (
        bool_shape,
        bool_subtract,
        poly_rect,
    )

    outer = poly_rect(x, y, w, h)
    inner = poly_rect(x + border, y + border, w - 2 * border, h - 2 * border)
    geom = bool_subtract(outer, inner)
    return bool_shape(geom, slide, x, y, w, h, fill=fill or "#1D78FA", line=line, C=C)


def bool_clipped_card(slide, x, y, w, h, clip_corners, clip_size=0.3, fill=None, line=None, C=None):
    if not _bool_import():
        return rrect(slide, x, y, w, h, fill or "#1D78FA", line, C)
    from pptx_designer.renderer.boolean_shapes import (
        Polygon,
        bool_shape,
        bool_subtract,
        bool_union,
        poly_rect,
    )

    base = poly_rect(x, y, w, h)
    clips = []
    if "tl" in clip_corners:
        clips.append(Polygon([(x, y), (x + clip_size, y), (x, y + clip_size)]))
    if "tr" in clip_corners:
        clips.append(Polygon([(x + w - clip_size, y), (x + w, y), (x + w, y + clip_size)]))
    if "bl" in clip_corners:
        clips.append(Polygon([(x, y + h - clip_size), (x + clip_size, y + h), (x, y + h)]))
    if "br" in clip_corners:
        clips.append(Polygon([(x + w, y + h - clip_size), (x + w, y + h), (x + w - clip_size, y + h)]))
    if not clips:
        return rrect(slide, x, y, w, h, fill or "#1D78FA", line, C)
    all_clips = bool_union(*clips)
    geom = bool_subtract(base, all_clips)
    return bool_shape(geom, slide, x, y, w, h, fill=fill or "#1D78FA", line=line, C=C)


def bool_neon_tube(slide, x, y, w, h, wall=0.06, fill=None, C=None):
    if not _bool_import():
        return neon_border(slide, x, y, w, h, color=fill or "#8B5CF6")
    from pptx_designer.renderer.boolean_shapes import (
        bool_shape,
        bool_subtract,
        poly_rounded_rect,
    )

    outer = poly_rounded_rect(x, y, w, h, wall * 2)
    inner = poly_rounded_rect(x + wall, y + wall, w - 2 * wall, h - 2 * wall, wall)
    geom = bool_subtract(outer, inner)
    return bool_shape(geom, slide, x, y, w, h, fill=fill or "#8B5CF6", C=C)


def bool_star(slide, cx, cy, r, points=5, inner_ratio=0.4, fill=None, line=None, C=None):
    if not _bool_import():
        mso_map = {5: MSO_SHAPE.STAR_5_POINT, 6: MSO_SHAPE.STAR_6_POINT, 8: MSO_SHAPE.STAR_8_POINT}
        mso = mso_map.get(points, MSO_SHAPE.STAR_5_POINT)
        return _centered_shape(slide, mso, cx, cy, r * 2, r * 2, fill or "#1D78FA", line, C)
    from pptx_designer.renderer.boolean_shapes import bool_shape, poly_star

    geom = poly_star(cx, cy, r, points=points, inner_ratio=inner_ratio)
    return bool_shape(geom, slide, cx - r, cy - r, r * 2, r * 2, fill=fill or "#1D78FA", line=line, C=C)


def bool_cross(slide, cx, cy, w, h, bar_ratio=0.33, fill=None, line=None, C=None):
    if not _bool_import():
        return cross(slide, cx, cy, max(w, h), fill or "#1D78FA", line, C)
    from pptx_designer.renderer.boolean_shapes import bool_shape, bool_union, poly_rect

    bar_w = w * bar_ratio
    bar_h = h * bar_ratio
    v_bar = poly_rect(cx - bar_w / 2, cy - h / 2, bar_w, h)
    h_bar = poly_rect(cx - w / 2, cy - bar_h / 2, w, bar_h)
    geom = bool_union(v_bar, h_bar)
    return bool_shape(geom, slide, cx - w / 2, cy - h / 2, w, h, fill=fill or "#1D78FA", line=line, C=C)


# ── Typography & Spacing presets ──────────────────────────────────────


class Typography:
    """Font size scale for a design style."""

    def __init__(self, hero=44, h1=28, h2=20, h3=16, body=12, caption=10, micro=8):
        self.hero = hero
        self.h1 = h1
        self.h2 = h2
        self.h3 = h3
        self.body = body
        self.caption = caption
        self.micro = micro

    def scale(self, level):
        return getattr(self, level, self.body)


class Spacing:
    """Spacing system (margins, padding, gaps, rhythm)."""

    def __init__(
        self, page_margin=0.65, section_gap=0.5, card_gap=0.35, card_padding=0.2, line_height=1.4, bar_gap=0.2
    ):
        self.page_margin = page_margin
        self.section_gap = section_gap
        self.card_gap = card_gap
        self.card_padding = card_padding
        self.line_height = line_height
        self.bar_gap = bar_gap


TYPOGRAPHY = {
    "mckinsey": Typography(hero=44, h1=28, h2=20, h3=16, body=12, caption=10),
    "cyberpunk": Typography(hero=48, h1=32, h2=22, h3=16, body=11, caption=9),
    "creative": Typography(hero=42, h1=26, h2=18, h3=14, body=12, caption=10),
    "professional": Typography(hero=40, h1=28, h2=20, h3=16, body=12, caption=10),
    "minimal": Typography(hero=36, h1=24, h2=18, h3=14, body=11, caption=9),
}

SPACING = {
    "mckinsey": Spacing(page_margin=0.65, section_gap=0.5, card_gap=0.35),
    "cyberpunk": Spacing(page_margin=0.5, section_gap=0.4, card_gap=0.3),
    "creative": Spacing(page_margin=0.7, section_gap=0.6, card_gap=0.4),
    "professional": Spacing(page_margin=0.65, section_gap=0.5, card_gap=0.35),
    "minimal": Spacing(page_margin=0.8, section_gap=0.6, card_gap=0.4),
}
