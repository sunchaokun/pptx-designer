"""Deterministic extraction of template evidence for VI Build Mode."""

from __future__ import annotations

import hashlib
from collections import Counter
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_designer.enterprise.vi_context import normalize_design_context

_EMU_PER_INCH = 914400
_EXTRACTOR_VERSION = "1.0"


@dataclass
class TextZone:
    """Text zone in a slide, including directly stored run evidence."""

    text: str = ""
    x: float = 0
    y: float = 0
    width: float = 0
    height: float = 0
    font_size: float = 12
    font_name: str = ""
    color: str = ""


@dataclass
class SlideDNA:
    """Design evidence for one slide."""

    slide_index: int = 0
    title: str = ""
    text_zones: list[TextZone] = field(default_factory=list)
    images: list[str] = field(default_factory=list)
    shapes: list[dict[str, Any]] = field(default_factory=list)


@dataclass
class DesignDNA:
    """Legacy-compatible, JSON-serializable template evidence."""

    title: str = ""
    slides: list[SlideDNA] = field(default_factory=list)
    colors: dict[str, str] = field(default_factory=dict)
    fonts: dict[str, str] = field(default_factory=dict)
    brand_spec: Any = None
    warnings: list[str] = field(default_factory=list)


@dataclass
class PagePlan:
    """Page plan for rendering."""

    page_type: str = "content"
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    images: list[str] = field(default_factory=list)


def _inches(emu: int) -> float:
    return round(emu / _EMU_PER_INCH, 4)


def _rgb(value: Any) -> str:
    """Read direct RGB evidence without treating theme values as direct RGB."""
    try:
        rgb = value.rgb
    except (AttributeError, TypeError, ValueError):
        return ""
    return f"#{rgb}" if rgb else ""


def _shape_fill(shape: Any) -> str:
    try:
        return _rgb(shape.fill.fore_color)
    except (AttributeError, TypeError, ValueError):
        return ""


def _shape_line(shape: Any) -> str:
    try:
        return _rgb(shape.line.color)
    except (AttributeError, TypeError, ValueError):
        return ""


def _first_run(shape: Any) -> Any | None:
    try:
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.runs:
                return paragraph.runs[0]
    except (AttributeError, TypeError):
        return None
    return None


def _text_zone(shape: Any) -> TextZone | None:
    try:
        value = shape.text_frame.text.strip()
    except (AttributeError, TypeError):
        return None
    if not value:
        return None
    run = _first_run(shape)
    font = run.font if run is not None else None
    size = getattr(font, "size", None)
    return TextZone(
        text=value,
        x=_inches(shape.left),
        y=_inches(shape.top),
        width=_inches(shape.width),
        height=_inches(shape.height),
        font_size=round(size.pt, 2) if size is not None else 12,
        font_name=getattr(font, "name", "") or "",
        color=_rgb(getattr(font, "color", None)),
    )


class DesignDNAExtractor:
    """Extract only directly observable template information from a PPTX."""

    def extract(self, pptx_path: str) -> DesignDNA:
        """Extract legacy-compatible DNA without swallowing source failures."""
        prs = Presentation(pptx_path)
        dna = DesignDNA()
        color_counts: Counter[str] = Counter()
        font_counts: Counter[str] = Counter()

        for slide_index, slide in enumerate(prs.slides):
            slide_dna = SlideDNA(slide_index=slide_index)
            for shape_index, shape in enumerate(slide.shapes):
                fill = _shape_fill(shape)
                line = _shape_line(shape)
                if fill:
                    color_counts[fill] += 1
                if line:
                    color_counts[line] += 1
                slide_dna.shapes.append(
                    {
                        "index": shape_index,
                        "type": str(shape.shape_type),
                        "name": shape.name,
                        "x": _inches(shape.left),
                        "y": _inches(shape.top),
                        "width": _inches(shape.width),
                        "height": _inches(shape.height),
                        "fill": fill,
                        "line": line,
                    }
                )

                if getattr(shape, "has_text_frame", False):
                    zone = _text_zone(shape)
                    if zone is not None:
                        slide_dna.text_zones.append(zone)
                        if zone.font_name:
                            font_counts[zone.font_name] += 1
                        if zone.color:
                            color_counts[zone.color] += 1
                        if not slide_dna.title:
                            slide_dna.title = zone.text
                        if not dna.title:
                            dna.title = zone.text

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_hash = hashlib.sha256(image.blob).hexdigest()
                    slide_dna.images.append(f"image-{image_hash[:12]}.{image.ext}")

            dna.slides.append(slide_dna)

        dna.colors = {
            f"color_{index + 1}": color for index, (color, _) in enumerate(color_counts.most_common())
        }
        dna.fonts = {f"font_{index + 1}": font for index, (font, _) in enumerate(font_counts.most_common())}
        return dna


def _extract_image_references(prs: Presentation) -> list[dict[str, Any]]:
    references = []
    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            image = shape.image
            references.append(
                {
                    "id": f"slide-{slide_index + 1}-image-{shape_index + 1}",
                    "kind": "image",
                    "sha256": hashlib.sha256(image.blob).hexdigest(),
                    "extension": image.ext,
                    "source_slide": slide_index + 1,
                    "bounds": {
                        "left": _inches(shape.left),
                        "top": _inches(shape.top),
                        "width": _inches(shape.width),
                        "height": _inches(shape.height),
                    },
                }
            )
    return references


def _extract_archetypes(prs: Presentation) -> tuple[list[dict[str, Any]], dict[str, dict[str, Any]]]:
    archetypes = []
    components: dict[str, dict[str, Any]] = {}
    slide_area = _inches(prs.slide_width) * _inches(prs.slide_height)
    for slide_index, slide in enumerate(prs.slides):
        pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        component_ids: list[str] = []
        for image_index, picture in enumerate(pictures, start=1):
            component_id = f"photo-panel-{slide_index + 1}-{image_index}"
            component_ids.append(component_id)
            components[component_id] = {
                "kind": "photo_panel",
                "reference_slide": slide_index + 1,
                "image_mode": "cover",
                "bounds": {
                    "left": _inches(picture.left),
                    "top": _inches(picture.top),
                    "width": _inches(picture.width),
                    "height": _inches(picture.height),
                },
            }
        for shape_index, shape in enumerate(slide.shapes, start=1):
            fill = _shape_fill(shape)
            area_ratio = (_inches(shape.width) * _inches(shape.height)) / slide_area if slide_area else 0
            if not fill or area_ratio < 0.1:
                continue
            component_id = f"color-panel-{slide_index + 1}-{shape_index}"
            component_ids.append(component_id)
            components[component_id] = {
                "kind": "color_panel",
                "reference_slide": slide_index + 1,
                "fill": fill,
                "bounds": {
                    "left": _inches(shape.left),
                    "top": _inches(shape.top),
                    "width": _inches(shape.width),
                    "height": _inches(shape.height),
                },
            }
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type != MSO_SHAPE_TYPE.LINE:
                continue
            component_id = f"rule-{slide_index + 1}-{shape_index}"
            component_ids.append(component_id)
            components[component_id] = {
                "kind": "rule",
                "reference_slide": slide_index + 1,
                "fill": _shape_line(shape) or "#FFFFFF",
                "bounds": {
                    "left": _inches(shape.left),
                    "top": _inches(shape.top),
                    "width": _inches(shape.width),
                    "height": max(0.02, _inches(shape.height)),
                },
            }
        archetypes.append(
            {
                "id": f"slide-{slide_index + 1}-{'photo' if pictures else 'text'}",
                "reference_slide": slide_index + 1,
                "required_assets": ["supporting_photo"] if pictures else [],
                "permitted_components": component_ids,
            }
        )
    return archetypes, components


def _typography_from_dna(dna: DesignDNA) -> dict[str, str]:
    fonts = list(dna.fonts.values())
    if not fonts:
        return {}
    return {"heading": fonts[0], "body": fonts[1] if len(fonts) > 1 else fonts[0]}


def _semantic_roles_from_presentation(prs: Presentation, dna: DesignDNA) -> dict[str, str]:
    """Separate text ink from large surface fills before assigning roles."""
    text_colors: Counter[str] = Counter()
    fill_colors: Counter[str] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            fill = _shape_fill(shape)
            if fill:
                fill_colors[fill] += 1
            if getattr(shape, "has_text_frame", False):
                zone = _text_zone(shape)
                if zone is not None and zone.color:
                    text_colors[zone.color] += 1

    fallback = next(iter(dna.colors.values()), "#1D78FA")
    primary = fill_colors.most_common(1)[0][0] if fill_colors else fallback
    ink = text_colors.most_common(1)[0][0] if text_colors else "#111827"
    return {
        "background": "#FFFFFF",
        "surface": "#FFFFFF",
        "ink": ink,
        "primary": primary,
        "data-series-1": primary,
    }


def extract_design_context(pptx_path: str) -> dict[str, Any]:
    """Extract a versioned template context for direct consumption by VI Build."""
    path = Path(pptx_path)
    prs = Presentation(path)
    dna = DesignDNAExtractor().extract(str(path))
    references = _extract_image_references(prs)
    archetypes, components = _extract_archetypes(prs)
    fingerprint = hashlib.sha256(path.read_bytes()).hexdigest()
    semantic_roles = _semantic_roles_from_presentation(prs, dna)
    context = {
        "source": {
            "kind": "template",
            "template_path": str(path),
            "template_fingerprint": fingerprint,
            "extractor_version": _EXTRACTOR_VERSION,
            "confidence": 1.0,
            "warnings": list(dna.warnings),
        },
        "colors": {"primary": semantic_roles["primary"], "text_dark": semantic_roles["ink"]},
        "semantic_roles": semantic_roles,
        "typography": _typography_from_dna(dna),
        "assets": {"references": references, "image_grammar": {"required": False}},
        "components": components,
        "archetypes": archetypes,
        "content_slots": [],
        "locks": [],
        "diagnostics": {"warnings": list(dna.warnings)},
    }
    return normalize_design_context(context)


def extract_design_dna(pptx_path: str) -> dict[str, Any]:
    """Extract legacy DNA as a JSON-serializable compatibility projection."""
    return asdict(DesignDNAExtractor().extract(pptx_path))


__all__ = [
    "DesignDNA",
    "DesignDNAExtractor",
    "PagePlan",
    "SlideDNA",
    "TextZone",
    "extract_design_context",
    "extract_design_dna",
]
