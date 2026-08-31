"""Template-local compiler from VI context to the stable Build contract."""

from __future__ import annotations

from collections.abc import Mapping, Sequence
from copy import deepcopy
from typing import Any

from pptx_designer.enterprise.vi_context import normalize_design_context


class VITemplateAdapter:
    """Compile reviewed template recipes into renderer-neutral ``BuildSpec`` data.

    The adapter owns template selection and provenance.  It deliberately does
    not render slides or import Build component implementations, so adding a
    template cannot add template-specific branches to Build Core.
    """

    def __init__(self, design_context: Mapping[str, Any], *, assets: Mapping[str, Any] | None = None):
        self.context = normalize_design_context(design_context)
        self.assets = dict(assets or {})

    def compile(
        self,
        *,
        page_role: str,
        page_goal: str = "",
        content: Mapping[str, Any] | None = None,
    ) -> dict[str, Any]:
        """Compile a confirmed framework-page rebind contract.

        Content pages deliberately have no archetype route. Their complete
        composition must be Build-owned and pass through :meth:`compile_atomic`.
        """
        content = dict(content or {})
        if page_role == "content":
            raise ValueError("content pages must use compile_atomic(Build-owned atomic_build_plan)")
        framework_page = next(
            (
                item for item in self.context.get("framework_pages", [])
                if isinstance(item, Mapping) and item.get("role") == page_role
            ),
            None,
        )
        if not framework_page:
            raise ValueError(f"no reviewed framework page for page_role={page_role}")
        spec = {
            "kind": "BuildSpec",
            "page_role": page_role,
            "page_goal": page_goal,
            "reference_slide": framework_page.get("reference_slide"),
            "framework_id": framework_page.get("id", f"framework-{page_role}"),
            "render_strategy": "prototype",
            "components": [],
            "assets": deepcopy(content.get("assets", {})),
            "locks": deepcopy(self.context.get("locks", [])),
            "source": {"template_fingerprint": self.context.get("source", {}).get("template_fingerprint")},
        }
        confirmed = {
                str(slot.get("id")): slot
                for slot in self.context.get("content_slots", [])
                if isinstance(slot, Mapping)
                and slot.get("id")
                and slot.get("page_role", page_role) == page_role
            }
        slot_values = content.get("slots", {})
        unknown = sorted(set(slot_values) - set(confirmed))
        if unknown:
            raise ValueError(f"unconfirmed framework slots: {', '.join(unknown)}")
        text_contract = deepcopy(framework_page.get("text_contract", {}))
        if text_contract and not isinstance(text_contract, Mapping):
            raise ValueError(f"framework text contract must be a mapping: {page_role}")
        if text_contract.get("strict", False):
            required = {str(slot.get("id")) for slot in confirmed.values() if slot.get("required", True)}
            missing = sorted(required - set(slot_values))
            if missing:
                raise ValueError(f"unbound required framework slots: {', '.join(missing)}")
        spec["slot_instances"] = [{"slot_id": slot_id, "value": deepcopy(value)} for slot_id, value in slot_values.items()]
        spec["framework_text_contract"] = text_contract
        spec["delivery_origin"] = "framework_rebound"
        return spec

    def compile_atomic(
        self,
        *,
        page_role: str,
        atomic_build_plan: Mapping[str, Any],
        page_goal: str = "",
    ) -> dict[str, Any]:
        """Resolve a Build-authored atomic composition under VI constraints.

        This is intentionally not a second page planner.  ``atomic_build_plan``
        already contains Build's component choices, exact geometry, z-order and
        relation bindings.  The adapter only resolves reviewed visual styles and
        rejects plans that violate the template contract.
        """
        if page_role != "content":
            raise ValueError("atomic Build plans are supported for content pages only")
        plan = deepcopy(dict(atomic_build_plan))
        atoms = plan.get("atoms")
        if not isinstance(atoms, Sequence) or isinstance(atoms, (str, bytes)) or not atoms:
            raise ValueError("atomic content page requires one or more atoms")

        grammar = self.context.get("visual_grammar", {})
        grammar = grammar if isinstance(grammar, Mapping) else {}
        allowed_kinds = set(grammar.get("allowed_atom_kinds", []) or [])
        styles = self.context.get("atom_styles", {})
        styles = styles if isinstance(styles, Mapping) else {}
        atom_ids: set[str] = set()
        components: list[dict[str, Any]] = []
        for ordinal, raw_atom in enumerate(atoms):
            if not isinstance(raw_atom, Mapping):
                raise ValueError("atomic plan contains a non-mapping atom")
            atom = dict(raw_atom)
            atom_id = str(atom.get("id", ""))
            kind = str(atom.get("kind", ""))
            if not atom_id or not kind:
                raise ValueError("atomic atom requires id and kind")
            if atom_id in atom_ids:
                raise ValueError(f"duplicate_atom_id:{atom_id}")
            atom_ids.add(atom_id)
            if allowed_kinds and kind not in allowed_kinds:
                raise ValueError(f"atom_kind_not_allowed:{atom_id}:{kind}")
            bounds = self._bounds(atom.get("geometry", {}))
            self._validate_atomic_bounds(atom_id, bounds, grammar)

            style_ref = atom.get("style_ref")
            style: dict[str, Any] = {}
            if style_ref is not None:
                raw_style = styles.get(str(style_ref))
                if not isinstance(raw_style, Mapping):
                    raise ValueError(f"unknown_atom_style:{style_ref}")
                style = deepcopy(dict(raw_style))
                styled_kind = style.get("kind")
                if styled_kind is not None and str(styled_kind) != kind:
                    raise ValueError(f"atom_style_kind_mismatch:{atom_id}")
            inline_style = atom.get("style", {})
            if inline_style is not None and not isinstance(inline_style, Mapping):
                raise ValueError(f"atom_style_must_be_mapping:{atom_id}")
            style.update(deepcopy(dict(inline_style or {})))
            style["kind"] = kind
            style["bounds"] = bounds
            self._validate_atomic_style(atom_id, kind, style, grammar)
            components.append(
                {
                    "atom_id": atom_id,
                    "component_id": None,
                    "recipe": style,
                    "data": deepcopy(atom.get("data")),
                    "z_index": int(atom.get("z_index", ordinal)),
                    "source": {"style_ref": style_ref, "build_owned": True},
                }
            )

        components.sort(key=lambda item: item["z_index"])
        spec: dict[str, Any] = {
            "kind": "BuildSpec",
            "atomic_build": True,
            "page_role": page_role,
            "page_goal": page_goal,
            "render_strategy": "components",
            "components": components,
            "content_model": deepcopy(plan.get("content_model")),
            "relation_bindings": deepcopy(plan.get("relation_bindings", [])),
            "delivery_origin": "build_atomic",
            "locks": deepcopy(self.context.get("locks", [])),
            "source": {"template_fingerprint": self.context.get("source", {}).get("template_fingerprint")},
        }
        base_id = plan.get("base_id")
        if base_id is not None:
            spec["base_id"] = str(base_id)
            spec["fixed_base"] = self._reviewed_fixed_base(str(base_id))
        return spec

    def render(self, spec: Mapping[str, Any], presentation: Any) -> Any:
        """Render a component BuildSpec through the existing public helpers.

        This bridge is intentionally generic: template-specific decisions are
        already resolved in ``spec``.  It does not add template branches to
        ``core`` or to the reusable Build component implementations.
        """
        if spec.get("kind") != "BuildSpec":
            raise ValueError("renderer requires a BuildSpec")
        if spec.get("render_strategy") == "prototype":
            from pptx_designer.enterprise.prototype import clone_slide_prototype

            reference_slide = int(spec.get("reference_slide", 0))
            if not 1 <= reference_slide <= len(presentation.slides):
                raise ValueError(f"framework reference_slide is unavailable: {reference_slide}")
            slide = clone_slide_prototype(presentation, presentation.slides[reference_slide - 1])
            for instance in spec.get("slot_instances", []):
                slot_id = str(instance["slot_id"])
                slot = next(
                    (item for item in self.context.get("content_slots", [])
                     if isinstance(item, Mapping) and item.get("id") == slot_id),
                    None,
                )
                if slot is None:
                    raise ValueError(f"framework slot is not confirmed: {slot_id}")
                self._bind_framework_text(slide, slot_id, slot, str(instance.get("value", "")))
            self._enforce_framework_text_contract(
                slide,
                spec.get("slot_instances", []),
                spec.get("framework_text_contract", {}),
            )
            return slide
        if spec.get("render_strategy") != "components":
            raise ValueError("unsupported BuildSpec render strategy")
        from pptx_designer.core.build_spec import render_build_spec

        return render_build_spec(spec, presentation, self.context)

    def _reviewed_fixed_base(self, base_id: str) -> dict[str, Any]:
        """Resolve a fixed visual layer without selecting a page composition."""
        base = self.context.get("fixed_bases", {}).get(base_id)
        if base is None:
            return {"id": base_id, "status": "review_required"}
        if not base.get("safe_to_clone", False):
            raise ValueError(f"unsafe fixed base: {base_id}")
        if "reference_slide" in base and "fixed_shape_indices" not in base:
            raise ValueError(f"fixed base {base_id} requires explicit fixed_shape_indices")
        fixed_base = {
            "id": base_id,
            "fixed_objects": deepcopy(base.get("fixed_objects", [])),
            "removed_objects": deepcopy(base.get("removed_objects", [])),
            "unsupported_objects": deepcopy(base.get("unsupported_objects", [])),
        }
        for key in ("reference_slide", "fixed_shape_indices", "exclude_shape_indices"):
            if key in base:
                fixed_base[key] = deepcopy(base[key])
        return fixed_base

    @classmethod
    def _validate_atomic_bounds(
        cls, atom_id: str, bounds: Mapping[str, float], grammar: Mapping[str, Any]
    ) -> None:
        left, top, width, height = (bounds[key] for key in ("left", "top", "width", "height"))
        if width <= 0 or height <= 0 or left < 0 or top < 0 or left + width > 13.333 or top + height > 7.5:
            raise ValueError(f"atom_out_of_bounds:{atom_id}")
        safe_area = grammar.get("safe_area")
        if isinstance(safe_area, Mapping):
            safe = cls._bounds(safe_area)
            if left < safe["left"] or top < safe["top"] or left + width > safe["left"] + safe["width"] or top + height > safe["top"] + safe["height"]:
                raise ValueError(f"atom_outside_safe_area:{atom_id}")
        for zone in grammar.get("forbidden_zones", []) or []:
            if isinstance(zone, Mapping) and cls._boxes_intersect(bounds, cls._bounds(zone)):
                raise ValueError(f"atom_in_forbidden_zone:{atom_id}")

    @staticmethod
    def _boxes_intersect(first: Mapping[str, float], second: Mapping[str, float]) -> bool:
        return (
            max(first["left"], second["left"]) < min(first["left"] + first["width"], second["left"] + second["width"])
            and max(first["top"], second["top"]) < min(first["top"] + first["height"], second["top"] + second["height"])
        )

    @staticmethod
    def _validate_atomic_style(
        atom_id: str, kind: str, style: Mapping[str, Any], grammar: Mapping[str, Any]
    ) -> None:
        minimum = grammar.get("min_font_size")
        if kind in {"text", "multiline_text"} and minimum is not None and float(style.get("font_size", 18)) < float(minimum):
            raise ValueError(f"atom_font_size_below_minimum:{atom_id}")

    @staticmethod
    def _bounds(bounds: Mapping[str, Any]) -> dict[str, float]:
        required = ("left", "top", "width", "height")
        if any(field not in bounds for field in required):
            raise ValueError("component bounds are incomplete")
        return {field: float(bounds[field]) for field in required}

    @classmethod
    def _text_bounds(cls, bounds: Mapping[str, Any]) -> dict[str, float]:
        return cls._bounds(bounds)

    @staticmethod
    def _asset_path(asset: Any) -> str:
        if isinstance(asset, Mapping):
            asset = asset.get("path") or asset.get("url")
        if not isinstance(asset, str) or not asset:
            raise ValueError("BuildSpec image component requires a path or URL")
        return asset

    @staticmethod
    def _bind_framework_text(slide: Any, slot_id: str, slot: Mapping[str, Any], value: str) -> None:
        target = slot.get("target", {})
        shape = None
        if isinstance(target, Mapping) and "shape_index" in target:
            index = int(target["shape_index"])
            if 0 <= index < len(slide.shapes):
                shape = slide.shapes[index]
        elif isinstance(target, Mapping) and target.get("object_id"):
            object_id = str(target["object_id"])
            shape = next((item for item in slide.shapes if item.name == object_id), None)
        if shape is None or not shape.has_text_frame or not shape.text_frame.paragraphs:
            raise ValueError(f"framework slot {slot_id} does not target a text shape")
        paragraph = shape.text_frame.paragraphs[0]
        if paragraph.runs:
            paragraph.runs[0].text = value
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = value
        # A reviewed slot replaces the whole text object. Clear stale runs in
        # later paragraphs, otherwise multi-line template placeholders leak
        # their original sample text into the generated page.
        for extra in shape.text_frame.paragraphs[1:]:
            for run in extra.runs:
                run.text = ""
        text_style = slot.get("text_style", {})
        if text_style:
            if not isinstance(text_style, Mapping):
                raise ValueError(f"framework slot {slot_id} text_style must be a mapping")
            for item in shape.text_frame.paragraphs:
                for run in item.runs:
                    if "font_size" in text_style:
                        from pptx.util import Pt

                        run.font.size = Pt(float(text_style["font_size"]))
                    if "font_name" in text_style:
                        run.font.name = str(text_style["font_name"])

    def _enforce_framework_text_contract(
        self,
        slide: Any,
        instances: Sequence[Mapping[str, Any]],
        contract: Mapping[str, Any] | Any,
    ) -> None:
        """Fail closed when a strict framework page leaves template text unowned.

        Framework pages are allowed to keep explicitly declared decorative text,
        but every other non-empty text shape must be replaced or cleared. This
        prevents template placeholders from silently becoming delivery copy.
        """
        if not isinstance(contract, Mapping) or not contract.get("strict", False):
            return
        replaced: set[int] = set()
        for instance in instances:
            slot_id = str(instance.get("slot_id", ""))
            slot = next(
                (
                    item
                    for item in self.context.get("content_slots", [])
                    if isinstance(item, Mapping) and item.get("id") == slot_id
                ),
                None,
            )
            if slot is not None:
                target = slot.get("target", {})
                if isinstance(target, Mapping) and "shape_index" in target:
                    replaced.add(int(target["shape_index"]))
        clear = {int(index) for index in contract.get("clear_shape_indices", [])}
        preserve = {int(index) for index in contract.get("preserve_shape_indices", [])}
        for index in clear:
            if not 0 <= index < len(slide.shapes):
                raise ValueError(f"framework clear shape index is unavailable: {index}")
            self._clear_framework_text(slide.shapes[index])
        unowned = [
            index
            for index, shape in enumerate(slide.shapes)
            if getattr(shape, "has_text_frame", False)
            and shape.text.strip()
            and index not in replaced | clear | preserve
        ]
        if unowned:
            raise ValueError(f"unowned framework text shapes: {', '.join(map(str, unowned))}")

    @staticmethod
    def _clear_framework_text(shape: Any) -> None:
        if not getattr(shape, "has_text_frame", False):
            raise ValueError("framework clear target is not a text shape")
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = ""

    @staticmethod
    def _blank_layout(presentation: Any) -> Any:
        for layout in presentation.slide_layouts:
            if not layout.placeholders:
                return layout
        return presentation.slide_layouts[0]

__all__ = ["VITemplateAdapter"]
