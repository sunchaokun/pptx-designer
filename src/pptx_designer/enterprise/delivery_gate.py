"""Delivery gate — quality checks for PPT output."""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any


@dataclass
class CheckItem:
    """Quality check item."""

    category: str
    check_id: str
    severity: str  # "fatal", "warning", "review"
    slide_index: int | None = None
    message: str = ""
    detail: str = ""
    auto_fixable: bool = False


@dataclass
class QualityReport:
    """Quality check report."""

    total_slides: int = 0
    total_checks: int = 0
    passed: int = 0
    fatals: int = 0
    warnings: int = 0
    auto_fixed: int = 0
    checks: list[CheckItem] = field(default_factory=list)


class DeliveryGate:
    """Quality gate for PPT delivery."""

    def __init__(self):
        pass

    def check(self, pptx_path: str, dna: Any = None, plans: list | None = None) -> QualityReport:
        """Run quality checks on a PPTX file.

        Args:
            pptx_path: Path to .pptx file
            dna: Design DNA (optional)
            plans: Page plans (optional)

        Returns:
            QualityReport
        """
        from pptx import Presentation

        report = QualityReport()

        try:
            prs = Presentation(pptx_path)
            report.total_slides = len(prs.slides)

            # Check for blank slides
            for i, slide in enumerate(prs.slides):
                if len(slide.shapes) == 0:
                    report.checks.append(CheckItem(
                        category="content",
                        check_id="blank_page",
                        severity="fatal",
                        slide_index=i,
                        message=f"Slide {i + 1} is blank",
                    ))
                    report.fatals += 1
                else:
                    report.passed += 1

            report.total_checks = len(report.checks)

        except Exception as e:
            report.checks.append(CheckItem(
                category="error",
                check_id="read_error",
                severity="fatal",
                message=f"Error reading PPTX: {str(e)}",
            ))
            report.fatals += 1

        return report

    def auto_fix(self, pptx_path: str, dna: Any = None, plans: list | None = None,
                 report: QualityReport | None = None) -> None:
        """Auto-fix quality issues.

        Args:
            pptx_path: Path to .pptx file
            dna: Design DNA (optional)
            plans: Page plans (optional)
            report: Quality report to update
        """
        pass

    def format_report(self, report: QualityReport) -> str:
        """Format quality report as string.

        Args:
            report: Quality report

        Returns:
            Formatted string
        """
        lines = [
            f"Quality Report: {report.total_checks} checks",
            f"  Passed: {report.passed}",
            f"  Fatals: {report.fatals}",
            f"  Warnings: {report.warnings}",
        ]
        for check in report.checks:
            lines.append(f"  [{check.severity}] {check.message}")
        return "\n".join(lines)
