"""Reusable slide-prototype operations for VI Build.

PowerPoint relationships belong to a slide part, not to individual shape XML.
Consequently a copied picture element cannot retain its source ``rId`` on a
new slide.  This module creates a page from a template prototype while
rebuilding those relationships for the destination slide.
"""

from __future__ import annotations

import os
import posixpath
import re
import tempfile
import zipfile
from copy import deepcopy
from typing import Any
from xml.etree import ElementTree

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


def clone_slide_prototype(presentation: Any, source_slide: Any) -> Any:
    """Append an editable copy of ``source_slide`` to ``presentation``.

    Shapes retain their native OOXML, preserving typography, geometry, and
    decorative elements. Picture relationships are re-created on the target
    slide, so copied images render correctly rather than pointing at a missing
    relationship id from the source page.

    The source slide must belong to ``presentation``. Embedded pictures are
    supported; other relationship-bearing objects are deliberately left for
    a later capability extension rather than silently altering them.
    """
    if source_slide.part.package is not presentation.part.package:
        raise ValueError("source_slide must belong to presentation")

    target_slide = presentation.slides.add_slide(_blank_layout(presentation))
    _remove_layout_placeholders(target_slide)
    copy_slide_shapes(presentation, source_slide, target_slide)
    return target_slide


def copy_slide_shapes(
    presentation: Any,
    source_slide: Any,
    target_slide: Any,
    shape_indices: list[int] | None = None,
    *,
    exclude_indices: list[int] | None = None,
) -> list[str]:
    """Copy selected top-level template shapes into an existing slide.

    This is the fixed-base primitive used by VI Build. It copies native
    objects, rebinds picture relationships, preserves source order, and
    rejects groups because nested relationships cannot yet be safely rebuilt.
    """
    if source_slide.part.package is not presentation.part.package:
        raise ValueError("source_slide must belong to presentation")
    excluded = set(exclude_indices or [])
    selected = range(len(source_slide.shapes)) if shape_indices is None else shape_indices
    copied: list[str] = []
    for index in selected:
        if index in excluded:
            continue
        if not 0 <= int(index) < len(source_slide.shapes):
            raise ValueError(f"fixed base shape index is unavailable: {index}")
        source_shape = source_slide.shapes[int(index)]
        if source_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            raise ValueError(f"fixed base cannot safely clone grouped shape: {source_shape.name}")
        copied_element = deepcopy(source_shape.element)
        if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _rebind_picture_relationship(source_slide, target_slide, source_shape, copied_element)
        target_slide.shapes._spTree.insert_element_before(copied_element, "p:extLst")
        copied.append(source_shape.name)
    return copied


def prune_unreferenced_slide_parts(path: str) -> None:
    """Remove orphaned slide XML parts left after deleting source slides.

    ``python-pptx`` updates the presentation slide list but intentionally does
    not garbage-collect package parts. Some previewers reject those orphaned
    parts even though PowerPoint itself can reopen the file.
    """
    source = os.path.abspath(path)
    with zipfile.ZipFile(source, "r") as archive:
        entries = {name: archive.read(name) for name in archive.namelist()}
    rels_name = "ppt/_rels/presentation.xml.rels"
    if rels_name not in entries:
        return
    root = ElementTree.fromstring(entries[rels_name])
    namespace = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    referenced: set[str] = set()
    for relation in root.findall("r:Relationship", namespace):
        if relation.get("Type", "").endswith("/slide"):
            target = relation.get("Target", "")
            referenced.add(posixpath.normpath(posixpath.join("ppt", target)).lstrip("/"))
    orphaned = {
        name for name in entries
        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name) and name not in referenced
    }
    if not orphaned:
        return
    for slide in orphaned:
        entries.pop(slide, None)
        entries.pop(f"ppt/slides/_rels/{posixpath.basename(slide)}.rels", None)
    folder = os.path.dirname(source)
    fd, temporary = tempfile.mkstemp(prefix="pptx-clean-", suffix=".pptx", dir=folder)
    os.close(fd)
    try:
        with zipfile.ZipFile(temporary, "w", zipfile.ZIP_DEFLATED) as archive:
            for name, payload in entries.items():
                archive.writestr(name, payload)
        os.replace(temporary, source)
    finally:
        if os.path.exists(temporary):
            os.unlink(temporary)


def _rebind_picture_relationship(source_slide: Any, target_slide: Any, source_shape: Any, copied_element: Any) -> None:
    """Point a copied picture element at an image relationship on its new slide."""
    old_rid = source_shape._element.blipFill.blip.rEmbed
    image_part = source_slide.part.related_part(old_rid)
    new_rid = target_slide.part.relate_to(image_part, RT.IMAGE)
    copied_element.blipFill.blip.rEmbed = new_rid




def _blank_layout(presentation: Any) -> Any:
    """Prefer a blank layout so copied shape indexes remain deterministic."""
    for layout in presentation.slide_layouts:
        if not layout.placeholders:
            return layout
    return presentation.slide_layouts[0]


def _remove_layout_placeholders(slide: Any) -> None:
    """Remove fallback-layout placeholders before inserting prototype shapes."""
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        element = shape.element
        element.getparent().remove(element)
