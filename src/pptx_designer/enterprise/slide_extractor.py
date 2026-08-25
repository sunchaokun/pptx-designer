"""Slide extractor — extract content from existing PPT files."""

from __future__ import annotations

import os


class SlideExtractor:
    """Extracts content and layout from existing .pptx files."""

    def __init__(self, temp_dir: str | None = None):
        self._temp_dir = temp_dir or os.path.join(os.path.expanduser("~"), ".pptx-designer", "temp")

    def extract(self, pptx_path: str) -> list[dict]:
        """Extract content from a PPTX file.

        Args:
            pptx_path: Path to .pptx file

        Returns:
            List of page dicts
        """
        from pptx import Presentation

        prs = Presentation(pptx_path)
        pages = []

        for i, slide in enumerate(prs.slides):
            page = {
                "title": "",
                "bullets": [],
                "goal": "content",
                "slide_index": i,
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        if shape.shape_type == 13:  # Placeholder
                            if hasattr(shape, "placeholder_format"):
                                if shape.placeholder_format.idx == 0:
                                    page["title"] = text
                                else:
                                    page["bullets"].append(text)
                        else:
                            if not page["title"]:
                                page["title"] = text
                            else:
                                page["bullets"].append(text)

            pages.append(page)

        return pages

    def extract_all(self, pptx_path: str) -> list[dict]:
        """Extract all slides from a PPTX file.

        Args:
            pptx_path: Path to .pptx file

        Returns:
            List of page dicts
        """
        return self.extract(pptx_path)
