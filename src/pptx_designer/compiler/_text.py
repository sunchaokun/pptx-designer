"""SVG text rendering with baseline-aware positioning and <tspan> support.

Extracts text rendering from _compiler.py and adds:
- Pillow-based text measurement (with graceful fallback to estimate_text_size)
- SVG baseline (dominant-baseline, alignment-baseline) → PPT vertical anchor mapping
- Multi-line <tspan> support — each tspan with x/y becomes a new paragraph;
  tspans without x/y are inline runs within the same paragraph
- font-family → PPT font mapping
"""

from __future__ import annotations

import re
from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from pptx_designer.diagrams.text_measurer import estimate_text_size

from ._affine import Affine

_FONT_MAP: dict[str, str] = {
    "arial": "Arial",
    "helvetica": "Arial",
    "sans-serif": "Arial",
    "times new roman": "Times New Roman",
    "times": "Times New Roman",
    "serif": "Times New Roman",
    "georgia": "Georgia",
    "courier": "Courier New",
    "courier new": "Courier New",
    "monospace": "Courier New",
    "consolas": "Consolas",
    "verdana": "Verdana",
    "tahoma": "Tahoma",
    "calibri": "Calibri",
    "impact": "Impact",
    "comic sans ms": "Comic Sans MS",
}

_ANCHOR_MAP: dict[str, PP_ALIGN] = {
    "middle": PP_ALIGN.CENTER,
    "end": PP_ALIGN.RIGHT,
    "start": PP_ALIGN.LEFT,
}


def _opacity_pct(el, paint: str = "fill") -> int:
    """Return effective element/paint opacity as a PPTX percentage."""
    try:
        element_opacity = float(el.get("opacity", "1"))
        paint_opacity = float(el.get(f"{paint}-opacity", "1"))
    except (TypeError, ValueError):
        return 100
    return max(0, min(100, round(element_opacity * paint_opacity * 100)))


def _apply_run_alpha(run, alpha_pct: int) -> None:
    """Apply alpha to a run's existing RGB solid fill."""
    if alpha_pct >= 100:
        return
    r_pr = run._r.get_or_add_rPr()
    solid_fill = r_pr.find(qn("a:solidFill"))
    if solid_fill is None:
        return
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    existing = srgb.find(qn("a:alpha"))
    if existing is not None:
        srgb.remove(existing)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(alpha_pct * 1000))
    srgb.append(alpha)


@dataclass
class TextMetrics:
    width_inches: float
    height_inches: float
    ascent_ratio: float = 0.8
    descent_ratio: float = 0.2


@dataclass
class _SpanSpec:
    text: str = ""
    x: float | None = None
    y: float | None = None
    dx: float = 0.0
    dy: float = 0.0
    font_size: float | None = None
    font_family: str | None = None
    fill: str | None = None
    stroke: str | None = None
    stroke_width: float = 0.0
    bold: bool = False
    italic: bool = False
    is_new_line: bool = False


def _resolve_font_family(raw: str | None) -> str:
    if not raw:
        return "Calibri"
    families = re.split(r"[,\s]+", raw.strip("'\""))
    for f in families:
        key = f.strip().lower()
        if key in _FONT_MAP:
            return _FONT_MAP[key]
    return families[0].strip().strip("'\"")


_BASELINE_MAP: dict[str, MSO_ANCHOR] = {
    "auto": MSO_ANCHOR.MIDDLE,
    "alphabetic": MSO_ANCHOR.BOTTOM,
    "text-after-edge": MSO_ANCHOR.BOTTOM,
    "text-before-edge": MSO_ANCHOR.TOP,
    "central": MSO_ANCHOR.MIDDLE,
    "middle": MSO_ANCHOR.MIDDLE,
    "hanging": MSO_ANCHOR.TOP,
    "mathematical": MSO_ANCHOR.MIDDLE,
}

_BASELINE_OFFSET: dict[str, str] = {
    "auto": "middle",
    "alphabetic": "baseline",
    "text-after-edge": "descent",
    "text-before-edge": "top",
    "central": "middle",
    "middle": "middle",
    "hanging": "top",
    "mathematical": "middle",
}


def _resolve_baseline(el) -> MSO_ANCHOR:
    db = el.get("dominant-baseline", el.get("alignment-baseline", "auto"))
    return _BASELINE_MAP.get(db.lower(), MSO_ANCHOR.MIDDLE)


def _resolve_baseline_mode(el) -> str:
    """Return y-anchor mode for precise vertical positioning."""
    db = el.get("dominant-baseline", el.get("alignment-baseline", "auto"))
    return _BASELINE_OFFSET.get(db.lower(), "middle")


def _compute_text_top(iy: float, metrics: TextMetrics, v_anchor, baseline_mode: str = "middle") -> float:
    """Compute textbox top position from SVG y coordinate and baseline.

    SVG y coordinate semantics by dominant-baseline:
      - alphabetic: y is at the baseline (ascent above, descent below)
      - text-before-edge / hanging: y is at the top of the em box
      - central / middle / auto: y is at the vertical center
      - text-after-edge: y is at the bottom of the em box

    PPT textbox top = iy minus the portion of the glyph above the y anchor.
    """
    asc = metrics.height_inches * metrics.ascent_ratio
    desc = metrics.height_inches * metrics.descent_ratio

    if baseline_mode == "top":
        return iy
    elif baseline_mode == "baseline":
        return iy - asc
    elif baseline_mode == "descent":
        return iy - asc - desc
    else:
        return iy - asc / 2.0


def _has_cjk(text: str) -> bool:
    for ch in text:
        cp = ord(ch)
        if 0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF or 0x3000 <= cp <= 0x303F or 0xFF00 <= cp <= 0xFFEF:
            return True
    return False


def _measure_text(content: str, font_size_pt: float, font_family: str, max_width_inches: float) -> TextMetrics:
    font = None
    try:
        from PIL import ImageFont

        try:
            font = ImageFont.truetype(font_family + ".ttf", int(font_size_pt))
        except OSError:
            # Pillow's default bitmap font has a fixed, unrelated size. It
            # cannot measure text that will be emitted at the requested
            # Office point size.
            font = None
    except ImportError:
        pass

    if font is not None:
        bbox = font.getbbox(content)
        w_px = bbox[2] - bbox[0]
        h_px = bbox[3] - bbox[1]
        ascent, descent = font.getmetrics()
        total_h = ascent + descent

        px_per_inch = 96.0
        w_in = w_px / px_per_inch
        h_in = (total_h if total_h > 0 else h_px) / px_per_inch

        return TextMetrics(
            width_inches=max(w_in, 0.5),
            height_inches=max(h_in, 0.3),
            ascent_ratio=ascent / total_h if total_h > 0 else 0.8,
            descent_ratio=descent / total_h if total_h > 0 else 0.2,
        )

    w_est, h_est = estimate_text_size(content, max(8, int(font_size_pt)), max_width_inches, font_family)
    return TextMetrics(
        # ``estimate_text_size`` is intentionally compact for card labels.
        # SVG text has no implicit wrapping, and Office can use wider glyph
        # metrics than that estimator (notably Arial on Windows/LibreOffice).
        # Reserve a conservative width when a matching font file is absent.
        width_inches=max(w_est * 1.4, 0.5),
        height_inches=max(h_est, 0.3),
        ascent_ratio=0.75,
        descent_ratio=0.25,
    )


def _resolve_fill(raw: str | None, C: dict, resolve_color_fn) -> str:  # noqa: N803
    if not raw or raw == "none":
        return "#000000"
    if raw.startswith("url(#"):
        return "#000000"
    resolved = resolve_color_fn(raw, C, "")
    return resolved if resolved else "#000000"


def _parse_font_size(raw: str | None, parent_fs: float) -> float:
    if not raw:
        return parent_fs
    clean = re.sub(r"[^\d.]", "", raw)
    return float(clean) if clean else parent_fs


def _parse_font_weight(raw: str | None) -> bool:
    if not raw:
        return False
    return raw not in ("normal", "100", "200", "300")


def _collect_spans(
    el,
    parent_fs: float,
    parent_ff: str,
    parent_fill: str,
    C: dict,
    resolve_color_fn,  # noqa: N803
    parent_bold: bool = False,
    parent_italic: bool = False,
    parent_stroke: str | None = None,
    parent_stroke_width: float = 0.0,
) -> list[_SpanSpec]:
    spans: list[_SpanSpec] = []

    # Check parent element's own font-weight/font-style for direct text content
    parent_bold = parent_bold or _parse_font_weight(el.get("font-weight"))
    parent_italic = parent_italic or el.get("font-style") == "italic"
    stroke_raw = el.get("stroke")
    if stroke_raw and not parent_stroke:
        parent_stroke = _resolve_fill(stroke_raw, C, resolve_color_fn)
    parent_stroke_width = float(el.get("stroke-width", parent_stroke_width))

    direct_text = el.text
    if direct_text and direct_text.strip():
        spans.append(
            _SpanSpec(
                text=direct_text.strip(),
                font_size=parent_fs,
                font_family=parent_ff,
                fill=parent_fill,
                stroke=parent_stroke,
                stroke_width=parent_stroke_width,
                bold=parent_bold,
                italic=parent_italic,
            )
        )

    for child in el:
        tag = child.tag.split("}")[-1] if isinstance(child.tag, str) else ""
        if tag != "tspan":
            tail = child.tail
            if tail and tail.strip():
                spans.append(
                    _SpanSpec(
                        text=tail.strip(),
                        font_size=parent_fs,
                        font_family=parent_ff,
                        fill=parent_fill,
                        stroke=parent_stroke,
                        stroke_width=parent_stroke_width,
                        bold=parent_bold,
                        italic=parent_italic,
                    )
                )
            continue

        fs = _parse_font_size(child.get("font-size"), parent_fs)
        ff = _resolve_font_family(child.get("font-family")) if child.get("font-family") else parent_ff
        fill = _resolve_fill(child.get("fill", parent_fill), C, resolve_color_fn)
        child_bold = _parse_font_weight(child.get("font-weight"))
        child_italic = child.get("font-style") == "italic"
        # Inherit from parent if child doesn't specify
        bold = child_bold or parent_bold
        italic = child_italic or parent_italic
        stroke = _resolve_fill(child.get("stroke"), C, resolve_color_fn) if child.get("stroke") else parent_stroke
        stroke_width = float(child.get("stroke-width", parent_stroke_width))

        has_x = child.get("x") is not None
        has_y = child.get("y") is not None
        dy_val = float(child.get("dy", "0"))
        is_new_line = has_x or has_y or abs(dy_val) > parent_fs * 0.8

        span = _SpanSpec(
            text=(child.text or "").strip(),
            x=float(child.get("x")) if has_x else None,
            y=float(child.get("y")) if has_y else None,
            dx=float(child.get("dx", "0")),
            dy=dy_val,
            font_size=fs,
            font_family=ff,
            fill=fill,
            stroke=stroke,
            stroke_width=stroke_width,
            bold=bold,
            italic=italic,
            is_new_line=is_new_line,
        )
        spans.append(span)

        tail = child.tail
        if tail and tail.strip():
            spans.append(
                _SpanSpec(
                    text=tail.strip(),
                    font_size=parent_fs,
                    font_family=parent_ff,
                    fill=parent_fill,
                    stroke=parent_stroke,
                    stroke_width=parent_stroke_width,
                    bold=parent_bold,
                    italic=parent_italic,
                )
            )

    return spans


def _group_spans_into_lines(spans: list[_SpanSpec]) -> list[list[_SpanSpec]]:
    lines: list[list[_SpanSpec]] = []
    current: list[_SpanSpec] = []
    for sp in spans:
        if sp.is_new_line and current:
            lines.append(current)
            current = []
        current.append(sp)
    if current:
        lines.append(current)
    return lines if lines else [[]]


def render_svg_text(
    el,
    tf: Affine,
    to_inches_fn,
    slide,
    C: dict,  # noqa: N803
    resolve_color_fn,
    features: set,
    svg_w: float = 0.0,
    svg_h: float = 0.0,
    slide_w: float = 0.0,
    slide_h: float = 0.0,
) -> None:
    features.add("text")

    x = float(el.get("x", 0))
    y = float(el.get("y", 0))
    ix, iy = to_inches_fn(*tf.apply(x, y))

    parent_fs = _parse_font_size(el.get("font-size"), 14.0)
    parent_ff = _resolve_font_family(el.get("font-family"))
    parent_fill = _resolve_fill(el.get("fill"), C, resolve_color_fn)
    parent_stroke = _resolve_fill(el.get("stroke"), C, resolve_color_fn) if el.get("stroke") else None
    parent_stroke_width = float(el.get("stroke-width", "0"))
    fill_alpha = _opacity_pct(el, "fill")
    stroke_alpha = _opacity_pct(el, "stroke")
    anchor = el.get("text-anchor", "start")
    v_anchor = _resolve_baseline(el)
    baseline_mode = _resolve_baseline_mode(el)

    scale = slide_w / svg_w if svg_w > 0 and slide_w > 0 else 1.0
    scaled_fs = max(parent_fs * scale * 72.0, 6.0)

    spans = _collect_spans(
        el,
        parent_fs,
        parent_ff,
        parent_fill,
        C,
        resolve_color_fn,
        parent_stroke=parent_stroke,
        parent_stroke_width=parent_stroke_width,
    )

    has_tspan_children = any(isinstance(c.tag, str) and c.tag.split("}")[-1] == "tspan" for c in el)

    if not has_tspan_children:
        content = "".join(el.itertext()).strip()
        if not content:
            return
        _render_simple_text(
            content,
            ix,
            iy,
            scaled_fs,
            parent_ff,
            parent_fill,
            anchor,
            v_anchor,
            baseline_mode,
            el,
            slide,
            C,
            resolve_color_fn,
            stroke=parent_stroke,
            stroke_width=parent_stroke_width,
            fill_alpha=fill_alpha,
            stroke_alpha=stroke_alpha,
        )
        return

    if not any(s.text for s in spans):
        return

    _render_tspan_text(
        spans,
        ix,
        iy,
        scaled_fs,
        parent_ff,
        parent_fill,
        anchor,
        v_anchor,
        baseline_mode,
        el,
        slide,
        to_inches_fn,
        tf,
        C,
        resolve_color_fn,
        fill_alpha=fill_alpha,
        stroke_alpha=stroke_alpha,
        svg_to_pt=scale * 72.0,
    )


def _render_simple_text(
    content: str,
    ix: float,
    iy: float,
    fs: float,
    ff: str,
    fill: str,
    anchor: str,
    v_anchor,
    baseline_mode: str,
    el,
    slide,
    C: dict,
    resolve_color_fn,  # noqa: N803
    stroke: str | None = None,
    stroke_width: float = 0.0,
    fill_alpha: int = 100,
    stroke_alpha: int = 100,
) -> None:
    metrics = _measure_text(content, fs, ff, 8.0)

    # SVG text is allowed to overflow viewBox; textbox width uses natural measurement.
    width = metrics.width_inches + 0.2

    if anchor == "middle":
        left = ix - width / 2
    elif anchor == "end":
        left = ix - width
    else:
        left = ix - 0.1

    top = _compute_text_top(iy, metrics, v_anchor, baseline_mode)

    height = metrics.height_inches * 1.5

    # Render outline layer first (will be behind fill layer)
    if stroke and stroke_width > 0:
        # The outline layer is larger and bolder so it peeks out from behind the fill
        outline_fs = fs + stroke_width * 2  # Larger font size for visible outline
        outline_bold = True  # Bold for thicker outline effect
        tb_outline = slide.shapes.add_textbox(
            Inches(left - 0.05), Inches(top - 0.05), Inches(width + 0.1), Inches(height + 0.1)
        )
        tf_outline = tb_outline.text_frame
        tf_outline.word_wrap = False
        tf_outline.vertical_anchor = v_anchor
        p_outline = tf_outline.paragraphs[0]
        p_outline.alignment = _ANCHOR_MAP.get(anchor, PP_ALIGN.LEFT)
        run_outline = p_outline.add_run()
        run_outline.text = content
        run_outline.font.size = Pt(outline_fs)
        run_outline.font.color.rgb = RGBColor.from_string(stroke.lstrip("#"))
        _apply_run_alpha(run_outline, stroke_alpha)
        run_outline.font.name = ff
        run_outline.font.bold = outline_bold
        bold = el.get("font-weight")
        if bold and bold not in ("normal", "100", "200", "300"):
            run_outline.font.bold = True
        italic = el.get("font-style")
        if italic == "italic":
            run_outline.font.italic = True

    # Render fill layer on top (will cover most of outline, but edges peek out)
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf_el = tb.text_frame
    tf_el.word_wrap = False
    tf_el.vertical_anchor = v_anchor
    p = tf_el.paragraphs[0]
    p.alignment = _ANCHOR_MAP.get(anchor, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = content
    run.font.size = Pt(fs)
    run.font.color.rgb = RGBColor.from_string(fill.lstrip("#"))
    _apply_run_alpha(run, fill_alpha)
    run.font.name = ff

    bold = el.get("font-weight")
    if bold and bold not in ("normal", "100", "200", "300"):
        run.font.bold = True

    italic = el.get("font-style")
    if italic == "italic":
        run.font.italic = True


def _render_tspan_text(
    spans: list[_SpanSpec],
    ix: float,
    iy: float,
    parent_fs: float,
    parent_ff: str,
    parent_fill: str,
    anchor: str,
    v_anchor,
    baseline_mode: str,
    el,
    slide,
    to_inches_fn,
    tf: Affine,
    C: dict,
    resolve_color_fn,  # noqa: N803
    fill_alpha: int = 100,
    stroke_alpha: int = 100,
    svg_to_pt: float = 1.0,
) -> None:
    lines = _group_spans_into_lines(spans)

    all_text = " ".join(s.text for line in lines for s in line if s.text)
    metrics = _measure_text(all_text, parent_fs, parent_ff, 8.0)
    line_h = metrics.height_inches * 1.3

    width = metrics.width_inches + 0.4

    if anchor == "middle":
        left = ix - width / 2
    elif anchor == "end":
        left = ix - width
    else:
        left = ix - 0.1

    top = _compute_text_top(iy, metrics, v_anchor, baseline_mode)
    if baseline_mode == "descent":
        top = top - line_h * (len(lines) - 1)

    height = line_h * len(lines) + 0.2

    # Check if any span has stroke
    has_stroke = any(sp.stroke and sp.stroke_width > 0 for sp in spans)

    # Render outline layer first (if stroke is set)
    if has_stroke:
        max_sw = max((sp.stroke_width for sp in spans if sp.stroke and sp.stroke_width > 0), default=0)
        max_stroke = next((sp.stroke for sp in spans if sp.stroke and sp.stroke_width > 0), None)
        if max_stroke and max_sw > 0:
            offset = max_sw * 0.015
            tb_outline = slide.shapes.add_textbox(
                Inches(left - offset), Inches(top - offset), Inches(width + offset * 2), Inches(height + offset * 2)
            )
            tf_outline = tb_outline.text_frame
            tf_outline.word_wrap = True
            tf_outline.vertical_anchor = v_anchor

            first_para_outline = True
            for line in lines:
                if first_para_outline:
                    p_o = tf_outline.paragraphs[0]
                    first_para_outline = False
                else:
                    p_o = tf_outline.add_paragraph()

                p_o.alignment = _ANCHOR_MAP.get(anchor, PP_ALIGN.LEFT)
                p_o.space_after = Pt(0)

                for sp in line:
                    if not sp.text:
                        continue
                    run_o = p_o.add_run()
                    run_o.text = sp.text
                    run_o.font.size = Pt(max((sp.font_size or 0) * svg_to_pt, 6.0))
                    run_o.font.color.rgb = RGBColor.from_string((sp.stroke or max_stroke).lstrip("#"))
                    _apply_run_alpha(run_o, stroke_alpha)
                    run_o.font.name = sp.font_family or parent_ff
                    if sp.bold:
                        run_o.font.bold = True
                    if sp.italic:
                        run_o.font.italic = True

    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf_el = tb.text_frame
    tf_el.word_wrap = True
    tf_el.vertical_anchor = v_anchor

    first_para = True
    prev_line_iy = None
    for line in lines:
        if first_para:
            p = tf_el.paragraphs[0]
            first_para = False
        else:
            p = tf_el.add_paragraph()

        p.alignment = _ANCHOR_MAP.get(anchor, PP_ALIGN.LEFT)
        p.space_after = Pt(0)

        first_sp = next((s for s in line if s.text), None)
        if first_sp and first_sp.y is not None and prev_line_iy is not None:
            line_iy = to_inches_fn(*tf.apply(first_sp.x or 0, first_sp.y))[1]
            gap = line_iy - prev_line_iy
            if gap > 0:
                p.space_before = Pt(max(0, gap * 72 - line_h * 72 * 0.3))
        elif first_sp and first_sp.dy != 0.0 and prev_line_iy is not None:
            dy_inches = first_sp.dy * (line_h / parent_fs) if parent_fs > 0 else 0
            if dy_inches > 0:
                p.space_before = Pt(max(0, dy_inches * 72 - line_h * 72 * 0.3))
        if first_sp and first_sp.y is not None:
            prev_line_iy = to_inches_fn(*tf.apply(first_sp.x or 0, first_sp.y))[1]
        else:
            prev_line_iy = None

        for sp in line:
            if not sp.text:
                continue

            if sp.dx != 0.0:
                dx_pt = sp.dx * (parent_fs / 14.0) * 0.5
                spacer = p.add_run()
                spacer.text = " "
                spacer.font.size = Pt(max((sp.font_size or 0) * svg_to_pt, 6.0))
                spacer_r_pr = spacer._r.get_or_add_rPr()
                spacer_r_pr.set("spc", str(int(dx_pt * 100)))

            run = p.add_run()
            run.text = sp.text
            run.font.size = Pt(max((sp.font_size or 0) * svg_to_pt, 6.0))
            run.font.color.rgb = RGBColor.from_string((sp.fill or parent_fill).lstrip("#"))
            _apply_run_alpha(run, fill_alpha)
            run.font.name = sp.font_family or parent_ff
            if sp.bold:
                run.font.bold = True
            if sp.italic:
                run.font.italic = True
