"""SVGCompiler — compile SVG subset to native editable PPTX shapes.

Public API::

    from pptx_designer.compiler import SVGCompiler, SVGCompileError, SVGResult
    result = SVGCompiler(C=context).compile(svg_text, slide, rect)
"""

from __future__ import annotations

import math
import re
import time
from dataclasses import dataclass, field

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from pptx_designer.effects.shape_effects import set_solid_fill_with_alpha
from pptx_designer.renderer.boolean_shapes import bool_shape
from pptx_designer.renderer.freeform_builder import FreeformBuilder

from ._affine import Affine, parse_transform
from ._dash import apply_stroke_style, parse_stroke_style
from ._errors import SVGCompileError
from ._ir import SVGIRDocument, build_svg_ir
from ._paint import GradientDef
from ._paint import _parse_percent_or_float as _parse_pct
from ._paint import apply_gradient as _apply_gradient
from ._paint import resolve_paint as _resolve_paint
from ._path import parse_path, to_beziers
from ._sanitizer import sanitize
from ._text import render_svg_text as _render_svg_text

SVG_NS = "http://www.w3.org/2000/svg"
SVG = f"{{{SVG_NS}}}"


@dataclass
class _LineCmd:
    x: float
    y: float


@dataclass
class _CubicCmd:
    x1: float
    y1: float
    x2: float
    y2: float
    x3: float
    y3: float


_PathCmd = _LineCmd | _CubicCmd

# ─────────────────────────── color helpers ───────────────────────────

_NAMED_COLORS: dict[str, str] = {
    "aliceblue": "F0F8FF",
    "antiquewhite": "FAEBD7",
    "aqua": "00FFFF",
    "aquamarine": "7FFFD4",
    "azure": "F0FFFF",
    "beige": "F5F5DC",
    "bisque": "FFE4C4",
    "black": "000000",
    "blanchedalmond": "FFEBCD",
    "blue": "0000FF",
    "blueviolet": "8A2BE2",
    "brown": "A52A2A",
    "burlywood": "DEB887",
    "cadetblue": "5F9EA0",
    "chartreuse": "7FFF00",
    "chocolate": "D2691E",
    "coral": "FF7F50",
    "cornflowerblue": "6495ED",
    "cornsilk": "FFF8DC",
    "crimson": "DC143C",
    "cyan": "00FFFF",
    "darkblue": "00008B",
    "darkcyan": "008B8B",
    "darkgoldenrod": "B8860B",
    "darkgray": "A9A9A9",
    "darkgreen": "006400",
    "darkgrey": "A9A9A9",
    "darkkhaki": "BDB76B",
    "darkmagenta": "8B008B",
    "darkolivegreen": "556B2F",
    "darkorange": "FF8C00",
    "darkorchid": "9932CC",
    "darkred": "8B0000",
    "darksalmon": "E9967A",
    "darkseagreen": "8FBC8F",
    "darkslateblue": "483D8B",
    "darkslategray": "2F4F4F",
    "darkslategrey": "2F4F4F",
    "darkturquoise": "00CED1",
    "darkviolet": "9400D3",
    "deeppink": "FF1493",
    "deepskyblue": "00BFFF",
    "dimgray": "696969",
    "dimgrey": "696969",
    "dodgerblue": "1E90FF",
    "firebrick": "B22222",
    "floralwhite": "FFFAF0",
    "forestgreen": "228B22",
    "fuchsia": "FF00FF",
    "gainsboro": "DCDCDC",
    "ghostwhite": "F8F8FF",
    "gold": "FFD700",
    "goldenrod": "DAA520",
    "gray": "808080",
    "green": "008000",
    "greenyellow": "ADFF2F",
    "grey": "808080",
    "honeydew": "F0FFF0",
    "hotpink": "FF69B4",
    "indianred": "CD5C5C",
    "indigo": "4B0082",
    "ivory": "FFFFF0",
    "khaki": "F0E68C",
    "lavender": "E6E6FA",
    "lavenderblush": "FFF0F5",
    "lawngreen": "7CFC00",
    "lemonchiffon": "FFFACD",
    "lightblue": "ADD8E6",
    "lightcoral": "F08080",
    "lightcyan": "E0FFFF",
    "lightgoldenrodyellow": "FAFAD2",
    "lightgray": "D3D3D3",
    "lightgreen": "90EE90",
    "lightgrey": "D3D3D3",
    "lightpink": "FFB6C1",
    "lightsalmon": "FFA07A",
    "lightseagreen": "20B2AA",
    "lightskyblue": "87CEFA",
    "lightslategray": "778899",
    "lightslategrey": "778899",
    "lightsteelblue": "B0C4DE",
    "lightyellow": "FFFFE0",
    "lime": "00FF00",
    "limegreen": "32CD32",
    "linen": "FAF0E6",
    "magenta": "FF00FF",
    "maroon": "800000",
    "mediumaquamarine": "66CDAA",
    "mediumblue": "0000CD",
    "mediumorchid": "BA55D3",
    "mediumpurple": "9370DB",
    "mediumseagreen": "3CB371",
    "mediumslateblue": "7B68EE",
    "mediumspringgreen": "00FA9A",
    "mediumturquoise": "48D1CC",
    "mediumvioletred": "C71585",
    "midnightblue": "191970",
    "mintcream": "F5FFFA",
    "mistyrose": "FFE4E1",
    "moccasin": "FFE4B5",
    "navajowhite": "FFDEAD",
    "navy": "000080",
    "oldlace": "FDF5E6",
    "olive": "808000",
    "olivedrab": "6B8E23",
    "orange": "FFA500",
    "orangered": "FF4500",
    "orchid": "DA70D6",
    "palegoldenrod": "EEE8AA",
    "palegreen": "98FB98",
    "paleturquoise": "AFEEEE",
    "palevioletred": "DB7093",
    "papayawhip": "FFEFD5",
    "peachpuff": "FFDAB9",
    "peru": "CD853F",
    "pink": "FFC0CB",
    "plum": "DDA0DD",
    "powderblue": "B0E0E6",
    "purple": "800080",
    "rebeccapurple": "663399",
    "red": "FF0000",
    "rosybrown": "BC8F8F",
    "royalblue": "4169E1",
    "saddlebrown": "8B4513",
    "salmon": "FA8072",
    "sandybrown": "F4A460",
    "seagreen": "2E8B57",
    "seashell": "FFF5EE",
    "sienna": "A0522D",
    "silver": "C0C0C0",
    "skyblue": "87CEEB",
    "slateblue": "6A5ACD",
    "slategray": "708090",
    "slategrey": "708090",
    "snow": "FFFAFA",
    "springgreen": "00FF7F",
    "steelblue": "4682B4",
    "tan": "D2B48C",
    "teal": "008080",
    "thistle": "D8BFD8",
    "tomato": "FF6347",
    "turquoise": "40E0D0",
    "violet": "EE82EE",
    "wheat": "F5DEB3",
    "white": "FFFFFF",
    "whitesmoke": "F5F5F5",
    "yellow": "FFFF00",
    "yellowgreen": "9ACD32",
}


def _resolve_svg_color(raw: str | None, C: dict | None, fallback: str) -> str | None:  # noqa: N803
    """Resolve an SVG fill/stroke color value against a C (context) dict.

    Priority:
      1. var(--name) → look up C[name], then C[f"palette_{name}"]
      2. known C key
      3. hex
      4. rgb()/rgba()
      5. hsl()/hsla()
      6. named-color
      7. currentColor → C.get("text_dark", "#000000")

    Returns None for empty / "none" input (callers rely on None being falsy).
    """
    if not raw:
        return None
    v = raw.strip()
    if not v or v == "none":
        return None

    if v.startswith("var("):
        name = v[4:-1].strip().lstrip("-")
        if name in (C or {}):
            return C[name]
        pk = f"palette_{name}"
        if pk in (C or {}):
            return C[pk]
        raise SVGCompileError(f"unresolved color token: var({name})")

    if v in (C or {}):
        return C[v]

    if v == "currentColor":
        return C.get("text_dark", "#000000") if C else "#000000"

    if v.startswith("#"):
        h = v.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return "#" + h

    lower = v.lower()
    if lower in _NAMED_COLORS:
        return "#" + _NAMED_COLORS[lower]

    # rgb(…) / rgba(…)
    m = _RGB_RE.match(v)
    if m:
        r, g, b, a = int(m.group(1)), int(m.group(2)), int(m.group(3)), m.group(4)
        if a is not None:
            return f"#{r:02X}{g:02X}{b:02X}{int(a * 255):02X}"
        return f"#{r:02X}{g:02X}{b:02X}"

    # hsl(…) / hsla(…)
    m = _HSL_RE.match(v)
    if m:
        import colorsys

        h_val, s_val, l_val = float(m.group(1)) / 360.0, float(m.group(2)) / 100.0, float(m.group(3)) / 100.0
        r, g, b = colorsys.hls_to_rgb(h_val, l_val, s_val)
        a = float(m.group(4)) if m.group(4) else 1.0
        if a < 1.0:
            return f"#{int(r * 255):02X}{int(g * 255):02X}{int(b * 255):02X}{int(a * 255):02X}"
        return f"#{int(r * 255):02X}{int(g * 255):02X}{int(b * 255):02X}"

    raise SVGCompileError(f"unsupported color value: {v!r}")


_RGB_RE = re.compile(r"rgba?\(\s*(\d{1,3})\s*,\s*(\d{1,3})\s*,\s*(\d{1,3})\s*(?:,\s*([\d.]+)\s*)?\)")
_HSL_RE = re.compile(
    r"hsla?\(\s*(\d+(?:\.\d+)?)\s*,\s*(\d+(?:\.\d+)?)%\s*,\s*(\d+(?:\.\d+)?)%\s*(?:,\s*([\d.]+)\s*)?\)"
)


# ─────────────────────────── data classes ───────────────────────────


@dataclass
class SVGResult:
    """Structured result of an SVG compilation.

    ``SVGResult`` remains the public return type for compatibility, while its
    report fields make a compilation diagnosable without inspecting the slide.
    """

    shapes: list = field(default_factory=list)
    native_shapes: list = field(default_factory=list)
    fallback_shapes: list = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)
    features: set = field(default_factory=set)
    feature_levels: dict[str, str] = field(default_factory=dict)
    ir_document: SVGIRDocument | None = None
    source_to_output: dict[str, list] = field(default_factory=dict)
    metrics: dict[str, float | int] = field(default_factory=dict)
    compile_ms: float = 0.0
    shape_count: int = 0


# Descriptive alias for callers that use the result as a compilation report.
SVGRenderReport = SVGResult


_DEFAULT_LIMITS: dict[str, int] = {
    "max_svg_bytes": 2_000_000,
    "max_nodes": 10_000,
    "max_path_commands": 100_000,
    "max_tree_depth": 200,
}


# ─────────────────────────── compiler ───────────────────────────


class SVGCompiler:
    """Compile a subset of SVG to native editable PPTX shapes."""

    def __init__(
        self,
        C: dict | None = None,  # noqa: N803
        limits: dict[str, int] | None = None,
        text_style: dict[str, dict] | None = None,
        group_opacity: str = "strict",
        layout: dict | None = None,
    ) -> None:
        self.C = C or {}
        self.limits = {**_DEFAULT_LIMITS, **(limits or {})}
        self.text_style = text_style or {}
        if group_opacity not in {"strict", "distribute"}:
            raise ValueError("group_opacity must be 'strict' or 'distribute'")
        self.group_opacity = group_opacity
        self.layout = layout or {}
        if self.layout.get("text_collision", "warning") not in {"warning", "error"}:
            raise ValueError("layout.text_collision must be 'warning' or 'error'")

    def compile(
        self,
        svg_text: str,
        slide,
        rect: tuple[float, float, float, float],
        vb: tuple[float, float, float, float] | None = None,
        scaling: str = "contain",
    ) -> SVGResult:
        result = SVGResult()
        t0 = time.perf_counter()

        raw_bytes = len(svg_text.encode("utf-8"))
        self._check_limit("max_svg_bytes", raw_bytes, "SVG document size")

        sanitize_t0 = time.perf_counter()
        root = sanitize(svg_text)
        if self.group_opacity == "distribute":
            self._distribute_group_opacity(root)
        sanitize_ms = (time.perf_counter() - sanitize_t0) * 1000
        node_count = sum(1 for _ in root.iter() if isinstance(_.tag, str))
        self._check_limit("max_nodes", node_count, "SVG node count")
        tree_depth = self._tree_depth(root)
        self._check_limit("max_tree_depth", tree_depth, "SVG tree depth")
        path_command_count = self._count_path_commands(root)
        self._check_limit("max_path_commands", path_command_count, "SVG path command count")
        ir_t0 = time.perf_counter()
        ir_document = build_svg_ir(root)
        ir_build_ms = (time.perf_counter() - ir_t0) * 1000
        self._reject_unrenderable_features(ir_document)

        if vb is None:
            vb_el = root.get("viewBox")
            if vb_el:
                vals = [float(v) for v in vb_el.replace(",", " ").split()]
                if len(vals) >= 4:
                    vb = (vals[0], vals[1], vals[2], vals[3])
            if vb is None:
                vb = (0.0, 0.0, 400.0, 300.0)

        self._vb = vb
        self._rect = rect
        self._scaling = scaling if scaling in ("contain", "cover", "stretch") else "contain"
        self._slide = slide
        self._grads: dict[str, GradientDef] = {}
        self._clips: dict[str, list[list[tuple[float, float]]]] = {}
        self._defs: dict[str, etree._Element] = {}
        self._shape_count = 0
        self._features: set[str] = set()
        self._warnings: list[str] = []
        self._text_rects: list[tuple[float, float, float, float]] = []
        self._source_to_output: dict[str, list] = {}
        if self.group_opacity == "distribute":
            self._warnings.append(
                "group opacity distributed to child elements; overlapping descendants may differ from SVG compositing"
            )
        self._validate_layout(root)

        pre_shape_count = len(slide.shapes)

        defs_t0 = time.perf_counter()
        self._collect_defs(root)
        defs_ms = (time.perf_counter() - defs_t0) * 1000
        render_t0 = time.perf_counter()
        self._walk(root, Affine(), [])
        render_ms = (time.perf_counter() - render_t0) * 1000

        result.features = self._features | set(ir_document.features)
        result.warnings = self._warnings
        result.shape_count = self._shape_count

        # Collect shapes created during compilation
        result.shapes = [slide.shapes[i] for i in range(pre_shape_count, len(slide.shapes))]
        result.native_shapes = list(result.shapes)
        result.source_to_output = self._source_to_output
        result.feature_levels = self._feature_levels(result.features)
        result.ir_document = ir_document

        self._detect_text_overlaps(slide, pre_shape_count, result)
        result.compile_ms = (time.perf_counter() - t0) * 1000
        result.metrics = {
            "svg_bytes": raw_bytes,
            "node_count": node_count,
            "tree_depth": tree_depth,
            "path_command_count": path_command_count,
            "sanitize_ms": sanitize_ms,
            "ir_build_ms": ir_build_ms,
            "ir_node_count": len(ir_document.nodes),
            "defs_ms": defs_ms,
            "render_ms": render_ms,
            "total_ms": result.compile_ms,
            "native_shape_count": len(result.native_shapes),
            "fallback_shape_count": 0,
        }

        return result

    def _check_limit(self, name: str, value: int, label: str) -> None:
        limit = self.limits.get(name)
        if limit is not None and value > limit:
            raise SVGCompileError(f"{label} exceeds {name} limit ({value} > {limit})")

    @staticmethod
    def _count_path_commands(root: etree._Element) -> int:
        count = 0
        for el in root.iter():
            if not isinstance(el.tag, str) or el.tag.split("}")[-1] != "path":
                continue
            count += len(re.findall(r"[MmLlHhVvCcSsQqTtAaZz]", el.get("d", "")))
        return count

    @staticmethod
    def _tree_depth(root: etree._Element) -> int:
        max_depth = 0
        stack = [(root, 1)]
        while stack:
            element, depth = stack.pop()
            if not isinstance(element.tag, str):
                continue
            max_depth = max(max_depth, depth)
            stack.extend((child, depth + 1) for child in element)
        return max_depth

    @staticmethod
    def _reject_unrenderable_features(ir_document: SVGIRDocument) -> None:
        """Fail safely when native OOXML cannot preserve a group compositing rule."""
        for node in ir_document.nodes:
            if "group_opacity" in node.features and "hidden" not in node.features:
                label = f" id='{node.source_id}'" if node.source_id else ""
                raise SVGCompileError(
                    f"group opacity on <{node.tag}{label}> requires raster fallback; "
                    "native rendering would be visually incorrect"
                )

    @staticmethod
    def _distribute_group_opacity(root: etree._Element) -> None:
        """Push group opacity to descendants for native editable rendering."""
        def parse_opacity(raw: str | None) -> float:
            value = (raw or "1").strip()
            if value.endswith("%"):
                return float(value[:-1]) / 100.0
            return float(value)

        def walk(el: etree._Element, inherited: float = 1.0) -> None:
            own = parse_opacity(el.get("opacity"))
            effective = max(0.0, min(1.0, inherited * own))
            tag = el.tag.split("}")[-1] if isinstance(el.tag, str) else ""
            if tag in {"g", "svg"}:
                if own != 1.0:
                    el.set("data-group-opacity-distributed", str(own))
                el.set("opacity", "1")
                for child in el:
                    walk(child, effective)
            else:
                if effective != 1.0:
                    el.set("opacity", f"{effective:.6g}")
                for child in el:
                    walk(child, effective)

        walk(root)

    def _validate_layout(self, root: etree._Element) -> None:
        """Validate declared text zones in root viewBox coordinates."""
        zones = self.layout.get("zones", {})
        margin = float(self.layout.get("safe_margin", 0))
        if not zones and margin <= 0:
            return
        x0, y0, vb_w, vb_h = self._vb
        violations: list[str] = []
        def walk(el: etree._Element, tf: Affine) -> None:
            if not isinstance(el.tag, str):
                return
            tf = tf.compose(parse_transform(el.get("transform")))
            if el.tag.split("}")[-1] == "text":
                try:
                    local_x = float((el.get("x") or "0").split(",")[0])
                    local_y = float((el.get("y") or "0").split(",")[0])
                except ValueError:
                    local_x = local_y = None
                if local_x is not None and local_y is not None:
                    x, y = tf.apply(local_x, local_y)
                    role = el.get("class")
                    label = f"class {role!r}" if role else "unclassified text"
                    if margin and (
                        x < x0 + margin
                        or y < y0 + margin
                        or x > x0 + vb_w - margin
                        or y > y0 + vb_h - margin
                    ):
                        violations.append(f"{label} is inside the {margin:g}px unsafe SVG margin")
                    if role in zones:
                        zx, zy, zw, zh = zones[role]
                        if not (zx <= x <= zx + zw and zy <= y <= zy + zh):
                            violations.append(f"{label} anchor ({x:g},{y:g}) is outside zone {role!r}")
            for child in el:
                walk(child, tf)

        walk(root, Affine())
        if not violations:
            return
        mode = self.layout.get("text_collision", "warning")
        if mode == "error":
            raise SVGCompileError("SVG layout contract failed: " + "; ".join(violations))
        self._warnings.extend(f"layout: {item}" for item in violations)

    def _feature_levels(self, features: set[str]) -> dict[str, str]:
        levels: dict[str, str] = {}
        for feature in features:
            if feature in {"gradient"}:
                levels[feature] = "OOXML_EFFECT"
            elif feature in {"clipPath", "evenodd"}:
                levels[feature] = "NATIVE_APPROX"
            elif feature in {
                "image",
                "filter",
                "mask",
                "pattern",
                "marker",
                "group_opacity",
                "raster_fallback_candidate",
            }:
                levels[feature] = "RASTER_FALLBACK_CANDIDATE"
            else:
                levels[feature] = "NATIVE"
        return levels

    # ── coordinate transforms ────────────────────────────────

    def _to_inches(self, x: float, y: float) -> tuple[float, float]:
        lx, ly, w, h = self._rect
        vx, vy, vw, vh = self._vb
        if self._scaling == "stretch":
            sx = w / vw if vw > 0 else 1.0
            sy = h / vh if vh > 0 else 1.0
            return lx + (x - vx) * sx, ly + (y - vy) * sy
        if self._scaling == "cover":
            s = max(w / vw, h / vh) if vw > 0 and vh > 0 else 1.0
        else:
            s = min(w / vw, h / vh) if vw > 0 and vh > 0 else 1.0
        ox = lx + (w - vw * s) / 2.0
        oy = ly + (h - vh * s) / 2.0
        return ox + (x - vx) * s, oy + (y - vy) * s

    # ── defs collection ──────────────────────────────────────

    def _collect_defs(self, root: etree._Element) -> None:
        for g in root.iter(f"{SVG}linearGradient"):
            stops = []
            for s in g.iter(f"{SVG}stop"):
                off = s.get("offset", "0")
                pos = float(off.rstrip("%")) / 100.0 if off.endswith("%") else float(off)
                # Get color from attribute or style
                col = s.get("stop-color")
                op_str = s.get("stop-opacity")
                if col is None:
                    # Parse from style attribute
                    style = s.get("style", "")
                    for prop in style.split(";"):
                        prop = prop.strip()
                        if prop.startswith("stop-color:"):
                            col = prop.split(":", 1)[1].strip()
                        elif prop.startswith("stop-opacity:"):
                            op_str = prop.split(":", 1)[1].strip()
                if col is None:
                    col = "#000000"
                col = _resolve_svg_color(col, self.C, "#000000")
                op = float(op_str) if op_str else 1.0
                stops.append((pos, col, op))
            x1 = _parse_pct(g.get("x1", "0"))
            y1 = _parse_pct(g.get("y1", "0"))
            x2 = _parse_pct(g.get("x2", "1"))
            y2 = _parse_pct(g.get("y2", "0"))
            self._grads[g.get("id")] = GradientDef(stops=stops, x1=x1, y1=y1, x2=x2, y2=y2)

        for g in root.iter(f"{SVG}radialGradient"):
            stops = []
            for s in g.iter(f"{SVG}stop"):
                off = s.get("offset", "0")
                pos = float(off.rstrip("%")) / 100.0 if off.endswith("%") else float(off)
                # Get color from attribute or style
                col = s.get("stop-color")
                op_str = s.get("stop-opacity")
                if col is None:
                    style = s.get("style", "")
                    for prop in style.split(";"):
                        prop = prop.strip()
                        if prop.startswith("stop-color:"):
                            col = prop.split(":", 1)[1].strip()
                        elif prop.startswith("stop-opacity:"):
                            op_str = prop.split(":", 1)[1].strip()
                if col is None:
                    col = "#000000"
                col = _resolve_svg_color(col, self.C, "#000000")
                op = float(op_str) if op_str else 1.0
                stops.append((pos, col, op))
            cx = _parse_pct(g.get("cx", "50%"))
            cy = _parse_pct(g.get("cy", "50%"))
            r_ = _parse_pct(g.get("r", "50%"))
            self._grads[g.get("id")] = GradientDef(stops=stops, cx=cx, cy=cy, r=r_, gradient_type="radial")

        for c in root.iter(f"{SVG}clipPath"):
            polys = []
            for child in c:
                tf = Affine()
                sub = self._svg_polygon(child, tf)
                if sub:
                    for path_cmds in sub:
                        flat: list[tuple[float, float]] = []
                        for cmd in path_cmds:
                            if isinstance(cmd, _LineCmd):
                                flat.append((cmd.x, cmd.y))
                            else:
                                flat.append((cmd.x1, cmd.y1))
                                flat.append((cmd.x2, cmd.y2))
                                flat.append((cmd.x3, cmd.y3))
                        polys.append(flat)
            self._clips[c.get("id")] = polys

        for el in root.iter():
            if not isinstance(el.tag, str):
                continue
            el_id = el.get("id")
            if (
                el_id
                and el.tag.split("}")[-1]
                in (
                    "symbol",
                    "g",
                    "rect",
                    "circle",
                    "ellipse",
                    "path",
                    "polygon",
                    "polyline",
                    "line",
                )
                and el_id not in self._defs
            ):
                self._defs[el_id] = el

    # ── SVG element → polygon points (SVG-space) ─────────────

    def _svg_polygon(self, el: etree._Element, tf: Affine) -> list[list[_PathCmd]]:
        if not isinstance(el.tag, str):
            return []
        tag = el.tag.split("}")[-1]
        if tag == "rect":
            x = float(el.get("x", 0))
            y = float(el.get("y", 0))
            w = float(el.get("width", 0))
            h = float(el.get("height", 0))
            rx = float(el.get("rx", 0))
            ry = float(el.get("ry", 0))
            if rx > 0 or ry > 0:
                rx = min(rx, w / 2) if rx > 0 else max(0, ry)
                ry = min(ry, h / 2) if ry > 0 else max(0, rx)
                return [self._cubics_to_path(tf, self._rounded_rect_cubics(x, y, w, h, rx, ry))]
            pts = [(x, y), (x + w, y), (x + w, y + h), (x, y + h)]
        elif tag == "circle":
            cx = float(el.get("cx", 0))
            cy = float(el.get("cy", 0))
            r = float(el.get("r", 0))
            return [self._cubics_to_path(tf, self._circle_cubics(cx, cy, r))]
        elif tag == "ellipse":
            cx = float(el.get("cx", 0))
            cy = float(el.get("cy", 0))
            rx = float(el.get("rx", 0))
            ry = float(el.get("ry", 0))
            return [self._cubics_to_path(tf, self._ellipse_cubics(cx, cy, rx, ry))]
        elif tag in ("polygon", "polyline"):
            flat = [float(v) for v in el.get("points", "").replace(",", " ").split()]
            pts = list(zip(flat[0::2], flat[1::2], strict=False))
        elif tag == "line":
            pts = [
                (float(el.get("x1")), float(el.get("y1"))),
                (float(el.get("x2")), float(el.get("y2"))),
            ]
        elif tag == "path":
            return self._path_to_cmds(el, tf)
        else:
            return []
        path: list[_PathCmd] = []
        for px, py in pts:
            tx, ty = tf.apply(px, py)
            path.append(_LineCmd(tx, ty))
        return [path]

    def _cubics_to_path(self, tf: Affine, pts: list[tuple[float, float]]) -> list[_PathCmd]:
        """Convert (start, c1, c2, end)×N flat list to _PathCmd list."""
        path: list[_PathCmd] = []
        i = 0
        while i < len(pts):
            if i % 4 == 0 and i + 3 < len(pts):
                sx, sy = tf.apply(pts[i][0], pts[i][1])
                c1x, c1y = tf.apply(pts[i + 1][0], pts[i + 1][1])
                c2x, c2y = tf.apply(pts[i + 2][0], pts[i + 2][1])
                ex, ey = tf.apply(pts[i + 3][0], pts[i + 3][1])
                if not path:
                    path.append(_LineCmd(sx, sy))
                path.append(_CubicCmd(c1x, c1y, c2x, c2y, ex, ey))
                i += 4
            else:
                tx, ty = tf.apply(pts[i][0], pts[i][1])
                path.append(_LineCmd(tx, ty))
                i += 1
        return path

    def _path_to_cmds(self, el: etree._Element, tf: Affine) -> list[list[_PathCmd]]:
        cmds, _ = parse_path(el.get("d", ""))
        subs = to_beziers(cmds)
        out: list[list[_PathCmd]] = []
        for sub in subs:
            if not sub:
                continue
            start = sub[0][0]
            sx, sy = tf.apply(start[0], start[1])
            path: list[_PathCmd] = [_LineCmd(sx, sy)]
            for seg in sub:
                (_x0, _y0), (x1, y1), (x2, y2), (x3, y3) = seg
                c1x, c1y = tf.apply(x1, y1)
                c2x, c2y = tf.apply(x2, y2)
                ex, ey = tf.apply(x3, y3)
                path.append(_CubicCmd(c1x, c1y, c2x, c2y, ex, ey))
            out.append(path)
        return out

    @staticmethod
    def _rounded_rect_cubics(
        x: float,
        y: float,
        w: float,
        h: float,
        rx: float,
        ry: float,
    ) -> list[tuple[float, float]]:
        """Return cubic Bezier points for a rounded rectangle.

        4 corners × 4 points (start, c1, c2, end) = 16 points,
        plus 4 straight-edge start points = 20 points total.
        """
        k = (4 / 3) * (math.sqrt(2) - 1)
        pts: list[tuple[float, float]] = []

        # Top-left corner → top-right corner
        pts.append((x + rx, y))
        pts.extend(
            [
                (x + rx - k * rx, y),
                (x, y + ry - k * ry),
                (x, y + ry),
            ]
        )
        # Top-right corner → bottom-right corner
        pts.append((x, y + h - ry))
        pts.extend(
            [
                (x, y + h - ry + k * ry),
                (x + rx - k * rx, y + h),
                (x + rx, y + h),
            ]
        )
        # Bottom-right corner → bottom-left corner
        pts.append((x + w - rx, y + h))
        pts.extend(
            [
                (x + w - rx + k * rx, y + h),
                (x + w, y + h - ry + k * ry),
                (x + w, y + h - ry),
            ]
        )
        # Bottom-left corner → top-left corner
        pts.append((x + w, y + ry))
        pts.extend(
            [
                (x + w, y + ry - k * ry),
                (x + w - rx + k * rx, y),
                (x + w - rx, y),
            ]
        )
        return pts

    @staticmethod
    def _circle_cubics(cx: float, cy: float, r: float) -> list[tuple[float, float]]:
        """Return 4 cubic Bezier control points for a circle (16 points total).

        Uses standard circle approximation with k = 4/3 * tan(pi/8).
        Each quadrant is approximated by one cubic Bezier curve.
        """
        pts: list[tuple[float, float]] = []
        k = (4 / 3) * math.tan(math.pi / 8)  # ≈ 0.5523

        # Start at angle 0 (rightmost point) and go counter-clockwise
        for i in range(4):
            a0 = i * math.pi / 2
            a1 = (i + 1) * math.pi / 2

            # Start and end points
            p0x = cx + r * math.cos(a0)
            p0y = cy - r * math.sin(a0)  # Note: SVG y-axis is inverted
            p3x = cx + r * math.cos(a1)
            p3y = cy - r * math.sin(a1)

            # Control points (tangent to circle at start/end)
            c1x = p0x - k * r * math.sin(a0)
            c1y = p0y - k * r * math.cos(a0)
            c2x = p3x + k * r * math.sin(a1)
            c2y = p3y + k * r * math.cos(a1)

            pts.extend([(p0x, p0y), (c1x, c1y), (c2x, c2y), (p3x, p3y)])

        return pts

    @staticmethod
    def _ellipse_cubics(cx: float, cy: float, rx: float, ry: float) -> list[tuple[float, float]]:
        """Return 4 cubic Bezier control points for an ellipse (16 points total).

        Uses standard ellipse approximation.
        Note: SVG y-axis is inverted (down is positive), so we use cy - ry*sin(a).
        """
        pts: list[tuple[float, float]] = []
        k = (4 / 3) * (math.sqrt(2) - 1)  # ≈ 0.5523

        for i in range(4):
            a0 = i * math.pi / 2
            a1 = (i + 1) * math.pi / 2

            # Start and end points (SVG y-axis inverted)
            p0x = cx + rx * math.cos(a0)
            p0y = cy - ry * math.sin(a0)
            p3x = cx + rx * math.cos(a1)
            p3y = cy - ry * math.sin(a1)

            # Control points (tangent to ellipse at start/end)
            c1x = p0x - k * rx * math.sin(a0)
            c1y = p0y - k * ry * math.cos(a0)
            c2x = p3x + k * rx * math.sin(a1)
            c2y = p3y + k * ry * math.cos(a1)

            pts.extend([(p0x, p0y), (c1x, c1y), (c2x, c2y), (p3x, p3y)])

        return pts

    @staticmethod
    def _flatten_cubic(
        seg: tuple[tuple[float, float], tuple[float, float], tuple[float, float], tuple[float, float]],
        n: int = 6,
    ) -> list[tuple[float, float]]:
        (x0, y0), (x1, y1), (x2, y2), (x3, y3) = seg
        out: list[tuple[float, float]] = []
        for i in range(n + 1):
            t = i / n
            mt = 1.0 - t
            px = mt**3 * x0 + 3 * mt * mt * t * x1 + 3 * mt * t * t * x2 + t**3 * x3
            py = mt**3 * y0 + 3 * mt * mt * t * y1 + 3 * mt * t * t * y2 + t**3 * y3
            out.append((px, py))
        return out

    # ── paint resolution ─────────────────────────────────────

    def _paint(self, el: etree._Element, which: str) -> tuple[str, object | None, int]:
        return _resolve_paint(el, which, self._grads, self.C, _resolve_svg_color, self._features)

    # ── rendering ────────────────────────────────────────────

    def _walk(self, el: etree._Element, tf: Affine, clip_stack: list) -> None:
        # Skip non-element nodes (e.g., lxml Comment, ProcessingInstruction)
        if not isinstance(el.tag, str):
            return
        if el.get("display") == "none" or el.get("visibility") in {"hidden", "collapse"}:
            return
        pre_node_shape_count = len(self._slide.shapes)
        tag = el.tag.split("}")[-1] if el.tag else ""
        tf = tf.compose(parse_transform(el.get("transform")))

        c = el.get("clip-path")
        if c and c.startswith("url(#"):
            gid = c[c.index("#") + 1 : -1]
            polys = self._clips.get(gid, [])
            clip_stack = clip_stack + polys

        if tag == "g" or tag == "svg":
            for child in el:
                self._walk(child, tf, clip_stack)
        elif tag == "defs":
            pass  # already collected
        elif tag == "use":
            self._render_use(el, tf, clip_stack)
        elif tag in ("image", "filter", "mask"):
            self._features.add(tag)
            self._warnings.append(f"unsupported SVG feature: <{tag}> element (skipped)")
        elif tag == "text":
            self._render_text(el, tf)
        elif tag in (
            "rect",
            "circle",
            "ellipse",
            "polygon",
            "polyline",
            "line",
            "path",
        ):
            self._render_shape(el, tag, tf, clip_stack)

        source_id = el.get("id")
        if source_id:
            self._record_source_output(source_id, pre_node_shape_count)

    def _record_source_output(self, source_id: str | None, start_index: int) -> None:
        if not source_id:
            return
        new_shapes = [self._slide.shapes[i] for i in range(start_index, len(self._slide.shapes))]
        if new_shapes:
            self._source_to_output.setdefault(source_id, []).extend(new_shapes)

    def _render_use(self, el: etree._Element, tf: Affine, clip_stack: list) -> None:
        pre_use_shape_count = len(self._slide.shapes)
        self._features.add("use")
        href = el.get("href") or el.get(f"{SVG}href") or el.get("{http://www.w3.org/1999/xlink}href") or ""
        href = href.removeprefix("#")
        ref = self._defs.get(href)
        if ref is None:
            self._warnings.append(f"<use> references unknown id '{href}' (skipped)")
            return

        ux = float(el.get("x", 0))
        uy = float(el.get("y", 0))
        use_tf = tf.compose(Affine(1, 0, 0, 1, ux, uy))

        ref_tag = ref.tag.split("}")[-1] if isinstance(ref.tag, str) else ""
        if ref_tag in ("g", "symbol", "svg"):
            for child in ref:
                self._walk(child, use_tf, clip_stack)
        elif ref_tag == "text":
            self._render_text(ref, use_tf)
        elif ref_tag in (
            "rect",
            "circle",
            "ellipse",
            "polygon",
            "polyline",
            "line",
            "path",
        ):
            self._render_shape(ref, ref_tag, use_tf, clip_stack)
        else:
            self._warnings.append(f"<use> references unsupported element <{ref_tag}> (skipped)")
        self._record_source_output(ref.get("id"), pre_use_shape_count)

    def _render_shape(self, el: etree._Element, tag: str, tf: Affine, clip_stack: list) -> None:
        self._features.add(tag)

        fkind, fval, fa = self._paint(el, "fill")
        skind, sval, _ = self._paint(el, "stroke")
        sw = float(el.get("stroke-width", "1"))
        stroke_style = parse_stroke_style(el)

        subpaths = self._svg_polygon(el, tf)
        if not subpaths:
            return

        all_pts: list[tuple[float, float]] = []
        for sub in subpaths:
            for cmd in sub:
                if isinstance(cmd, _LineCmd):
                    all_pts.append((cmd.x, cmd.y))
                else:
                    all_pts.append((cmd.x1, cmd.y1))
                    all_pts.append((cmd.x2, cmd.y2))
                    all_pts.append((cmd.x3, cmd.y3))
        if not all_pts:
            return

        minx = min(p[0] for p in all_pts)
        maxx = max(p[0] for p in all_pts)
        miny = min(p[1] for p in all_pts)
        maxy = max(p[1] for p in all_pts)

        # map bounding box to slide inches
        corners_in = [self._to_inches(minx, miny), self._to_inches(maxx, maxy)]
        ix0, iy0 = corners_in[0]
        ix1, iy1 = corners_in[1]
        iw = ix1 - ix0
        ih = iy1 - iy0

        # rebuild subpaths in slide inches, local coords
        local: list[list[_PathCmd]] = []
        for sub in subpaths:
            local_sub: list[_PathCmd] = []
            for cmd in sub:
                if isinstance(cmd, _LineCmd):
                    px, py = self._to_inches(cmd.x, cmd.y)
                    local_sub.append(_LineCmd(px - ix0, py - iy0))
                else:
                    p1x, p1y = self._to_inches(cmd.x1, cmd.y1)
                    p2x, p2y = self._to_inches(cmd.x2, cmd.y2)
                    p3x, p3y = self._to_inches(cmd.x3, cmd.y3)
                    local_sub.append(_CubicCmd(p1x - ix0, p1y - iy0, p2x - ix0, p2y - iy0, p3x - ix0, p3y - iy0))
            local.append(local_sub)

        # clip polygons in slide inches, local coords (relative to shape bbox)
        clip_local: list[list[tuple[float, float]]] = []
        for clip_subpath in clip_stack:
            cpts = [self._to_inches(px, py) for px, py in clip_subpath]
            clip_local.append([(px - ix0, py - iy0) for px, py in cpts])

        needs_bool = bool(clip_stack) or el.get("fill-rule") == "evenodd"

        if needs_bool:
            flat_local: list[list[tuple[float, float]]] = []
            for sub in local:
                flat: list[tuple[float, float]] = []
                for cmd in sub:
                    if isinstance(cmd, _LineCmd):
                        flat.append((cmd.x, cmd.y))
                    else:
                        flat.extend(
                            self._flatten_cubic(((0, 0), (cmd.x1, cmd.y1), (cmd.x2, cmd.y2), (cmd.x3, cmd.y3)), n=6)[1:]
                        )
                flat_local.append(flat)
            self._render_bool(
                flat_local,
                clip_local,
                ix0,
                iy0,
                iw,
                ih,
                fkind,
                fval,
                fa,
                skind,
                sval,
                stroke_style,
            )
            return

        # Fast path: rect with solid or gradient fill and no stroke → native shape
        if (
            tag == "rect"
            and not el.get("rx")
            and not el.get("ry")
            and fkind in ("solid", "grad")
            and skind == "none"
            and not clip_stack
        ):
            if fkind == "solid":
                self._add_native(ix0, iy0, iw, ih, fval, fa)
            else:
                # Gradient on native shape
                elem = self._add_native(ix0, iy0, iw, ih, None, fa)
                self._apply_gradient_to_elem(elem, fval)
            return

        # Freeform path
        fill_hex = _resolve_svg_color(fval, self.C, "") if fkind == "solid" else None
        stroke_hex = _resolve_svg_color(sval, self.C, "") if skind == "solid" else None
        elem = self._add_freeform(local, ix0, iy0, iw, ih, fill_hex or "#FFFFFF", fa, stroke_hex, sw)
        if stroke_style and elem is not None:
            apply_stroke_style(elem, stroke_style)
        if fkind == "grad":
            self._apply_gradient_to_elem(elem, fval)

    def _render_bool(
        self,
        local_subpaths: list[list[tuple[float, float]]],
        clip_local: list[list[tuple[float, float]]],
        ix0: float,
        iy0: float,
        iw: float,
        ih: float,
        fkind: str,
        fval,
        fa: int,
        skind: str,
        sval,
        stroke_style=None,
    ) -> None:
        """Boolean-compute shapes and render, intersecting with clip regions."""
        try:
            from shapely.errors import GEOSException, TopologicalError
            from shapely.geometry import Polygon as ShapelyPoly
            from shapely.validation import make_valid
        except ImportError:
            raise SVGCompileError("shapely is required for boolean operations") from None

        poly = None
        for sub in local_subpaths:
            pts = [(p[0] + ix0, p[1] + iy0) for p in sub]
            if len(pts) < 3:
                continue
            try:
                p = ShapelyPoly(pts).buffer(0)
                p = make_valid(p)
            except (TopologicalError, ValueError, GEOSException):
                p = None
            if p is not None and not p.is_empty:
                poly = p if poly is None else poly.union(p)

        if poly is None or poly.is_empty:
            return

        # Apply clip regions: intersect shape with clip in slide inches
        for clip_poly in clip_local:
            cpts = [(p[0] + ix0, p[1] + iy0) for p in clip_poly]
            if len(cpts) < 3:
                continue
            try:
                cpoly = ShapelyPoly(cpts).buffer(0)
                cpoly = make_valid(cpoly)
            except (TopologicalError, ValueError, GEOSException):
                cpoly = None
            if cpoly is not None and not cpoly.is_empty:
                poly = poly.intersection(cpoly)

        if poly.is_empty:
            return

        stroke_hex = _resolve_svg_color(sval, self.C, "") if skind == "solid" else None

        if fkind == "grad":
            elem = bool_shape(poly, self._slide, ix0, iy0, iw, ih, fill="#FFFFFF", line=stroke_hex, alpha=fa)
            if elem is not None:
                self._apply_gradient_to_elem(elem, fval)
        else:
            fill_hex = _resolve_svg_color(fval, self.C, "") if fkind == "solid" else "#FFFFFF"
            elem = bool_shape(poly, self._slide, ix0, iy0, iw, ih, fill=fill_hex, line=stroke_hex, alpha=fa)
        if elem is not None:
            if stroke_style:
                apply_stroke_style(elem, stroke_style)
            self._shape_count += 1

    def _add_native(self, x: float, y: float, w: float, h: float, fill: str | None, alpha: int) -> object:
        sh = self._slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip("#"))
            if alpha < 100:
                set_solid_fill_with_alpha(sh, fill, alpha)
        else:
            sh.fill.background()
        sh.line.fill.background()
        self._shape_count += 1
        return sh

    def _add_freeform(
        self,
        local_subs: list[list[_PathCmd]],
        ix0: float,
        iy0: float,
        iw: float,
        ih: float,
        fill: str | None,
        alpha: int,
        stroke: str | None,
        sw: float,
    ) -> object:
        fb = FreeformBuilder()
        for sub in local_subs:
            if not sub:
                continue
            first = sub[0]
            if isinstance(first, _LineCmd):
                fb.move_to(first.x, first.y)
            else:
                fb.move_to(first.x1, first.y1)
            for cmd in sub[1:]:
                if isinstance(cmd, _LineCmd):
                    fb.line_to(cmd.x, cmd.y)
                else:
                    fb.cubic_bezier_to(cmd.x1, cmd.y1, cmd.x2, cmd.y2, cmd.x3, cmd.y3)
            fb.close()
        elem = fb.build(
            self._slide,
            ix0,
            iy0,
            iw,
            ih,
            fill_color=fill,
            line_color=stroke,
            line_width_pt=max(sw, 0.5),
            no_fill=(fill is None),
        )
        if alpha < 100 and fill:
            self._apply_alpha_to_elem(elem, fill, alpha)
        self._shape_count += 1
        return elem

    def _apply_gradient_to_elem(self, elem: object, grad: GradientDef) -> None:
        _apply_gradient(elem, grad, self._wrap_elem)

    def _apply_alpha_to_elem(self, elem: object, fill: str, alpha: int) -> None:
        if alpha >= 100:
            return
        wrapper = self._wrap_elem(elem)
        set_solid_fill_with_alpha(wrapper, fill, alpha)

    @staticmethod
    def _wrap_elem(elem):
        """Wrap a raw lxml element so visual_effects can access shape._element."""
        if hasattr(elem, "_element"):
            return elem

        class _ShapeProxy:
            def __init__(self, el):
                self._element = el

        return _ShapeProxy(elem)

    # ── text ─────────────────────────────────────────────────

    def _render_text(self, el: etree._Element, tf: Affine) -> None:
        # scale for font-size: svg region size / viewBox size
        _lx, _ly, rw, rh = self._rect
        if self.text_style and not el.get("class"):
            self._warnings.append("SVG text has no class; add text_style['role'] with an explicit font_size")
        elif self.text_style and el.get("class") not in self.text_style:
            self._warnings.append(f"SVG text class {el.get('class')!r} has no matching text_style entry")
        svg_w, svg_h = self._vb[2], self._vb[3]
        _render_svg_text(
            el=el,
            tf=tf,
            to_inches_fn=self._to_inches,
            slide=self._slide,
            C=self.C,
            resolve_color_fn=_resolve_svg_color,
            features=self._features,
            svg_w=svg_w,
            svg_h=svg_h,
            slide_w=rw,
            slide_h=rh,
            text_style=self.text_style,
        )
        self._shape_count += 1

    # ── collision detection ──────────────────────────────────

    @staticmethod
    def _detect_text_overlaps(slide, pre_count: int, result: SVGResult) -> None:
        """Detect overlapping text boxes among shapes added during this compile.

        Walks all text boxes added since pre_count, collects their bounding
        rects in slide inches, and emits a warning for each pair that overlaps
        by more than a small tolerance (avoiding false positives from adjacent
        text boxes that touch at edges).
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        text_rects: list[tuple[float, float, float, float]] = []
        for shape in slide.shapes:
            if slide.shapes.index(shape) < pre_count:
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            try:
                x = shape.left / 914400.0
                y = shape.top / 914400.0
                w = shape.width / 914400.0
                h = shape.height / 914400.0
            except (AttributeError, TypeError):
                continue
            text_rects.append((x, y, x + w, y + h))

        n = len(text_rects)
        if n < 2:
            return

        tol = 0.05
        for i in range(n):
            for j in range(i + 1, n):
                ax0, ay0, ax1, ay1 = text_rects[i]
                bx0, by0, bx1, by1 = text_rects[j]
                dx = min(ax1, bx1) - max(ax0, bx0)
                dy = min(ay1, by1) - max(ay0, by0)
                if dx > tol and dy > tol:
                    h1 = ay1 - ay0
                    h2 = by1 - by0
                    # Same-left stacks (title/subtitle/badge) where one box
                    # is vertically offset from the other by >= 0.1" — these
                    # are intentional vertical layouts, not real overlaps.
                    if abs(ax0 - bx0) < 0.4:
                        top_offset = abs(ay0 - by0) if h1 <= h2 else abs(by0 - ay0)
                        if top_offset >= 0.1:
                            continue
                    result.warnings.append(
                        f"text box overlap detected: "
                        f"rect1=({ax0:.2f},{ay0:.2f},{ax1:.2f},{ay1:.2f}) "
                        f"rect2=({bx0:.2f},{by0:.2f},{bx1:.2f},{by1:.2f}) "
                        f"overlap=({dx:.2f}x{dy:.2f})"
                    )
