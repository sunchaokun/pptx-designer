"""Text Effects — gradient fills, outlines, shadows, glow on text runs.

All functions operate on a:rPr (run properties), enabling text-level
visual effects that visual_effects.py only provides at shape level.
"""

from __future__ import annotations

from lxml import etree
from pptx.oxml.ns import qn

TEXT_GRADIENT_PRESETS: dict[str, list[tuple[str, int]]] = {
    "gold-shine": [("F5AF19", 0), ("FFA300", 50000), ("FF6B00", 100000)],
    "blue-deep": [("1D78FA", 0), ("0165FF", 100000)],
    "purple-neon": [("8B5CF6", 0), ("6366F1", 50000), ("3B82F6", 100000)],
    "ink-wash": [("2C2C2C", 0), ("8B7355", 50000), ("D4C5A9", 100000)],
    "cyber-cyan": [("22D3EE", 0), ("06B6D4", 50000), ("0891B2", 100000)],
    "sunset": [("FF6B6B", 0), ("F59E0B", 50000), ("FF5500", 100000)],
    "emerald": [("10B981", 0), ("059669", 50000), ("047857", 100000)],
    "rose-gold": [("B76E79", 0), ("E8B4B8", 50000), ("F5D0D0", 100000)],
    "seal-red": [("C41E3A", 0), ("8B0000", 100000)],
    "steel": [("64748B", 0), ("94A3B8", 50000), ("CBD5E1", 100000)],
}

_VERT_MAP = {
    "ea": "eaVert",
    "mongolian": "mongolianVert",
    "270": "vert270",
}


def _remove_run_fill(rPr: etree._Element) -> None:
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)


def _ensure_effect_list(rPr: etree._Element) -> etree._Element:
    effectLst = rPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = etree.SubElement(rPr, qn("a:effectLst"))
    return effectLst


def apply_text_gradient(run, stops: list[tuple[str, int]], angle: int = 5400000) -> None:
    rPr = run._r.get_or_add_rPr()
    _remove_run_fill(rPr)
    gradFill = etree.SubElement(rPr, qn("a:gradFill"))
    gradFill.set("rotWithShape", "1")
    gradFill.set("flip", "none")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for color_hex, pos in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", color_hex.lstrip("#"))
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", str(angle))
    lin.set("scaled", "1")


def apply_text_gradient_preset(run, preset: str) -> None:
    if preset not in TEXT_GRADIENT_PRESETS:
        raise KeyError(f"Unknown text gradient preset: {preset!r}. Available: {list(TEXT_GRADIENT_PRESETS.keys())}")
    stops = TEXT_GRADIENT_PRESETS[preset]
    apply_text_gradient(run, stops)


def apply_text_outline(run, color: str, width_pt: float) -> None:
    rPr = run._r.get_or_add_rPr()
    existing = rPr.find(qn("a:ln"))
    if existing is not None:
        rPr.remove(existing)
    ln = etree.SubElement(rPr, qn("a:ln"))
    ln.set("w", str(int(width_pt * 12700)))
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", color.lstrip("#"))


def apply_text_shadow(
    run, blur: float = 6.0, dist: float = 3.0, direction: float = 90.0, color: str = "#000000", alpha: int = 25
) -> None:
    rPr = run._r.get_or_add_rPr()
    effectLst = _ensure_effect_list(rPr)
    existing = effectLst.find(qn("a:outerShdw"))
    if existing is not None:
        effectLst.remove(existing)
    shdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    shdw.set("blurRad", str(int(blur * 12700)))
    shdw.set("dist", str(int(dist * 12700)))
    shdw.set("dir", str(int(direction * 60000)))
    shdw.set("algn", "tl")
    shdw.set("rotWithShape", "0")
    srgb = etree.SubElement(shdw, qn("a:srgbClr"))
    srgb.set("val", color.lstrip("#"))
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(alpha * 1000))


def apply_text_glow(run, radius: float = 8.0, color: str = "#2563EB", alpha: int = 40) -> None:
    rPr = run._r.get_or_add_rPr()
    effectLst = _ensure_effect_list(rPr)
    existing = effectLst.find(qn("a:glow"))
    if existing is not None:
        effectLst.remove(existing)
    glow = etree.SubElement(effectLst, qn("a:glow"))
    glow.set("rad", str(int(radius * 12700)))
    srgb = etree.SubElement(glow, qn("a:srgbClr"))
    srgb.set("val", color.lstrip("#"))
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(alpha * 1000))


def apply_text_alpha(run, alpha_pct: int) -> None:
    rPr = run._r.get_or_add_rPr()
    _remove_run_fill(rPr)
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", "000000")
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(alpha_pct * 1000))


def apply_text_3d(run, depth_pt: float = 10.0, bevel: bool = True, material: str = "powder") -> None:
    rPr = run._r.get_or_add_rPr()
    existing = rPr.find(qn("a:sp3d"))
    if existing is not None:
        rPr.remove(existing)
    existing_scene = rPr.find(qn("a:scene3d"))
    if existing_scene is not None:
        rPr.remove(existing_scene)
    sp3d = etree.SubElement(rPr, qn("a:sp3d"))
    sp3d.set("z", str(int(depth_pt * 12700)))
    if bevel:
        bevelT = etree.SubElement(sp3d, qn("a:bevelT"))
        bevelT.set("w", str(int(4.0 * 12700)))
        bevelT.set("h", str(int(2.0 * 12700)))
    prstMat = etree.SubElement(sp3d, qn("a:prstMaterial"))
    prstMat.set("val", material)
    scene3d = etree.SubElement(rPr, qn("a:scene3d"))
    camera = etree.SubElement(scene3d, qn("a:camera"))
    camera.set("prst", "perspectiveFront")
    lightRig = etree.SubElement(scene3d, qn("a:lightRig"))
    lightRig.set("rig", "threePt")
    lightRig.set("dir", "t")


def set_vertical_text(text_frame, direction: str = "ea") -> None:
    vert_val = _VERT_MAP.get(direction, "eaVert")
    bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(text_frame._txBody, qn("a:bodyPr"))
    bodyPr.set("vert", vert_val)


def set_text_rotation(shape, degrees: float) -> None:
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is not None:
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is not None:
            xfrm.set("rot", str(int(degrees * 60000)))
            return
    xfrm = shape._element.find(f".//{qn('a:xfrm')}")
    if xfrm is not None:
        xfrm.set("rot", str(int(degrees * 60000)))


def apply_letter_spacing(run, tracking_em: float, font_size_pt: int) -> None:
    if tracking_em == 0.0:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is not None and rPr.get("spc") is not None:
            rPr.attrib.pop("spc", None)
        return
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(tracking_em * font_size_pt * 100)))
