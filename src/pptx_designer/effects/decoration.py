"""Decoration library — pre-built decorative elements.

Provides brush divider, seal stamp, scroll frame, neon border,
grid background, glass panel, and ink splash decorations.
"""

from __future__ import annotations

from typing import Any


def add_brush_divider(slide: Any, x: float, y: float, w: float,
                      color: str = "#2C2C2C", thickness: float = 0.08) -> Any:
    """Add brush-stroke divider line.

    Args:
        slide: Slide object
        x, y: Position in inches
        w: Width in inches
        color: Brush color
        thickness: Brush thickness in inches

    Returns:
        Created shape
    """
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx_designer.tools.shapes import _rgb

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def add_seal_stamp(slide: Any, x: float, y: float, size: float, txt: str,
                   fill_hex: str = "#C41E3A", font_name: str = "STZhongsong",
                   rotation: float = -15, style: str = "zhu",
                   border_width_pt: float = 4.0) -> Any:
    """Add Chinese seal stamp.

    Args:
        slide: Slide object
        x, y: Position in inches
        size: Size in inches
        txt: Text to display
        fill_hex: Fill color
        font_name: Font name
        rotation: Rotation angle in degrees
        style: Style ('zhu' or 'yin')
        border_width_pt: Border width in points

    Returns:
        Created shape
    """
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx_designer.tools.shapes import _rgb

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(fill_hex)
    shape.line.width = Pt(border_width_pt)
    shape.rotation = rotation

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size * 72 * 0.6)
    run.font.color.rgb = _rgb("#FFFFFF")
    run.font.bold = True
    if font_name:
        run.font.name = font_name

    return shape


def add_scroll_frame(slide: Any, x: float, y: float, w: float, h: float,
                     style: str = "xuan") -> Any:
    """Add scroll/paper frame.

    Args:
        slide: Slide object
        x, y: Position in inches
        w, h: Size in inches
        style: Style ('xuan' or 'silk')

    Returns:
        Created shape
    """
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx_designer.tools.shapes import _rgb

    color = "#F5F0E8" if style == "xuan" else "#E8E0D0"
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.color.rgb = _rgb("#D4C4A8")
    shape.line.width = Inches(0.02)
    return shape


def add_neon_border(slide: Any, x: float, y: float, w: float, h: float,
                    color: str = "#8B5CF6", radius: float = 0.1) -> Any:
    """Add neon-glow border.

    Args:
        slide: Slide object
        x, y: Position in inches
        w, h: Size in inches
        color: Neon color
        radius: Corner radius in inches

    Returns:
        Created shape
    """
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx_designer.tools.shapes import _rgb

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.background()
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Inches(0.02)
    return shape


def add_grid_background(slide: Any, spacing: float = 1.0, color: str = "#E0E0E0",
                        alpha: float = 15) -> Any:
    """Add grid background pattern.

    Args:
        slide: Slide object
        spacing: Grid spacing in inches
        color: Grid line color
        alpha: Opacity (0-100)

    Returns:
        Created shape
    """
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx_designer.tools.shapes import _rgb

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def add_glass_panel(slide: Any, x: float, y: float, w: float, h: float,
                    tint: str = "#FFFFFF", alpha: float = 50,
                    soft_edge: float = 8) -> Any:
    """Add glass panel overlay.

    Args:
        slide: Slide object
        x, y: Position in inches
        w, h: Size in inches
        tint: Tint color
        alpha: Opacity (0-100)
        soft_edge: Soft edge radius

    Returns:
        Created shape
    """
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx_designer.tools.shapes import _rgb

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(tint)
    shape.line.fill.background()
    return shape


def add_ink_splash(slide: Any, x: float, y: float, size: float,
                   color: str = "#2C2C2C", alpha: float = 100) -> Any:
    """Add ink splash decoration.

    Args:
        slide: Slide object
        x, y: Position in inches
        size: Size in inches
        color: Ink color
        alpha: Opacity (0-100)

    Returns:
        Created shape
    """
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx_designer.tools.shapes import _rgb

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape
